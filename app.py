import os
import io
import re

import streamlit as st
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
from matplotlib.ticker import FuncFormatter
from openai import OpenAI
from docx import Document
from docx.shared import Inches, Pt as DocxPt, RGBColor as DocxRGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from openpyxl import Workbook, load_workbook
from openpyxl.styles import Font, PatternFill, Border, Side, Alignment
from openpyxl.utils import get_column_letter
from openpyxl.drawing.image import Image as XLImage
from pptx import Presentation
from pptx.util import Inches as PPTInches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor


def clean_excel_text(value):
    """
    Remove Markdown-style formatting before writing text into Excel.
    Keeps the words but removes symbols like **bold**, __bold__, `code`, and markdown headers.
    """
    if value is None:
        return ""

    text = str(value)

    text = re.sub(r"\*\*(.*?)\*\*", r"\1", text)
    text = re.sub(r"__(.*?)__", r"\1", text)
    text = re.sub(r"`(.*?)`", r"\1", text)
    text = re.sub(r"^\s{0,3}#{1,6}\s*", "", text)

    text = text.replace("**", "")
    text = text.replace("__", "")
    text = text.replace("`", "")

    text = re.sub(r"\s+", " ", text).strip()

    return text


def clean_doc_text(value):
    """Remove Markdown artifacts before writing generated content to Word."""
    if value is None:
        return ""
    text = str(value)
    text = re.sub(r"\*\*(.*?)\*\*", r"\1", text)
    text = re.sub(r"__(.*?)__", r"\1", text)
    text = re.sub(r"`(.*?)`", r"\1", text)
    text = text.replace("**", "").replace("__", "").replace("`", "")
    text = re.sub(r"^\s*---+\s*$", "", text, flags=re.MULTILINE)
    text = re.sub(r"^\s{0,3}#{1,6}\s*", "", text, flags=re.MULTILINE)
    return re.sub(r"\s+", " ", text).strip()


def format_currency_doc(value):
    try:
        num = float(value)
        return f"(${abs(num):,.0f})" if num < 0 else f"${num:,.0f}"
    except Exception:
        return clean_doc_text(value)


def format_funding_ask_doc(value):
    try:
        num = float(value)
        if num >= 1_000_000 and num % 1_000_000 == 0:
            return f"${num / 1_000_000:.0f} million"
        if num >= 1_000_000:
            return f"${num / 1_000_000:.1f} million"
        return f"${num:,.0f}"
    except Exception:
        return str(value)


def format_price_doc(value):
    """Format a per-unit price without changing financial-model calculations."""
    try:
        return f"${float(value):,.2f}"
    except Exception:
        return clean_doc_text(value)


def format_percent_doc(value):
    try:
        num = float(value)
        if abs(num) <= 1:
            num *= 100
        return f"{num:.1f}%"
    except Exception:
        return clean_doc_text(value)


# ==================================================
# OPENAI CLIENT
# ==================================================
client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))


# ==================================================
# PAGE CONFIG
# ==================================================
st.set_page_config(page_title="TurboPitch", layout="wide")

# =========================================================
# TURBOPITCH UI FLOW SETUP
# =========================================================

if "tp_step" not in st.session_state:
    st.session_state.tp_step = 1

if "ai_review_done" not in st.session_state:
    st.session_state.ai_review_done = False

step_labels = {
    1: "Idea Input",
    2: "Financial Model",
    3: "AI Pitch",
    4: "Business Plan",
    5: "Export",
}

workflow_steps = [
    "1. Idea Input", "2. Financial Model", "3. AI Pitch",
    "4. Business Plan", "5. Export",
]

if "workflow_step" not in st.session_state:
    st.session_state.workflow_step = workflow_steps[st.session_state.tp_step - 1]

def select_workflow_step():
    st.session_state.tp_step = workflow_steps.index(st.session_state.workflow_step) + 1

def next_step():
    if st.session_state.tp_step == 1 and not st.session_state.get("idea", "").strip():
        st.warning("Please enter your startup idea before continuing.")
        return

    if st.session_state.tp_step == 4 and not st.session_state.get("business_plan_output", "").strip():
        st.warning("Please generate your investor materials before going to Downloads.")
        return

    if st.session_state.tp_step < 5:
        st.session_state.tp_step += 1
        st.session_state.workflow_step = workflow_steps[st.session_state.tp_step - 1]

def prev_step():
    if st.session_state.tp_step > 1:
        st.session_state.tp_step -= 1
        st.session_state.workflow_step = workflow_steps[st.session_state.tp_step - 1]

def render_nav(position="top"):
    nav_left, nav_right = st.columns(2)

    with nav_left:
        if st.session_state.tp_step > 1:
            st.button(
                "⬅ Back",
                on_click=prev_step,
                use_container_width=True,
                key=f"back_{position}"
            )

    with nav_right:
        if st.session_state.tp_step < 5:
            st.button(
                "Next ➜",
                on_click=next_step,
                use_container_width=True,
                key=f"next_{position}"
            )

st.markdown("# TurboPitch")
st.caption("Build investor-ready materials through a focused startup workflow.")
st.radio(
    "Startup workflow", workflow_steps, horizontal=True, key="workflow_step",
    on_change=select_workflow_step, label_visibility="collapsed",
)
st.progress(st.session_state.tp_step / 5)
st.caption(f"Step {st.session_state.tp_step} of 5 — {step_labels[st.session_state.tp_step]}")
st.divider()

# ==================================================
# CUSTOM CSS
# ==================================================
st.markdown("""
<style>
.kpi-card {
    padding: 18px;
    border-radius: 12px;
    color: white;
    margin-bottom: 10px;
    box-shadow: 0 4px 10px rgba(0,0,0,0.15);
}
.kpi-blue {
    background: linear-gradient(135deg, #1f4e78, #2f75b5);
}
.kpi-green {
    background: linear-gradient(135deg, #1f6b52, #2f8c6a);
}
.kpi-red {
    background: linear-gradient(135deg, #8a3434, #b85050);
}
.kpi-gold {
    background: linear-gradient(135deg, #6d5a2a, #9b8240);
    color: white;
}
.kpi-title {
    font-size: 14px;
    font-weight: 600;
    opacity: 0.9;
}
.kpi-value {
    font-size: 28px;
    font-weight: 700;
    margin-top: 6px;
}
.snapshot-card {
    padding: 16px;
    border-radius: 12px;
    background: rgba(255,255,255,0.03);
    border: 1px solid rgba(255,255,255,0.08);
}
.finance-card {
    padding: 18px;
    border-radius: 10px;
    border: 1px solid #d7e0ea;
    background: #f8fafc;
    color: #172b4d;
    margin: 8px 0 16px;
}
</style>
""", unsafe_allow_html=True)

if st.session_state.tp_step == 1:
    st.title("TurboPitch")
    st.subheader("Turn startup ideas into investor-ready materials — then pressure-test them through a VC lens.")
    st.caption(
        "TurboPitch combines founder inputs, structured financial modeling, benchmark logic, market reality logic, "
        "and AI reasoning to evaluate startup assumptions through an investor-readiness lens."
    )
# ==================================================
# SESSION STATE DEFAULTS
# ==================================================
defaults = {
    "idea": "",
    "industry": "SaaS",
    "price_per_unit": 5.0,
    "year1_units": 80000,
    "growth_y2": 0.60,
    "growth_y3": 0.40,
    "cost_per_unit": 1.20,
    "opex_pct": 0.25,
    "fixed_overhead": 300000.0,
    "starting_cash": 600000.0,
    "pushback_pct": 25,
    "investor_feedback": "Investor says revenue assumptions are too aggressive and early customer acquisition may be unrealistic.",
    "sanity_output": "",
    "business_plan_output": "",
    "interrogation_output": "",
    "answer_builder_output": "",
    "assumption_helper_output": "",
    "assumption_mode": "Manual",
}

for key, value in defaults.items():
    if key not in st.session_state:
        st.session_state[key] = value


# ==================================================
# BENCHMARK DATA
# ==================================================
INDUSTRY_BENCHMARKS = {
    "Vertical SaaS / AI Software": {
        "gross_margin": (0.50, 0.85),
        "growth_y2": (0.50, 1.50),
        "growth_y3": (0.30, 1.00),
        "opex_pct": (0.40, 0.75),
        "year1_units": (1000, 75000),
    },
    "SaaS": {
        "gross_margin": (0.70, 0.90),
        "growth_y2": (0.50, 1.50),
        "growth_y3": (0.30, 1.00),
        "opex_pct": (0.40, 0.70),
        "year1_units": (1000, 50000),
    },
    "Marketplace / Local Services": {
        "gross_margin": (0.40, 0.70),
        "growth_y2": (0.50, 1.20),
        "growth_y3": (0.30, 0.90),
        "opex_pct": (0.35, 0.60),
        # Local services require supply, dispatch capacity and route density;
        # a high transaction count is not comparable to SaaS seats.
        "year1_units": (5000, 50000),
    },
    "Consumer Product": {
        "gross_margin": (0.40, 0.60),
        "growth_y2": (0.30, 0.70),
        "growth_y3": (0.20, 0.50),
        "opex_pct": (0.25, 0.50),
        "year1_units": (10000, 150000),
    },
    "Food / Delivery": {
        "gross_margin": (0.20, 0.35),
        "growth_y2": (0.30, 0.90),
        "growth_y3": (0.20, 0.60),
        "opex_pct": (0.30, 0.50),
        "year1_units": (10000, 120000),
    },
    "AI Startup": {
        "gross_margin": (0.50, 0.80),
        "growth_y2": (0.50, 1.50),
        "growth_y3": (0.30, 1.00),
        "opex_pct": (0.40, 0.75),
        "year1_units": (1000, 75000),
    },
}


# ==================================================
# TRUST / METHODOLOGY CONTENT
# ==================================================
TRUST_DATA_SOURCES = [
    {
        "category": "Founder Inputs",
        "examples": "Idea description, industry choice, business model hints, manual assumptions",
        "use_case": "Used as the base context for all generated outputs and assumption logic."
    },
    {
        "category": "Internal Benchmark Ranges",
        "examples": "Gross margin bands, growth ranges, opex ranges, Year 1 unit ranges by industry",
        "use_case": "Used to build directional starting assumptions and compare founder inputs to typical business model patterns."
    },
    {
        "category": "Pricing / Business Model Heuristics",
        "examples": "Higher pricing for B2B / enterprise / AI software, lower pricing for food, retail, consumer, and job-seeker concepts",
        "use_case": "Used to adapt suggested assumptions based on keywords, customer type, and business model."
    },
    {
        "category": "Reality Engine Logic",
        "examples": "Customer segment checks, pricing fit checks, adoption realism checks, enterprise sales friction checks",
        "use_case": "Used to challenge assumptions that may be mathematically clean but unrealistic in the real world."
    },
    {
        "category": "Financial Modeling Logic",
        "examples": "Revenue = price × units, COGS derived from target margin, opex as % of revenue plus fixed overhead, tax and ending cash flow logic",
        "use_case": "Used to translate assumptions into a structured 3-year model."
    },
    {
        "category": "AI Interpretation",
        "examples": "Investor-readiness analysis, founder answer prep, suggested assumption rationale",
        "use_case": "Used to explain what the numbers mean, where investors may push back, and what should be validated next."
    }
]

TRUST_LIMITATIONS = [
    "TurboPitch is a decision-support tool, not a guarantee of startup success or funding.",
    "Outputs depend heavily on the quality and realism of founder inputs.",
    "Benchmark ranges are directional and should support judgment, not replace it.",
    "Some niche industries may require more custom market research and expert review.",
    "Reality Engine checks are heuristics and are meant to improve realism, not replace live market research.",
    "AI commentary is generated from structured assumptions, benchmark logic, market reality heuristics, and financial reasoning, but should still be reviewed by the founder."
]

BENCHMARK_SOURCE_LABELS = {
    "gross_margin": "Benchmark basis: internal industry margin ranges and business model heuristics.",
    "growth": "Benchmark basis: internal startup growth ranges and investor-readiness scaling heuristics.",
    "opex": "Benchmark basis: internal cost structure ranges and early-stage operating model heuristics.",
    "units": "Benchmark basis: internal Year 1 traction ranges and go-to-market feasibility assumptions."
}


# ==================================================
# HELPER FUNCTIONS
# ==================================================
def clean_ai_text(text: str) -> str:
    if not text:
        return ""
    text = text.replace("###", "")
    text = text.replace("##", "")
    text = text.replace("**", "")
    text = text.replace("*", "")
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()

def escape_streamlit_markdown(text):
    """Escape dynamic text that Streamlit Markdown may interpret as formatting."""
    if text is None:
        return ""
    return str(text).replace("$", "\\$")


def format_currency_ui(value):
    try:
        return f"${float(value):,.2f}"
    except (TypeError, ValueError):
        return str(value)

def currency_tick_formatter(x, pos):
    if abs(x) >= 1_000_000:
        return f"${x/1_000_000:.1f}M"
    if abs(x) >= 1_000:
        return f"${x/1_000:.0f}K"
    return f"${x:,.0f}"


def create_revenue_chart_image(df: pd.DataFrame):
    fig, ax = plt.subplots(figsize=(8.5, 4.8))
    ax.plot(df["Year"], df["Revenue"], marker="o", linewidth=2.5)
    ax.set_title("Revenue Projection Summary", fontsize=16, fontweight="bold", pad=14)
    ax.set_xlabel("Year", fontsize=11)
    ax.set_ylabel("Revenue ($)", fontsize=11)
    ax.yaxis.set_major_formatter(FuncFormatter(currency_tick_formatter))
    ax.grid(True, alpha=0.25)
    ax.set_facecolor("white")
    fig.patch.set_facecolor("white")

    for spine in ["top", "right"]:
        ax.spines[spine].set_visible(False)

    img_stream = io.BytesIO()
    fig.tight_layout()
    fig.savefig(img_stream, format="png", dpi=220, bbox_inches="tight")
    img_stream.seek(0)
    plt.close(fig)
    return img_stream


def create_ppt_financial_chart_image(df: pd.DataFrame, theme=None):
    fig, ax = plt.subplots(figsize=(9.2, 4.9))
    x = np.arange(len(df["Year"]))
    width = 0.24
    chart_colors = (theme or {}).get("chart", ["#2563EB", "#38BDF8", "#0F1E3D"])

    ax.bar(x - width, df["Revenue"], width=width, label="Revenue", color=chart_colors[0])
    ax.bar(x, df["Net Income"], width=width, label="Net Income", color=chart_colors[1])
    ax.bar(x + width, df["Ending Cash"], width=width, label="Ending Cash", color=chart_colors[2])

    ax.set_title("Financial Projection Overview", fontsize=15, fontweight="bold", pad=14)
    ax.set_xlabel("Year", fontsize=10)
    ax.set_ylabel("USD", fontsize=10)
    ax.set_xticks(x)
    ax.set_xticklabels(df["Year"])
    ax.yaxis.set_major_formatter(FuncFormatter(currency_tick_formatter))
    ax.grid(True, axis="y", alpha=0.25)
    ax.set_axisbelow(True)
    ax.set_facecolor("white")
    fig.patch.set_facecolor("white")

    for spine in ["top", "right"]:
        ax.spines[spine].set_visible(False)

    ax.legend(
        loc="upper center",
        bbox_to_anchor=(0.5, -0.14),
        ncol=3,
        frameon=False,
        fontsize=9,
    )

    img_stream = io.BytesIO()
    fig.tight_layout()
    fig.savefig(img_stream, format="png", dpi=220, bbox_inches="tight")
    img_stream.seek(0)
    plt.close(fig)
    return img_stream


def create_excel_projection_chart_image(df: pd.DataFrame):
    fig, ax = plt.subplots(figsize=(9.5, 4.8))

    ax.plot(df["Year"], df["Revenue"], marker="o", linewidth=2.8, label="Revenue")
    ax.plot(df["Year"], df["Net Income"], marker="o", linewidth=2.8, label="Net Income")
    ax.plot(df["Year"], df["Ending Cash"], marker="o", linewidth=2.8, label="Ending Cash")

    ax.set_title("Revenue, Net Income & Cash by Year", fontsize=13, fontweight="bold", pad=14)
    ax.set_xlabel("Forecast Year", fontsize=10)
    ax.set_ylabel("USD", fontsize=10)
    ax.yaxis.set_major_formatter(FuncFormatter(currency_tick_formatter))
    ax.grid(True, axis="y", alpha=0.25)
    ax.set_facecolor("white")
    fig.patch.set_facecolor("white")

    for spine in ["top", "right"]:
        ax.spines[spine].set_visible(False)

    ax.legend(
        loc="upper center",
        bbox_to_anchor=(0.5, -0.18),
        ncol=3,
        frameon=False,
        fontsize=9
    )

    img_stream = io.BytesIO()
    fig.tight_layout(rect=[0.03, 0.08, 0.98, 0.95])
    fig.savefig(img_stream, format="png", dpi=240, bbox_inches="tight", facecolor="white")
    img_stream.seek(0)
    plt.close(fig)
    return img_stream


def create_excel_compare_chart_image(df: pd.DataFrame):
    fig, ax = plt.subplots(figsize=(9.5, 4.8))

    x = range(len(df["Year"]))
    width = 0.23

    ax.bar([i - width for i in x], df["Revenue"], width=width, label="Revenue")
    ax.bar(x, df["COGS"], width=width, label="COGS")
    ax.bar([i + width for i in x], df["Operating Expenses"], width=width, label="Operating Expenses")

    ax.set_xticks(list(x))
    ax.set_xticklabels(df["Year"])
    ax.set_title("Revenue vs COGS vs Operating Expenses", fontsize=13, fontweight="bold", pad=14)
    ax.set_xlabel("Year", fontsize=10)
    ax.set_ylabel("USD", fontsize=10)
    ax.yaxis.set_major_formatter(FuncFormatter(currency_tick_formatter))
    ax.grid(True, axis="y", alpha=0.25)
    ax.set_facecolor("white")
    fig.patch.set_facecolor("white")

    for spine in ["top", "right"]:
        ax.spines[spine].set_visible(False)

    ax.legend(
        loc="upper center",
        bbox_to_anchor=(0.5, -0.18),
        ncol=3,
        frameon=False,
        fontsize=9
    )

    img_stream = io.BytesIO()
    fig.tight_layout(rect=[0.03, 0.08, 0.98, 0.95])
    fig.savefig(img_stream, format="png", dpi=240, bbox_inches="tight", facecolor="white")
    img_stream.seek(0)
    plt.close(fig)
    return img_stream


def build_projection(price, year1_units, growth_y2, growth_y3, cost_per_unit, opex_pct, fixed_overhead, starting_cash):
    year2_units = int(year1_units * (1 + growth_y2))
    year3_units = int(year2_units * (1 + growth_y3))

    years = ["Year 1", "Year 2", "Year 3"]
    units = [year1_units, year2_units, year3_units]

    revenue = [u * price for u in units]
    cogs = [u * cost_per_unit for u in units]
    gross_profit = [r - c for r, c in zip(revenue, cogs)]
    gross_margin_pct = [(gp / r) if r != 0 else 0 for gp, r in zip(gross_profit, revenue)]

    operating_expenses = [(r * opex_pct) + fixed_overhead for r in revenue]
    operating_income = [gp - op for gp, op in zip(gross_profit, operating_expenses)]

    tax_rate = 0.21
    taxes = [max(0, oi * tax_rate) for oi in operating_income]
    net_income = [oi - tax for oi, tax in zip(operating_income, taxes)]

    ending_cash = []
    cash_balance = starting_cash
    for ni in net_income:
        cash_balance += ni
        ending_cash.append(cash_balance)

    return pd.DataFrame({
        "Year": years,
        "Units": units,
        "Revenue": revenue,
        "COGS": cogs,
        "Gross Profit": gross_profit,
        "Gross Margin %": gross_margin_pct,
        "Operating Expenses": operating_expenses,
        "Operating Income": operating_income,
        "Taxes": taxes,
        "Net Income": net_income,
        "Ending Cash": ending_cash,
    })


def get_benchmark_assumptions(industry: str, business_model: str) -> dict:
    """Return rule-based starter assumptions for a first-pass investor model."""
    industry_key, model_key = (industry or "").lower(), (business_model or "").lower()
    assumptions = {"pricing_period": "One-Time", "revenue_driver": "Customers", "price_per_unit": 29.0, "year_1_units": 500.0, "year_2_growth_rate": 1.00, "year_3_growth_rate": 0.75, "cost_per_unit": 6.0, "opex_pct_revenue": 0.55, "fixed_annual_overhead": 180000.0, "starting_cash": 250000.0}
    if industry == "Vertical SaaS / AI Software" and business_model == "Subscription":
        assumptions.update({"pricing_period": "Monthly", "revenue_driver": "Locations / Accounts", "price_per_unit": 249.0, "year_1_units": 150.0, "year_2_growth_rate": 1.00, "year_3_growth_rate": 0.70, "cost_per_unit": 35.0, "opex_pct_revenue": 0.55, "fixed_annual_overhead": 240000.0, "starting_cash": 300000.0})
    elif "saas" in industry_key or "software" in industry_key or "subscription" in model_key:
        assumptions.update({"pricing_period": "Monthly" if business_model == "Subscription" else "One-Time", "price_per_unit": 49.0, "year_1_units": 750.0, "year_2_growth_rate": 1.25, "year_3_growth_rate": 0.85, "cost_per_unit": 8.0, "opex_pct_revenue": 0.60, "fixed_annual_overhead": 240000.0, "starting_cash": 300000.0})
    elif "marketplace" in industry_key or "take rate" in model_key:
        assumptions.update({"price_per_unit": 15.0, "year_1_units": 2500.0, "year_2_growth_rate": 1.10, "year_3_growth_rate": 0.80, "cost_per_unit": 4.0, "opex_pct_revenue": 0.65, "fixed_annual_overhead": 220000.0, "starting_cash": 300000.0})
    elif "consumer" in industry_key or "product" in industry_key:
        assumptions.update({"price_per_unit": 35.0, "year_1_units": 3000.0, "year_2_growth_rate": 0.75, "year_3_growth_rate": 0.50, "cost_per_unit": 14.0, "opex_pct_revenue": 0.45, "fixed_annual_overhead": 150000.0, "starting_cash": 200000.0})
    elif "service" in industry_key or "agency" in model_key or "retainer" in model_key:
        assumptions.update({"price_per_unit": 2500.0, "year_1_units": 40.0, "year_2_growth_rate": 0.60, "year_3_growth_rate": 0.40, "cost_per_unit": 800.0, "opex_pct_revenue": 0.35, "fixed_annual_overhead": 120000.0, "starting_cash": 150000.0})
    elif "health" in industry_key:
        assumptions.update({"price_per_unit": 199.0, "year_1_units": 120.0, "year_2_growth_rate": 0.90, "year_3_growth_rate": 0.65, "cost_per_unit": 45.0, "opex_pct_revenue": 0.58, "fixed_annual_overhead": 260000.0, "starting_cash": 350000.0})
    elif "fintech" in industry_key or "finance" in industry_key:
        assumptions.update({"price_per_unit": 99.0, "year_1_units": 300.0, "year_2_growth_rate": 1.00, "year_3_growth_rate": 0.70, "cost_per_unit": 20.0, "opex_pct_revenue": 0.62, "fixed_annual_overhead": 300000.0, "starting_cash": 400000.0})
    return assumptions


def calculate_projection(assumptions: dict) -> pd.DataFrame:
    """Calculate a shared three-year projection from workflow assumptions."""
    price = float(assumptions.get("price_per_unit", 0)); y1 = float(assumptions.get("year_1_units", 0))
    pricing_period = assumptions.get("pricing_period", "One-Time")
    annualization_multiplier = 12 if pricing_period == "Monthly" else 1
    y2g = float(assumptions.get("year_2_growth_rate", 0)); y3g = float(assumptions.get("year_3_growth_rate", 0))
    cost = float(assumptions.get("cost_per_unit", 0)); opex = float(assumptions.get("opex_pct_revenue", 0))
    overhead = float(assumptions.get("fixed_annual_overhead", 0)); cash = float(assumptions.get("starting_cash", 0))
    rows = []
    for year, units in {"Year 1": y1, "Year 2": y1 * (1 + y2g), "Year 3": y1 * (1 + y2g) * (1 + y3g)}.items():
        revenue, cogs = units * price * annualization_multiplier, units * cost * annualization_multiplier; gross_profit = revenue - cogs; gross_margin = gross_profit / revenue if revenue else 0
        operating_expenses = revenue * opex + overhead; operating_income = gross_profit - operating_expenses; taxes = max(operating_income * 0.21, 0); net_income = operating_income - taxes; cash += net_income
        rows.append({"Year": year, "Customers / Units": round(units), "Units": round(units), "Revenue": revenue, "COGS": cogs, "Gross Profit": gross_profit, "Gross Margin": gross_margin, "Gross Margin %": gross_margin, "Operating Expenses": operating_expenses, "Operating Income": operating_income, "Taxes": taxes, "Net Income": net_income, "Ending Cash": cash})
    return pd.DataFrame(rows)


def apply_investor_haircut(assumptions: dict, revenue_reduction_pct: float = 25) -> dict:
    haircut = assumptions.copy(); reduction = 1 - (float(revenue_reduction_pct) / 100)
    haircut["year_1_units"] = float(haircut.get("year_1_units", 0)) * reduction
    haircut["year_2_growth_rate"] = float(haircut.get("year_2_growth_rate", 0)) * 0.85
    haircut["year_3_growth_rate"] = float(haircut.get("year_3_growth_rate", 0)) * 0.85
    haircut["opex_pct_revenue"] = min(float(haircut.get("opex_pct_revenue", 0)) + 0.05, 0.95)
    return haircut


def build_pnl_view(projection_df: pd.DataFrame) -> pd.DataFrame:
    pnl_df = pd.DataFrame({
        "Line Item": [
            "Units",
            "Revenue",
            "COGS",
            "Gross Profit",
            "Gross Margin %",
            "Operating Expenses",
            "Operating Income",
            "Taxes",
            "Net Income",
            "Ending Cash",
        ],
        "Year 1": [
            projection_df.loc[0, "Units"],
            projection_df.loc[0, "Revenue"],
            projection_df.loc[0, "COGS"],
            projection_df.loc[0, "Gross Profit"],
            projection_df.loc[0, "Gross Margin %"],
            projection_df.loc[0, "Operating Expenses"],
            projection_df.loc[0, "Operating Income"],
            projection_df.loc[0, "Taxes"],
            projection_df.loc[0, "Net Income"],
            projection_df.loc[0, "Ending Cash"],
        ],
        "Year 2": [
            projection_df.loc[1, "Units"],
            projection_df.loc[1, "Revenue"],
            projection_df.loc[1, "COGS"],
            projection_df.loc[1, "Gross Profit"],
            projection_df.loc[1, "Gross Margin %"],
            projection_df.loc[1, "Operating Expenses"],
            projection_df.loc[1, "Operating Income"],
            projection_df.loc[1, "Taxes"],
            projection_df.loc[1, "Net Income"],
            projection_df.loc[1, "Ending Cash"],
        ],
        "Year 3": [
            projection_df.loc[2, "Units"],
            projection_df.loc[2, "Revenue"],
            projection_df.loc[2, "COGS"],
            projection_df.loc[2, "Gross Profit"],
            projection_df.loc[2, "Gross Margin %"],
            projection_df.loc[2, "Operating Expenses"],
            projection_df.loc[2, "Operating Income"],
            projection_df.loc[2, "Taxes"],
            projection_df.loc[2, "Net Income"],
            projection_df.loc[2, "Ending Cash"],
        ],
    })
    return pnl_df


def build_display_pnl(pnl_df: pd.DataFrame) -> pd.DataFrame:
    display_pnl = pnl_df.copy()

    for col in ["Year 1", "Year 2", "Year 3"]:
        formatted_values = []
        for _, row in display_pnl.iterrows():
            value = row[col]
            line_item = row["Line Item"]

            if line_item == "Gross Margin %":
                formatted_values.append(f"{value:.1%}")
            elif line_item == "Units":
                formatted_values.append(f"{int(value):,}")
            else:
                formatted_values.append(f"${value:,.0f}")

        display_pnl[col] = formatted_values

    return display_pnl


# ==================================================
# REALITY ENGINE
# ==================================================
LOCAL_SERVICE_TERMS = (
    "marketplace", "mobile service", "booking", "on-demand",
    "detailer", "fleet", "route", "contractor", "provider", "vendor", "commission",
)

SOFTWARE_OPERATING_MODEL_TERMS = (
    "software platform", "saas", "subscription", "dashboard", "analytics",
    "ai-powered software", "workflow automation", "business operators",
    "independent operators", "car wash operators",
)


def is_software_industry(industry: str) -> bool:
    return industry in ("SaaS", "AI Startup", "B2B Software", "SaaS / AI Software", "Vertical SaaS / AI Software")


def classify_industry(idea: str, selected_industry: str = "") -> str:
    """Use the operating model implied by the idea, not a stale SaaS default."""
    idea_lower = (idea or "").lower()
    # Operating-model signals take precedence over the vertical served.  For
    # example, software sold to car washes is Vertical SaaS, not a car-wash
    # marketplace.
    if any(term in idea_lower for term in SOFTWARE_OPERATING_MODEL_TERMS):
        return "Vertical SaaS / AI Software"
    if any(term in idea_lower for term in LOCAL_SERVICE_TERMS):
        return "Marketplace / Local Services"
    return "Marketplace / Local Services" if selected_industry == "Marketplace" else (selected_industry or "SaaS")


def detect_customer_segment(idea: str, industry: str) -> str:
    idea_lower = (idea or "").lower()
    industry = classify_industry(idea, industry)

    if industry == "Marketplace / Local Services":
        return "Consumers / Fleets"

    if is_software_industry(industry) and any(term in idea_lower for term in (
        "operator", "operators", "business", "businesses", "company", "companies",
        "independent", "workflow", "dashboard", "analytics", "software platform",
    )):
        return "B2B / SMB Operators"

    enterprise_terms = [
        "enterprise", "b2b", "team", "company", "companies", "businesses",
        "workflow", "salesforce", "crm", "procurement", "operations",
        "analytics platform", "api", "infrastructure", "compliance", "hr software"
    ]
    consumer_terms = [
        "consumer", "job seeker", "resume", "dating", "fitness", "recipe", "pet owner",
        "student", "students", "creator", "freelancer", "marketplace app", "social app",
        "personal finance app", "subscription box", "food delivery", "local service"
    ]

    enterprise_hits = sum(1 for term in enterprise_terms if term in idea_lower)
    consumer_hits = sum(1 for term in consumer_terms if term in idea_lower)

    if enterprise_hits > consumer_hits:
        return "B2B / Enterprise"
    if consumer_hits > enterprise_hits:
        return "B2C / Consumer"

    if industry in ["Consumer Product", "Food / Delivery"]:
        return "B2C / Consumer"
    if is_software_industry(industry):
        return "Mixed / Unclear"
    return "Mixed / Unclear"


def pricing_market_check(idea: str, industry: str, price: float) -> dict:
    idea_lower = (idea or "").lower()
    industry = classify_industry(idea, industry)
    segment = detect_customer_segment(idea, industry)

    status = "Green"
    message = "Pricing appears broadly reasonable for the apparent customer segment."

    if segment == "B2C / Consumer":
        if any(word in idea_lower for word in ["resume", "job", "career", "job seeker"]):
            if price > 100:
                status = "Red"
                message = (
                    f"Price of ${price:,.2f} looks too high for a job-seeker / resume tool. "
                    "This audience is usually price-sensitive, so adoption may be far lower than the model assumes."
                )
            elif price > 50:
                status = "Yellow"
                message = (
                    f"Price of ${price:,.2f} may be high for a job-seeker / resume tool. "
                    "You may need stronger proof of premium value or a lower entry tier."
                )
        elif industry == "Food / Delivery":
            if price > 40:
                status = "Red"
                message = (
                    f"Price of ${price:,.2f} looks too high for a food / delivery concept unless it is premium or enterprise-driven."
                )
            elif price > 25:
                status = "Yellow"
                message = (
                    f"Price of ${price:,.2f} may be above what many food / delivery customers expect."
                )
        elif industry == "Consumer Product":
            if price > 150:
                status = "Red"
                message = (
                    f"Price of ${price:,.2f} may be too high for a consumer product unless you can clearly justify premium positioning."
                )
            elif price > 80:
                status = "Yellow"
                message = (
                    f"Price of ${price:,.2f} may limit adoption unless the product has strong perceived value."
                )
        else:
            if price > 100:
                status = "Red"
                message = (
                    f"Price of ${price:,.2f} looks high for a consumer-facing concept and may suppress demand."
                )
            elif price > 50:
                status = "Yellow"
                message = (
                    f"Price of ${price:,.2f} may be on the high side for a consumer-facing concept."
                )

    elif segment in ("B2B / Enterprise", "B2B / SMB Operators"):
        if price < 20 and is_software_industry(industry):
            status = "Yellow"
            message = (
                f"Price of ${price:,.2f} may be too low for a B2B operator software concept and could understate value."
            )
        elif price > 5000:
            status = "Yellow"
            message = (
                f"Price of ${price:,.2f} is very high and may require an enterprise sales motion, proof of ROI, and long sales cycles."
            )
    else:
        if is_software_industry(industry) and price > 300:
            status = "Yellow"
            message = (
                f"Price of ${price:,.2f} may be reasonable for some B2B software models, but the concept does not clearly prove enterprise willingness to pay."
            )

    return {
        "status": status,
        "message": message,
        "segment": segment,
    }


def volume_market_check(idea: str, industry: str, year1_units: int) -> dict:
    industry = classify_industry(idea, industry)
    segment = detect_customer_segment(idea, industry)

    status = "Green"
    message = "Initial Year 1 adoption looks directionally reasonable."

    if industry == "Marketplace / Local Services":
        _, unit_high = INDUSTRY_BENCHMARKS[industry]["year1_units"]
        if year1_units > unit_high:
            status = "Orange"
            message = (
                f"Year 1 booking target of {year1_units:,} is above the local-services benchmark and needs evidence of provider supply, booking frequency, route density, and repeat demand."
            )
    elif segment == "B2C / Consumer":
        if year1_units > 250000:
            status = "Red"
            message = (
                f"Year 1 unit target of {year1_units:,} looks extremely aggressive for a consumer startup without proven distribution."
            )
        elif year1_units > 100000:
            status = "Yellow"
            message = (
                f"Year 1 unit target of {year1_units:,} is ambitious for a consumer startup and likely requires major marketing or partnerships."
            )
    elif segment in ("B2B / Enterprise", "B2B / SMB Operators"):
        if is_software_industry(industry) and year1_units > 75000:
            status = "Red"
            message = (
                f"Year 1 unit target of {year1_units:,} is aggressive for a B2B software rollout and needs support from pipeline, conversion rate, paid-location count, retention, and rollout capacity."
            )
        elif year1_units > 20000:
            status = "Red"
            message = (
                f"Year 1 unit target of {year1_units:,} looks too high for a B2B / enterprise concept unless units represent very small transactions rather than customers."
            )
        elif year1_units > 5000:
            status = "Yellow"
            message = (
                f"Year 1 unit target of {year1_units:,} may be aggressive for a B2B / enterprise concept, especially if each unit implies a customer contract."
            )
    else:
        if year1_units > 150000:
            status = "Yellow"
            message = (
                f"Year 1 unit target of {year1_units:,} may be too aggressive unless customer acquisition is unusually efficient."
            )

    return {
        "status": status,
        "message": message,
    }


def growth_market_check(idea: str, industry: str, growth_y2: float, growth_y3: float) -> dict:
    industry = classify_industry(idea, industry)
    y2_low, y2_high = INDUSTRY_BENCHMARKS[industry]["growth_y2"]
    y3_low, y3_high = INDUSTRY_BENCHMARKS[industry]["growth_y3"]
    status = "Green"
    message = "Growth path is within the selected operating-model benchmark ranges."
    if growth_y2 > y2_high or growth_y3 > y3_high:
        status = "Orange"
        message = f"Growth of {growth_y2:.0%} in Year 2 and {growth_y3:.0%} in Year 3 is above the {industry} benchmark and needs evidence of repeat demand and delivery capacity."
    elif growth_y2 < y2_low or growth_y3 < y3_low:
        status = "Yellow"
        message = f"Growth is below the {industry} benchmark range and should be tied to an explicit operating plan."

    return {
        "status": status,
        "message": message,
    }


def opex_reality_check(industry: str, opex_pct: float, fixed_overhead: float, segment: str,
                       price: float, year1_units: int) -> dict:
    industry = classify_industry("", industry)
    total_opex_pct = opex_pct + (fixed_overhead / (price * year1_units) if price > 0 and year1_units > 0 else 0)
    opex_low, opex_high = INDUSTRY_BENCHMARKS[industry]["opex_pct"]
    status = "Green"
    message = f"Total OpEx is {total_opex_pct:.0%} of Year 1 revenue, including fixed overhead."

    if total_opex_pct > opex_high:
        status = "Orange"
        message = f"Total OpEx is {total_opex_pct:.0%} of Year 1 revenue (including ${fixed_overhead:,.0f} fixed overhead), above the {industry} benchmark range."
    elif total_opex_pct < opex_low:
        status = "Yellow"
        message = f"Total OpEx is {total_opex_pct:.0%} of Year 1 revenue (including ${fixed_overhead:,.0f} fixed overhead), below the {industry} benchmark range."
    if industry == "Marketplace / Local Services" and status != "Orange":
        message += " Validate provider operations, dispatch, insurance, refunds, payment fees, and support costs."
    elif segment == "B2C / Consumer":
        if total_opex_pct < 0.20:
            status = "Yellow"
            message = f"Total OpEx of {total_opex_pct:.0%} of revenue may be too low for a consumer startup that needs acquisition, support, and brand-building."
    elif segment == "B2B / Enterprise":
        if total_opex_pct < 0.25:
            status = "Yellow"
            message = (
                f"Total OpEx of {total_opex_pct:.0%} of revenue may be too low for a B2B / enterprise model that may require sales, onboarding, and account support."
            )

    if fixed_overhead < 100000 and is_software_industry(industry):
        status = "Yellow"
        message = (
            f"Fixed overhead of ${fixed_overhead:,.0f} may be too low for a software or AI startup once salaries, tooling, and infrastructure are considered."
        )

    return {
        "status": status,
        "message": message,
    }


def financial_assumption_reality_check(industry, price, cost_per_unit=None, ending_cash_final=None):
    industry = classify_industry("", industry)
    if price <= 0:
        return {
            "status": "Red",
            "message": "Price must be greater than zero before margin and cash assumptions can be evaluated."
        }

    if cost_per_unit is None:
        return {
            "status": "Green",
            "message": "Financial assumption realism check was limited because COGS was not provided."
        }

    gross_margin = (price - cost_per_unit) / price
    cogs_pct = cost_per_unit / price
    benchmark = INDUSTRY_BENCHMARKS.get(industry, {})
    gm_low, gm_high = benchmark.get("gross_margin", (0.40, 0.80))

    if ending_cash_final is not None and ending_cash_final < 0:
        return {
            "status": "Red",
            "message": f"Ending cash is negative at ${ending_cash_final:,.0f}. Investors will expect a clear funding plan, runway bridge, or burn reduction before trusting this model."
        }

    if cost_per_unit <= 0 or cogs_pct < 0.05:
        return {
            "status": "Red",
            "message": f"COGS is only {cogs_pct:.1%} of revenue, creating a {gross_margin:.1%} gross margin. Investors will likely challenge whether core delivery costs are missing."
        }

    if gross_margin > 0.95 or gross_margin > gm_high + 0.10:
        return {
            "status": "Red",
            "message": f"Gross margin is {gross_margin:.1%}, far above the typical {industry} range of {gm_low:.0%}-{gm_high:.0%}. This needs strong evidence or a more conservative COGS assumption."
        }

    if gross_margin > gm_high or gross_margin > 0.90:
        return {
            "status": "Yellow",
            "message": f"Gross margin is {gross_margin:.1%}, above the typical {industry} range of {gm_low:.0%}-{gm_high:.0%}. Investors may ask for proof that COGS is complete."
        }

    return {
        "status": "Green",
        "message": f"Gross margin is {gross_margin:.1%} and ending cash does not show an obvious funding gap."
    }


def run_reality_engine(idea, industry, price, year1_units, growth_y2, growth_y3, opex_pct, fixed_overhead, cost_per_unit=None, ending_cash_final=None):
    industry = classify_industry(idea, industry)
    pricing_check = pricing_market_check(idea, industry, price)
    volume_check = volume_market_check(idea, industry, year1_units)
    growth_check = growth_market_check(idea, industry, growth_y2, growth_y3)
    opex_check = opex_reality_check(industry, opex_pct, fixed_overhead, pricing_check["segment"], price, year1_units)
    financial_check = financial_assumption_reality_check(industry, price, cost_per_unit, ending_cash_final)

    checks = {
        "Customer Segment Detection": {
            "status": "Green",
            "message": f"Detected customer segment: {pricing_check['segment']}"
        },
        "Pricing Market Fit": {
            "status": pricing_check["status"],
            "message": pricing_check["message"]
        },
        "Adoption Realism": {
            "status": volume_check["status"],
            "message": volume_check["message"]
        },
        "Growth Realism": {
            "status": growth_check["status"],
            "message": growth_check["message"]
        },
        "Operating Model Reality": {
            "status": opex_check["status"],
            "message": opex_check["message"]
        },
        "Financial Assumption Reality": {
            "status": financial_check["status"],
            "message": financial_check["message"]
        },
    }

    reds = sum(1 for item in checks.values() if item["status"] == "Red")
    yellows = sum(1 for item in checks.values() if item["status"] in ("Yellow", "Orange"))

    if reds >= 2:
        overall = "Red"
        summary = "Reality Engine sees multiple real-world adoption, pricing, cash, or margin issues."
    elif reds == 1 or yellows >= 2:
        overall = "Yellow"
        summary = "Reality Engine sees some real-world issues that may weaken investor credibility."
    else:
        overall = "Green"
        summary = "Reality Engine sees no major market realism issues at a high level."

    return {
        "checks": checks,
        "overall": overall,
        "summary": summary,
    }


def reality_status_icon(status: str) -> str:
    if status == "Green":
        return "🟢"
    if status == "Yellow":
        return "🟡"
    return "🔴"


# ==================================================
# RULE-BASED / SCORECARD LOGIC
# ==================================================
def run_rule_based_sanity_check(price, year1_units, growth_y2, growth_y3, cost_per_unit, opex_pct, ending_cash_final=None):
    warnings = []

    if price <= 0:
        warnings.append("🔴 Price per unit must be greater than zero.")
        return warnings

    gross_margin = (price - cost_per_unit) / price if price else 0
    cogs_pct = cost_per_unit / price if price else 0

    if price < cost_per_unit:
        warnings.append("🔴 Price per unit is below cost per unit. This creates a structurally unprofitable business.")
    elif cost_per_unit <= 0:
        warnings.append("🔴 COGS is modeled at zero or below. Investors will likely view this as an unrealistic cost assumption unless there is clear evidence.")
    elif cogs_pct < 0.05:
        warnings.append("🔴 COGS is below 5% of revenue, implying an unusually high gross margin. Investors will likely challenge whether delivery, support, hosting, fulfillment, and payment costs are missing.")
    elif gross_margin > 0.90:
        warnings.append("🟠 Gross margin is above 90%. This may be defensible in some software models, but investors will expect strong proof that COGS is not understated.")
    elif gross_margin < 0.20:
        warnings.append("🟠 Gross margin is below 20%. Investors usually expect stronger margins for scalable startups.")
    elif gross_margin < 0.40:
        warnings.append("🟡 Gross margin is moderate. Investors may ask how margins improve as the company scales.")

    if ending_cash_final is not None and ending_cash_final < 0:
        warnings.append(
            f"🔴 Ending cash is negative at ${ending_cash_final:,.0f}. Investors will treat this as a funding gap and expect a clear plan for runway, burn reduction, or additional capital."
        )

    if year1_units > 300000:
        warnings.append("🔴 Year 1 sales volume is extremely high for an early-stage startup and may be unrealistic.")
    elif year1_units > 150000:
        warnings.append("🟠 Year 1 sales are ambitious. Investors may ask for proof of early demand or strong distribution.")

    if growth_y2 > 1.5:
        warnings.append("🔴 Year 2 growth above 150% is very aggressive and may be difficult to sustain.")
    elif growth_y2 > 1.0:
        warnings.append("🟠 Year 2 growth above 100% will likely require strong evidence of market demand.")

    if growth_y3 > 1.2:
        warnings.append("🟠 Year 3 growth above 120% suggests very rapid scaling and may be questioned by investors.")

    if opex_pct > 0.70:
        warnings.append("🔴 Operating expenses exceed 70% of revenue, which may make profitability difficult.")
    elif opex_pct > 0.50:
        warnings.append("🟠 Operating expenses are high relative to revenue. Investors may ask how efficiency improves over time.")

    if gross_margin < 0.30 and opex_pct > 0.40:
        warnings.append("🔴 Low gross margin combined with high operating costs may make this model difficult to scale profitably.")

    if year1_units > 150000 and growth_y2 > 1.0:
        warnings.append("🟠 High initial sales combined with very fast growth may raise investor skepticism about demand assumptions.")

    if not warnings:
        warnings.append("🟢 No major structural issues detected in the assumptions. The model appears reasonable at a high level.")

    return warnings


def build_warning_summary(warnings):
    red_phrases = (
        "Price per unit must be greater than zero",
        "COGS is modeled at zero",
        "COGS is below 5%",
        "Ending cash is negative",
        "Price per unit is below cost",
        "Year 1 sales volume is extremely high",
        "Year 2 growth above 150%",
        "Operating expenses exceed 70%",
        "Low gross margin combined with high operating costs",
    )
    amber_phrases = (
        "Gross margin is above 90%",
        "Gross margin is below 20%",
        "Gross margin is moderate",
        "Year 1 sales are ambitious",
        "Year 2 growth above 100%",
        "Year 3 growth above 120%",
        "Operating expenses are high",
        "High initial sales combined with very fast growth",
    )

    red_count = sum(1 for w in warnings if w.startswith("🔴"))
    amber_count = sum(1 for w in warnings if w.startswith("🟠") or w.startswith("🟡"))
    green_count = sum(1 for w in warnings if w.startswith("🟢"))

    red_count = sum(1 for w in warnings if w.startswith("🔴") or any(phrase in w for phrase in red_phrases))
    amber_count = sum(1 for w in warnings if w.startswith("🟠") or w.startswith("🟡") or any(phrase in w for phrase in amber_phrases))

    if red_count >= 2:
        overall = "High Assumption Risk"
        icon = "🔴"
    elif red_count == 1 or amber_count >= 2:
        overall = "Moderate Assumption Risk"
        icon = "🟠"
    elif green_count > 0 and red_count == 0 and amber_count == 0:
        overall = "Low Assumption Risk"
        icon = "🟢"
    else:
        overall = "Watchlist"
        icon = "🟡"

    return {
        "icon": icon,
        "overall": overall,
        "red_count": red_count,
        "amber_count": amber_count,
    }


def build_scorecard(idea, industry, projection_df, price, cost_per_unit, year1_units, growth_y2, growth_y3, reality_engine_output):
    industry = classify_industry(idea, industry)
    ending_cash_final = projection_df["Ending Cash"].iloc[-1]
    gross_margin = ((price - cost_per_unit) / price) if price else 0
    cogs_pct = cost_per_unit / price if price else 0
    benchmark = INDUSTRY_BENCHMARKS.get(industry, {})
    gm_high = benchmark.get("gross_margin", (None, 0.80))[1]

    structural_pricing_status = "Green" if price >= cost_per_unit * 1.5 else "Yellow" if price >= cost_per_unit else "Red"
    market_pricing_status = reality_engine_output["checks"]["Pricing Market Fit"]["status"]

    if structural_pricing_status == "Red" or market_pricing_status == "Red":
        pricing_status = "Red"
    elif structural_pricing_status == "Yellow" or market_pricing_status == "Yellow":
        pricing_status = "Yellow"
    else:
        pricing_status = "Green"

    sales_status = reality_engine_output["checks"]["Adoption Realism"]["status"]

    growth_status = reality_engine_output["checks"]["Growth Realism"]["status"]

    if cost_per_unit <= 0 or cogs_pct < 0.05 or gross_margin > 0.95 or gross_margin > gm_high + 0.10:
        margin_status = "Red"
    elif gross_margin > gm_high or gross_margin > 0.90:
        margin_status = "Yellow"
    else:
        margin_status = "Green" if gross_margin >= 0.50 else "Yellow" if gross_margin >= 0.25 else "Red"

    cash_status = "Green" if ending_cash_final > 0 else "Red"
    operating_model_status = reality_engine_output["checks"]["Operating Model Reality"]["status"]

    statuses = [pricing_status, sales_status, growth_status, margin_status, cash_status, operating_model_status]

    if statuses.count("Red") >= 2:
        overall = "Red"
    elif "Red" in statuses or sum(status in ("Yellow", "Orange") for status in statuses) >= 2:
        overall = "Orange"
    else:
        overall = "Green"

    return {
        "Pricing Realism": pricing_status,
        "Sales Volume": sales_status,
        "Growth Assumptions": growth_status,
        "Margin Quality": margin_status,
        "Cash Viability": cash_status,
        "Operating Model Reality": operating_model_status,
        "Overall Investor Readiness": overall,
    }


def score_metric(status):
    if status == "Green":
        return "🟢"
    if status == "Yellow":
        return "🟡"
    return "🔴"


def financial_summary_text(projection_df: pd.DataFrame) -> str:
    return "\n".join([
        f"{row['Year']}: Revenue ${row['Revenue']:,.0f}, "
        f"COGS ${row['COGS']:,.0f}, "
        f"Gross Profit ${row['Gross Profit']:,.0f}, "
        f"Gross Margin {row['Gross Margin %']:.1%}, "
        f"Operating Expenses ${row['Operating Expenses']:,.0f}, "
        f"Operating Income ${row['Operating Income']:,.0f}, "
        f"Taxes ${row['Taxes']:,.0f}, "
        f"Net Income ${row['Net Income']:,.0f}, "
        f"Ending Cash ${row['Ending Cash']:,.0f}"
        for _, row in projection_df.iterrows()
    ])


def build_benchmark_feedback(industry, price, cost_per_unit, year1_units, growth_y2, growth_y3, opex_pct, fixed_overhead=0):
    industry = classify_industry("", industry)
    benchmark = INDUSTRY_BENCHMARKS.get(industry)
    if not benchmark:
        return ["No benchmark data available for this industry."]

    feedback = []

    gross_margin = ((price - cost_per_unit) / price) if price > 0 else 0
    cogs_pct = (cost_per_unit / price) if price > 0 else 0

    gm_low, gm_high = benchmark["gross_margin"]
    y2_low, y2_high = benchmark["growth_y2"]
    y3_low, y3_high = benchmark["growth_y3"]
    opex_low, opex_high = benchmark["opex_pct"]
    unit_low, unit_high = benchmark["year1_units"]

    if gross_margin < gm_low:
        feedback.append(
            f"🔴 Gross margin is {gross_margin:.1%}, below the typical {industry} benchmark range of {gm_low:.0%}–{gm_high:.0%}. "
            f"Investors may question pricing power or scalability.\n{BENCHMARK_SOURCE_LABELS['gross_margin']}"
        )
    elif cost_per_unit <= 0 or cogs_pct < 0.05 or gross_margin > 0.95 or gross_margin > gm_high + 0.10:
        feedback.append(
            f"🔴 Gross margin is {gross_margin:.1%}, far above the typical {industry} benchmark range of {gm_low:.0%} to {gm_high:.0%}. "
            f"This may signal that COGS is understated or that delivery, support, hosting, fulfillment, and payment costs are missing.\n{BENCHMARK_SOURCE_LABELS['gross_margin']}"
        )
    elif gross_margin > gm_high:
        feedback.append(
            f"🟠 Gross margin is {gross_margin:.1%}, above the typical {industry} benchmark range of {gm_low:.0%} to {gm_high:.0%}. "
            f"Investors may ask whether COGS is complete and repeatable at scale.\n{BENCHMARK_SOURCE_LABELS['gross_margin']}"
        )
    elif gross_margin > gm_high:
        feedback.append(
            f"🟢 Gross margin is {gross_margin:.1%}, above the typical {industry} benchmark range of {gm_low:.0%}–{gm_high:.0%}.\n"
            f"{BENCHMARK_SOURCE_LABELS['gross_margin']}"
        )
    else:
        feedback.append(
            f"🟢 Gross margin is {gross_margin:.1%}, within the typical {industry} benchmark range of {gm_low:.0%}–{gm_high:.0%}.\n"
            f"{BENCHMARK_SOURCE_LABELS['gross_margin']}"
        )

    if growth_y2 > y2_high:
        feedback.append(
            f"🟠 Year 2 growth of {growth_y2:.0%} is above the usual {industry} benchmark range of {y2_low:.0%}–{y2_high:.0%}. "
            f"Investors may expect stronger proof of traction.\n{BENCHMARK_SOURCE_LABELS['growth']}"
        )
    elif growth_y2 < y2_low:
        feedback.append(
            f"🟡 Year 2 growth of {growth_y2:.0%} is below the usual {industry} benchmark range of {y2_low:.0%}–{y2_high:.0%}. "
            f"This may look conservative, but could also limit investor excitement.\n{BENCHMARK_SOURCE_LABELS['growth']}"
        )
    else:
        feedback.append(
            f"🟢 Year 2 growth of {growth_y2:.0%} is within the typical {industry} benchmark range of {y2_low:.0%}–{y2_high:.0%}.\n"
            f"{BENCHMARK_SOURCE_LABELS['growth']}"
        )

    if growth_y3 > y3_high:
        feedback.append(
            f"🟠 Year 3 growth of {growth_y3:.0%} is above the usual {industry} benchmark range of {y3_low:.0%}–{y3_high:.0%}. "
            f"Investors may view this as aggressive unless supported by strong momentum.\n{BENCHMARK_SOURCE_LABELS['growth']}"
        )
    elif growth_y3 < y3_low:
        feedback.append(
            f"🟡 Year 3 growth of {growth_y3:.0%} is below the usual {industry} benchmark range of {y3_low:.0%}–{y3_high:.0%}.\n"
            f"{BENCHMARK_SOURCE_LABELS['growth']}"
        )
    else:
        feedback.append(
            f"🟢 Year 3 growth of {growth_y3:.0%} is within the typical {industry} benchmark range of {y3_low:.0%}–{y3_high:.0%}.\n"
            f"{BENCHMARK_SOURCE_LABELS['growth']}"
        )

    total_opex_pct = opex_pct + (fixed_overhead / (price * year1_units) if price > 0 and year1_units > 0 else 0)
    if total_opex_pct > opex_high:
        feedback.append(
            f"🔴 Total operating expense ratio of {total_opex_pct:.0%} (including fixed overhead) is above the typical {industry} benchmark range of {opex_low:.0%}–{opex_high:.0%}. "
            f"Investors may question efficiency and burn discipline.\n{BENCHMARK_SOURCE_LABELS['opex']}"
        )
    elif total_opex_pct < opex_low:
        feedback.append(
            f"🟡 Total operating expense ratio of {total_opex_pct:.0%} (including fixed overhead) is below the typical {industry} benchmark range of {opex_low:.0%}–{opex_high:.0%}. "
            f"This may look efficient, but investors may ask whether growth investment is too light.\n{BENCHMARK_SOURCE_LABELS['opex']}"
        )
    else:
        feedback.append(
            f"🟢 Total operating expense ratio of {total_opex_pct:.0%} (including fixed overhead) is within the typical {industry} benchmark range of {opex_low:.0%}–{opex_high:.0%}.\n"
            f"{BENCHMARK_SOURCE_LABELS['opex']}"
        )

    if year1_units > unit_high:
        feedback.append(
            f"🟠 Year 1 volume of {year1_units:,} is above the typical {industry} benchmark range of {unit_low:,}–{unit_high:,}. "
            f"Investors may ask for stronger evidence of demand and distribution capacity.\n{BENCHMARK_SOURCE_LABELS['units']}"
        )
    elif year1_units < unit_low:
        feedback.append(
            f"🟡 Year 1 volume of {year1_units:,} is below the typical {industry} benchmark range of {unit_low:,}–{unit_high:,}.\n"
            f"{BENCHMARK_SOURCE_LABELS['units']}"
        )
    else:
        feedback.append(
            f"🟢 Year 1 volume of {year1_units:,} is within the typical {industry} benchmark range of {unit_low:,}–{unit_high:,}.\n"
            f"{BENCHMARK_SOURCE_LABELS['units']}"
        )

    suggestions = []

    if gross_margin < gm_low:
        target_price = cost_per_unit / (1 - gm_low) if gm_low < 1 else price
        target_cogs = price * (1 - gm_low)
        suggestions.append(
            f"Suggested fix: To reach the low end of the benchmark gross margin ({gm_low:.0%}), price may need to increase to about ${target_price:,.2f}, "
            f"or COGS may need to decrease to about ${target_cogs:,.2f} at the current price."
        )

    if growth_y2 > y2_high:
        suggestions.append(
            f"Suggested fix: Consider lowering Year 2 growth toward the benchmark range, such as {y2_low:.0%}–{y2_high:.0%}."
        )

    if growth_y3 > y3_high:
        suggestions.append(
            f"Suggested fix: Consider lowering Year 3 growth toward the benchmark range, such as {y3_low:.0%}–{y3_high:.0%}."
        )

    if total_opex_pct > opex_high:
        suggestions.append(
            f"Suggested fix: Consider reducing operating expenses toward the benchmark range of {opex_low:.0%}–{opex_high:.0%} of revenue."
        )

    if year1_units > unit_high:
        suggestions.append(
            f"Suggested fix: Consider reducing Year 1 units toward a more supportable range, such as {unit_low:,}–{unit_high:,}."
        )

    if not suggestions:
        suggestions.append("Current assumptions are broadly aligned with this industry's benchmark ranges.")

    return feedback + ["", "Suggested Adjustments"] + suggestions


def build_benchmark_feedback_rows(industry, price, cost_per_unit, year1_units,
                                  growth_y2, growth_y3, opex_pct, fixed_overhead=0):
    """Return export-ready benchmark results without presentation emojis."""
    industry = classify_industry("", industry)
    benchmark = INDUSTRY_BENCHMARKS.get(industry, INDUSTRY_BENCHMARKS["SaaS"])

    def assess(value, low, high, category, unit="%"):
        value_text = f"{value:.0%}" if unit == "%" else f"{value:,.0f}"
        range_text = f"{low:.0%}\u2013{high:.0%}" if unit == "%" else f"{low:,}\u2013{high:,}"
        if value > high:
            status = "Orange" if category != "Gross Margin" else "Yellow"
            assessment = "above"
        elif value < low:
            status = "Yellow"
            assessment = "below"
        else:
            status = "Green"
            assessment = "within"
        return status, value_text, range_text, assessment

    gross_margin = (price - cost_per_unit) / price if price else 0
    rows = []
    metric_specs = [
        ("Gross Margin", gross_margin, *benchmark["gross_margin"], "%"),
        ("Year 2 Growth", growth_y2, *benchmark["growth_y2"], "%"),
        ("Year 3 Growth", growth_y3, *benchmark["growth_y3"], "%"),
        ("Total OpEx Ratio", opex_pct + (fixed_overhead / (price * year1_units) if price > 0 and year1_units > 0 else 0), *benchmark["opex_pct"], "%"),
        ("Year 1 Volume", year1_units, *benchmark["year1_units"], "units"),
    ]
    for category, value, low, high, unit in metric_specs:
        status, value_text, range_text, assessment = assess(value, low, high, category, unit)
        if assessment == "within":
            feedback = f"{category} is within the typical {industry} benchmark range."
        elif category == "Total OpEx Ratio" and assessment == "below":
            feedback = "Total OpEx appears efficient, but may understate GTM, product, support, and fixed-overhead investment."
        elif category == "Year 1 Volume" and assessment == "above":
            feedback = "Year 1 volume is above the typical benchmark range and needs stronger demand evidence."
        elif assessment == "above":
            feedback = f"{category} of {value_text} is above the typical benchmark range and needs stronger support."
        else:
            feedback = f"{category} of {value_text} is below the typical benchmark range and should be supported in the plan."
        rows.append({
            "category": category,
            "status": status,
            "feedback": feedback,
            "basis": f"Internal range: {range_text}",
        })

    priorities = [row for row in rows if row["status"] in ("Orange", "Red", "Yellow")]
    if priorities:
        primary = priorities[0]
        suggested_fix = f"Address {primary['category'].lower()} with a revised assumption or stronger bottom-up evidence."
    else:
        suggested_fix = "Maintain the current assumptions and document the evidence supporting each benchmark-aligned input."
    return rows, suggested_fix


# ==================================================
# ASSUMPTION GENERATOR
# ==================================================
def generate_rule_based_assumptions(industry, idea_text=""):
    benchmark = INDUSTRY_BENCHMARKS.get(industry, INDUSTRY_BENCHMARKS["SaaS"])

    gm_low, gm_high = benchmark["gross_margin"]
    y2_low, y2_high = benchmark["growth_y2"]
    y3_low, y3_high = benchmark["growth_y3"]
    opex_low, opex_high = benchmark["opex_pct"]
    unit_low, unit_high = benchmark["year1_units"]

    target_gross_margin = round((gm_low + gm_high) / 2, 2)
    suggested_y2 = round((y2_low + y2_high) / 2, 2)
    suggested_y3 = round((y3_low + y3_high) / 2, 2)
    suggested_opex = round((opex_low + opex_high) / 2, 2)
    suggested_units = int((unit_low + unit_high) / 2)

    if industry == "SaaS":
        suggested_price = 99.0
        suggested_fixed_overhead = 250000.0
        suggested_starting_cash = 500000.0
    elif industry in ("Marketplace", "Marketplace / Local Services"):
        suggested_price = 35.0
        suggested_fixed_overhead = 300000.0
        suggested_starting_cash = 600000.0
    elif industry == "Consumer Product":
        suggested_price = 45.0
        suggested_fixed_overhead = 350000.0
        suggested_starting_cash = 500000.0
    elif industry == "Food / Delivery":
        suggested_price = 18.0
        suggested_fixed_overhead = 280000.0
        suggested_starting_cash = 400000.0
    elif industry == "AI Startup":
        suggested_price = 299.0
        suggested_fixed_overhead = 350000.0
        suggested_starting_cash = 700000.0
    else:
        suggested_price = 99.0
        suggested_fixed_overhead = 300000.0
        suggested_starting_cash = 500000.0

    idea_lower = (idea_text or "").lower()
    segment = detect_customer_segment(idea_text, industry)

    if any(word in idea_lower for word in ["enterprise", "b2b", "software", "platform", "analytics", "api", "workflow"]):
        suggested_price *= 1.35
        suggested_fixed_overhead *= 1.15

    if any(word in idea_lower for word in ["consumer", "app", "marketplace", "delivery", "subscription box"]):
        suggested_units = int(suggested_units * 1.20)

    if any(word in idea_lower for word in ["restaurant", "retail", "food", "cpg", "product"]):
        suggested_price *= 0.80

    if any(word in idea_lower for word in ["resume", "job", "career", "job seeker", "student", "consumer app"]):
        suggested_price = min(suggested_price, 39.0)
        suggested_units = int(suggested_units * 1.15)

    if segment == "B2C / Consumer" and industry in ["AI Startup", "SaaS"]:
        suggested_price = min(suggested_price, 49.0)

    if segment == "B2B / Enterprise" and suggested_price < 49:
        suggested_price = 49.0

    suggested_price = round(suggested_price, 2)
    suggested_cost = round(suggested_price * (1 - target_gross_margin), 2)
    suggested_fixed_overhead = round(suggested_fixed_overhead, 0)
    suggested_starting_cash = round(suggested_starting_cash, 0)

    return {
        "price_per_unit": suggested_price,
        "year1_units": suggested_units,
        "growth_y2": suggested_y2,
        "growth_y3": suggested_y3,
        "cost_per_unit": suggested_cost,
        "opex_pct": suggested_opex,
        "fixed_overhead": suggested_fixed_overhead,
        "starting_cash": suggested_starting_cash,
        "target_gross_margin": target_gross_margin,
        "detected_segment": segment,
    }


def build_reality_engine_explanation(idea, industry, suggested_values):
    projection = build_projection(
        suggested_values["price_per_unit"],
        suggested_values["year1_units"],
        suggested_values["growth_y2"],
        suggested_values["growth_y3"],
        suggested_values["cost_per_unit"],
        suggested_values["opex_pct"],
        suggested_values["fixed_overhead"],
        suggested_values["starting_cash"],
    )

    reality = run_reality_engine(
        idea,
        industry,
        suggested_values["price_per_unit"],
        suggested_values["year1_units"],
        suggested_values["growth_y2"],
        suggested_values["growth_y3"],
        suggested_values["opex_pct"],
        suggested_values["fixed_overhead"],
        suggested_values["cost_per_unit"],
        projection["Ending Cash"].iloc[-1],
    )

    lines = []
    lines.append("Reality Check Layer")
    lines.append(
        "TurboPitch also runs a Reality Engine after generating assumptions. This layer checks whether the assumptions look believable in the real world, not just whether they are mathematically consistent."
    )

    for label, item in reality["checks"].items():
        lines.append(f"{label}")
        lines.append(f"{reality_status_icon(item['status'])} {item['message']}")

    lines.append("Reality Engine Summary")
    lines.append(f"{reality_status_icon(reality['overall'])} {reality['summary']}")

    return "\n".join(lines)


def build_full_assumption_explanation(idea, industry, suggested_values):
    idea_text = idea if idea else "No startup idea was provided."

    price = suggested_values["price_per_unit"]
    units = suggested_values["year1_units"]
    g2 = suggested_values["growth_y2"]
    g3 = suggested_values["growth_y3"]
    cogs = suggested_values["cost_per_unit"]
    opex = suggested_values["opex_pct"]
    overhead = suggested_values["fixed_overhead"]
    cash = suggested_values["starting_cash"]
    target_margin = suggested_values["target_gross_margin"]
    segment = suggested_values.get("detected_segment", detect_customer_segment(idea, industry))

    benchmark = INDUSTRY_BENCHMARKS.get(industry, INDUSTRY_BENCHMARKS["SaaS"])
    gm_low, gm_high = benchmark["gross_margin"]
    y2_low, y2_high = benchmark["growth_y2"]
    y3_low, y3_high = benchmark["growth_y3"]
    opex_low, opex_high = benchmark["opex_pct"]
    unit_low, unit_high = benchmark["year1_units"]

    idea_lower = idea_text.lower()

    pricing_logic = []
    if any(word in idea_lower for word in ["enterprise", "b2b", "software", "platform", "analytics", "ai", "api"]):
        pricing_logic.append("The idea description suggests a B2B, software, platform, analytics, AI, or infrastructure-oriented concept, so the starting price was increased using a premium pricing heuristic.")
    if any(word in idea_lower for word in ["restaurant", "retail", "food", "cpg", "product"]):
        pricing_logic.append("The idea description suggests a lower-ticket or more price-sensitive product category, so the starting price was reduced using a consumer / product heuristic.")
    if any(word in idea_lower for word in ["resume", "job", "career", "job seeker", "student"]):
        pricing_logic.append("The idea description suggests an individual end-user and a price-sensitive audience, so the price was capped down using a consumer affordability heuristic.")
    if not pricing_logic:
        pricing_logic.append("The starting price was selected from the base pricing profile tied to the chosen industry, then checked against the target gross margin logic.")

    volume_logic = []
    if any(word in idea_lower for word in ["consumer", "app", "marketplace", "delivery", "subscription box"]):
        volume_logic.append("The idea description suggests a consumer or higher-volume model, so the Year 1 unit assumption was increased using a volume scaling heuristic.")
    else:
        volume_logic.append("The Year 1 unit assumption was anchored to the midpoint of the internal industry traction range so the model starts from a moderate rather than extreme adoption case.")

    reality_text = build_reality_engine_explanation(idea, industry, suggested_values)

    explanation = f"""
Suggested Assumption Rationale

Startup Idea Context
The current startup idea entered into TurboPitch is:
{idea_text}

Detected Customer Segment
TurboPitch interprets this concept as primarily:
{segment}

Why These Assumptions Make Sense
TurboPitch did not invent random numbers. The suggested assumptions were built in layers using a structured process.

First, the system looked at the selected industry, which is {industry}. That industry contains an internal set of benchmark ranges for gross margin, Year 2 growth, Year 3 growth, operating expense ratio, and Year 1 traction levels.

Second, TurboPitch created a baseline model from those benchmark ranges.
The system used the midpoint of the internal industry range for gross margin, which led to a target gross margin of about {target_margin:.0%}.
It also used the midpoint of the internal industry growth ranges to set Year 2 growth at about {g2:.0%} and Year 3 growth at about {g3:.0%}.
For Year 1 traction, it started from the midpoint of the internal range of {unit_low:,} to {unit_high:,}, which helped produce a starting Year 1 unit assumption of {units:,}.
For operating expenses, it used the midpoint of the internal opex range of {opex_low:.0%} to {opex_high:.0%}, which led to a starting opex ratio of {opex:.0%}.

Third, TurboPitch adjusted the base values using keyword and business model heuristics from the idea description.
{" ".join(pricing_logic)}
{" ".join(volume_logic)}

How The AI Came Up With Price
TurboPitch first selected a base starting price from the industry profile for {industry}.
It then adjusted that starting price using keyword-based business model rules tied to the idea description and customer segment.
After that, it recalculated cost per unit to keep the model aligned with the target gross margin.
The final suggested price is ${price:,.2f}.

How The AI Came Up With Cost Per Unit
TurboPitch did not guess cost per unit independently.
Instead, cost per unit was derived from the chosen price and the target gross margin.
The formula used is:
Cost per Unit = Price per Unit × (1 - Target Gross Margin)
With a suggested price of ${price:,.2f} and a target gross margin of about {target_margin:.0%}, the model produced a cost per unit of about ${cogs:,.2f}.
This makes cost structure consistent with the intended margin profile.

How The AI Came Up With Year 1 Units
TurboPitch used the internal industry Year 1 unit range for {industry}, which is {unit_low:,} to {unit_high:,}.
It started from the middle of that range to avoid making the first model too conservative or too aggressive.
If the idea description looked more consumer-oriented or volume-oriented, the model increased the Year 1 unit assumption using a simple scaling heuristic.
The resulting Year 1 unit assumption is {units:,}.

How The AI Came Up With Growth Rates
TurboPitch used the internal growth benchmark ranges for the selected industry.
For {industry}, the internal Year 2 growth range is {y2_low:.0%} to {y2_high:.0%}, and the internal Year 3 growth range is {y3_low:.0%} to {y3_high:.0%}.
The model used the middle of those ranges to create a balanced starting point.
That produced a Year 2 growth rate of {g2:.0%} and a Year 3 growth rate of {g3:.0%}.
These are intended to represent a believable early-stage scaling path rather than a maximum upside case.

How The AI Came Up With Operating Expense %
TurboPitch used the internal operating expense range for the selected industry.
For {industry}, the internal opex range is {opex_low:.0%} to {opex_high:.0%}.
The model used the midpoint of that range to generate a starting operating expense assumption of {opex:.0%}.
This is meant to reflect a practical early-stage operating burden before the founder fine-tunes hiring, marketing spend, support cost, and overhead structure.

How The AI Came Up With Fixed Overhead
Fixed overhead was not pulled from a live market database.
TurboPitch uses a baseline fixed overhead amount by industry type, then adjusts it upward when the idea appears more enterprise, software-heavy, or infrastructure-heavy.
That produced a starting fixed overhead value of ${overhead:,.0f}.
This number is meant to capture annual costs that do not directly scale with each unit sold, such as salaries, software, rent, and admin burden.

How The AI Came Up With Starting Cash
Starting cash is based on a baseline cash profile by industry type rather than on a live funding dataset.
TurboPitch uses higher starting cash for more capital-intensive or software-heavy concepts and lower starting cash for lighter models.
That produced a starting cash assumption of ${cash:,.0f}.
This is meant to represent a practical beginning runway assumption for an early-stage model, not a statement of what the founder must raise.

Reality Check Layer
{reality_text.replace("Reality Check Layer", "").strip()}

What The Founder Should Validate First
The founder should validate pricing first by checking what similar solutions, products, or services charge in the market.
Next, they should validate whether Year 1 volume is realistic based on access to customers, channels, partnerships, and demand generation capability.
Then they should validate costs, especially if the business depends on fulfillment, labor, onboarding, paid marketing, hosting, or AI infrastructure.
Finally, they should validate whether the chosen growth rates are actually possible given the go-to-market plan.

What May Need To Be Adjusted
These assumptions may need to change if the business is more premium, more commoditized, more service-heavy, more local, or more enterprise than the first draft suggests.
The founder may need to raise price, reduce Year 1 volume, increase cost per unit, increase opex, or raise starting cash depending on real-world conditions.
TurboPitch is giving a structured starting point, not a final truth.

How Conservative Or Aggressive This Starting Model Is
This starting model is intended to be moderately balanced.
It is not a worst-case stress test, and it is not a hyper-aggressive venture fantasy model either.
It is meant to give the founder a credible first-pass model that can be refined after more research and feedback.

Important Limitation
TurboPitch uses internal benchmark ranges, rule-based logic, Reality Engine heuristics, and AI explanation to build starter assumptions.
It does not currently pull live proprietary pricing databases or live real-time market feeds into this specific assumption generator.
The value of the feature is that it gives founders a rational starting model, challenges obvious real-world mismatches, and explains how it got there, so they are not guessing from scratch.
""".strip()

    return explanation


def render_ai_methodology_note():
    st.markdown("### How this analysis was generated")
    st.markdown(
        """
This review was generated using a combination of:

- founder-provided assumptions
- structured financial modeling
- rule-based benchmark logic
- market reality heuristics
- AI interpretation of realism, consistency, and investor risk

TurboPitch is designed to act like an investor-readiness filter, not a generic chatbot.
It evaluates whether assumptions look credible enough to survive investor scrutiny.
        """
    )


def render_trust_center():
    st.markdown("## Why Trust TurboPitch?")
    st.markdown(
        """
TurboPitch is designed to make AI output more credible by combining structured financial modeling,
industry benchmark logic, market reality checks, and investor-style analysis.

The goal is not to pretend the AI knows everything.
The goal is to show founders how their assumptions may look through the eyes of an investor.
        """
    )

    trust_tab1, trust_tab2, trust_tab3, trust_tab4 = st.tabs([
        "How the AI Works",
        "Data Sources",
        "Transparency",
        "Limitations"
    ])

    with trust_tab1:
        st.markdown("### How the AI Works")
        st.markdown(
            """
1. Founder inputs assumptions  
   The founder enters pricing, sales volume, growth, costs, and cash assumptions.

2. TurboPitch builds a structured financial model  
   Revenue, COGS, gross profit, operating expenses, taxes, net income, and ending cash are calculated.

3. Rule-based checks are applied  
   The model checks for structural issues such as negative margins, aggressive growth, high operating cost burden, and scale risk.

4. Benchmark logic is applied  
   Key assumptions are compared against industry ranges by business type.

5. Reality Engine checks are applied  
   The model checks for customer-segment fit, pricing plausibility, adoption realism, and operating model credibility.

6. AI interprets the results  
   The AI explains where investors may push back, what looks credible, and what should be improved.
            """
        )
        st.info("TurboPitch is designed to interpret structured financial logic and challenge obvious real-world mismatches.")

    with trust_tab2:
        st.markdown("### Data Sources Used for Benchmarking and Assumption Generation")
        for item in TRUST_DATA_SOURCES:
            st.markdown(f"**{item['category']}**")
            st.markdown(f"- Examples: {item['examples']}")
            st.markdown(f"- How it is used: {item['use_case']}")
        st.success("The purpose of this section is to show users that the analysis is grounded in recognizable business logic, benchmark structures, and explainable modeling steps.")

    with trust_tab3:
        st.markdown("### Trust and Transparency Principles")
        st.markdown(
            """
TurboPitch is built around a few simple trust principles:

- show the assumptions
- show the math
- show the benchmark ranges
- show the market-reality checks
- explain the reasoning
- acknowledge limitations

That means the platform should feel more like a structured venture analysis workflow and less like a black-box AI response.
            """
        )

        st.markdown("### Built on FP&A-Style and Investor-Style Thinking")
        st.markdown(
            """
TurboPitch uses logic similar to FP&A and business performance review frameworks:
pricing discipline, margin quality, growth realism, operating efficiency, cash viability, and customer willingness to pay.

This helps make the output more investor-relevant and more explainable.
            """
        )

    with trust_tab4:
        st.markdown("### Important Limitations")
        for item in TRUST_LIMITATIONS:
            st.warning(item)


# ==================================================
# AI FUNCTIONS
# ==================================================
def run_ai_assumption_helper(idea, industry, suggested_values):
    base_explanation = build_full_assumption_explanation(idea, industry, suggested_values)

    prompt = f"""
You are an FP&A analyst helping a first-time founder understand suggested startup assumptions.

You are given a structured internal explanation of how TurboPitch generated the starter assumptions.
Your job is to rewrite it so it sounds polished, clear, and founder-friendly while preserving the logic.

Startup idea:
{idea if idea else "No startup idea provided."}

Industry:
{industry}

Structured explanation:
{base_explanation}

Instructions:
- Keep the section title as Suggested Assumption Rationale
- Preserve the logic for how price, units, growth, cost per unit, opex, fixed overhead, and starting cash were derived
- Preserve the Reality Check Layer and explain it clearly
- Make it detailed and easy to understand
- Do not use markdown symbols like ##, ###, **, or bullet asterisks
- Do not invent live external data sources
- Make clear that TurboPitch uses internal benchmark ranges, heuristics, Reality Engine logic, and financial model logic
"""

    try:
        response = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[
                {
                    "role": "system",
                    "content": "You are a practical FP&A startup advisor helping founders understand starter business assumptions in a transparent way."
                },
                {"role": "user", "content": prompt},
            ],
            temperature=0.4,
        )

        content = response.choices[0].message.content if response.choices and response.choices[0].message else ""
        if content and content.strip():
            return content.strip()

        return base_explanation

    except Exception:
        return base_explanation


def run_ai_sanity_check(
    idea,
    industry,
    price_per_unit,
    year1_units,
    growth_y2,
    growth_y3,
    cost_per_unit,
    opex_pct,
    fixed_overhead,
    projection_df,
    reality_engine_output
):
    fin_summary = financial_summary_text(projection_df)
    reality_summary = "\n".join(
        [f"{k}: {v['status']} - {v['message']}" for k, v in reality_engine_output["checks"].items()]
    )

    prompt = f"""
You are a skeptical but constructive venture capital investor reviewing an early-stage startup.

Your job is to evaluate whether this startup looks investable based on the idea, business model, financial assumptions, and market realism checks.

Write a report with these exact section headings:

Overall Impression
What Investors Will Like
What Investors Will Challenge
Biggest Risks
Final Investor Verdict
Model Adjustments
Suggestions to Improve Investor Appeal
How to Reframe This Pitch for Investors

Instructions:
- Be practical, sharp, and investor-minded.
- Think like a real early-stage investor, not a cheerleader.
- Be honest if the model looks weak or unrealistic.
- Use both the financial assumptions and the Reality Engine checks.
- If pricing looks mathematically strong but unrealistic for the customer, say so clearly.
- Do not just critique the model — explain what should change.
- In Final Investor Verdict, answer clearly with one of these:
  Yes
  Maybe
  No
- In Final Investor Verdict, explain the reasoning in 2 to 4 sentences.
- In Model Adjustments, give 4 to 6 specific recommended changes using concrete numbers whenever possible.
- Consider the startup's industry when evaluating realism.
- Do NOT write a full business plan.
- Do NOT write a pitch deck.
- Do NOT use markdown symbols like ##, ###, **, or bullet asterisks.
- Keep formatting clean and readable.

Startup Idea:
{idea if idea else "No startup idea provided."}

Industry:
{industry}

Assumptions:
Price per unit: ${price_per_unit:,.2f}
Year 1 units sold: {year1_units:,}
Year 2 growth rate: {growth_y2:.2%}
Year 3 growth rate: {growth_y3:.2%}
Cost per unit: ${cost_per_unit:,.2f}
Operating expense percent of revenue: {opex_pct:.2%}
Fixed annual overhead: ${fixed_overhead:,.0f}

Projection Summary:
{fin_summary}

Reality Engine Summary:
{reality_summary}
"""

    try:
        response = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[
                {
                    "role": "system",
                    "content": "You are a disciplined early-stage VC investor who evaluates startup ideas for investability, realism, and investor appeal."
                },
                {"role": "user", "content": prompt},
            ],
            temperature=0.55,
        )
        return response.choices[0].message.content
    except Exception as e:
        return f"AI sanity check unavailable.\n\nError: {e}"


def run_ai_investor_interrogation(
    idea,
    industry,
    price_per_unit,
    year1_units,
    growth_y2,
    growth_y3,
    cost_per_unit,
    opex_pct,
    fixed_overhead,
    starting_cash,
    projection_df,
    reality_engine_output
):
    fin_summary = financial_summary_text(projection_df)
    reality_summary = "\n".join(
        [f"{k}: {v['status']} - {v['message']}" for k, v in reality_engine_output["checks"].items()]
    )

    prompt = f"""
You are a skeptical venture capitalist in a pitch meeting.

Your job is to challenge the founder with hard but realistic investor questions.

Write a section called:

Investor Questions

Then list 10 to 12 specific, tough questions an investor would ask this founder.

Instructions:
- Be sharp, skeptical, and practical.
- Consider the startup's industry when generating investor questions.
- Use the financial assumptions and the Reality Engine checks.
- Focus on weak assumptions, growth realism, market adoption, customer acquisition, competition, margins, operating costs, pricing realism, scaling risk, and funding use.
- Questions should sound like real VC pushback in a live meeting.
- Do NOT answer the questions.
- Do NOT write a business plan.
- Do NOT write in essay form.
- Keep each question concise and direct.
- Do NOT use markdown symbols like ##, ###, **, or bullet asterisks.

Startup Idea:
{idea if idea else "No startup idea provided."}

Industry:
{industry}

Assumptions:
Price per unit: ${price_per_unit:,.2f}
Year 1 units sold: {year1_units:,}
Year 2 growth rate: {growth_y2:.2%}
Year 3 growth rate: {growth_y3:.2%}
Cost per unit: ${cost_per_unit:,.2f}
Operating expense percent of revenue: {opex_pct:.2%}
Fixed annual overhead: ${fixed_overhead:,.0f}
Starting cash: ${starting_cash:,.0f}

Projection Summary:
{fin_summary}

Reality Engine Summary:
{reality_summary}
"""

    try:
        response = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[
                {
                    "role": "system",
                    "content": "You are a skeptical venture capitalist who challenges startup founders with tough investor questions."
                },
                {"role": "user", "content": prompt},
            ],
            temperature=0.7,
        )
        return response.choices[0].message.content
    except Exception as e:
        return f"Investor interrogation unavailable.\n\nError: {e}"


def run_ai_founder_answer_builder(
    idea,
    industry,
    price_per_unit,
    year1_units,
    growth_y2,
    growth_y3,
    cost_per_unit,
    opex_pct,
    fixed_overhead,
    starting_cash,
    projection_df,
    reality_engine_output
):
    fin_summary = financial_summary_text(projection_df)
    reality_summary = "\n".join(
        [f"{k}: {v['status']} - {v['message']}" for k, v in reality_engine_output["checks"].items()]
    )

    prompt = f"""
You are a startup pitch coach helping a founder prepare for tough investor meetings.

Create a section called:

Founder Answer Prep

Generate 8 investor-style questions the founder is likely to face.

For each question, include these exact sub-sections:

Investor Question
What Investors Are Really Asking
Strong Answer Framework
Weak Answer Example
How To Improve The Answer

Instructions:
- Be practical, realistic, and investor-focused.
- Consider the startup's industry when creating the questions and answer prep.
- Use the startup idea, assumptions, and Reality Engine summary to create the questions.
- In Strong Answer Framework, give concise talking points, not a long essay.
- In Weak Answer Example, show the kind of vague answer founders often give.
- In How To Improve The Answer, explain how to make the response more credible.
- Do NOT use markdown symbols like ##, ###, **, or bullet asterisks.
- Keep formatting clean and readable.

Startup Idea:
{idea if idea else "No startup idea provided."}

Industry:
{industry}

Assumptions:
Price per unit: ${price_per_unit:,.2f}
Year 1 units sold: {year1_units:,}
Year 2 growth rate: {growth_y2:.2%}
Year 3 growth rate: {growth_y3:.2%}
Cost per unit: ${cost_per_unit:,.2f}
Operating expense percent of revenue: {opex_pct:.2%}
Fixed annual overhead: ${fixed_overhead:,.0f}
Starting cash: ${starting_cash:,.0f}

Projection Summary:
{fin_summary}

Reality Engine Summary:
{reality_summary}
"""

    try:
        response = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[
                {
                    "role": "system",
                    "content": "You are a practical startup pitch coach who helps founders answer investor questions more effectively."
                },
                {"role": "user", "content": prompt},
            ],
            temperature=0.7,
        )
        return response.choices[0].message.content
    except Exception as e:
        return f"Founder answer builder unavailable.\n\nError: {e}"


def generate_business_plan_and_deck(
    idea,
    industry,
    price_per_unit,
    year1_units,
    growth_y2,
    growth_y3,
    cost_per_unit,
    opex_pct,
    fixed_overhead,
    projection_df,
    reality_engine_output
):
    fin_summary = financial_summary_text(projection_df)
    reality_summary = "\n".join(
        [f"{k}: {v['status']} - {v['message']}" for k, v in reality_engine_output["checks"].items()]
    )

    prompt = f"""
You are a startup strategist helping a founder prepare investor-ready materials.

Write a FULL BUSINESS PLAN followed by a separate section called Pitch Deck Content.

Use these exact BUSINESS PLAN sections:

Executive Summary
Problem
Solution
Market Opportunity
Product Overview
Business Model
Go-To-Market Strategy
Competitive Landscape
Financial Overview
Projection Assumptions & Investor Interpretation
Funding Ask
Key Risks

Then create a separate section called:

Pitch Deck Content

For Pitch Deck Content:
- Use numbered slides
- Each slide should have a slide title
- Each slide should include 3 to 4 short, persuasive bullet points
- Make the slide bullets investor-friendly and realistic
- Incorporate the Reality Engine concerns where appropriate

Slides to generate:
1. Problem
2. Solution
3. Product
4. Market Opportunity
5. Business Model
6. Competitive Advantage
7. Go-To-Market Strategy
8. Financial Highlights
9. Projection Assumptions & Investor Interpretation
10. Funding Ask

Instructions:
- Write clearly and professionally
- Consider the startup's industry in the business plan and deck content.
- Use the Reality Engine to avoid blindly endorsing unrealistic assumptions.
- For Projection Assumptions & Investor Interpretation, explain the revenue build, customer/unit assumptions, pricing logic, gross margin logic, operating expense logic, cash runway, funding need vs. milestones, and likely investor pushback using the Assumptions, Projection Summary, and Reality Engine Summary. Do not invent unsupported numbers.
- Do NOT use markdown symbols like ##, ###, **, or bullet asterisks
- Keep the business plan polished and readable
- Make the deck bullets concise and presentation-ready

Startup Idea:
{idea if idea else "No startup idea provided."}

Industry:
{industry}

Assumptions:
Price per unit: ${price_per_unit:,.2f}
Year 1 units sold: {year1_units:,}
Year 2 growth rate: {growth_y2:.2%}
Year 3 growth rate: {growth_y3:.2%}
Cost per unit: ${cost_per_unit:,.2f}
Operating expense percent of revenue: {opex_pct:.2%}
Fixed annual overhead: ${fixed_overhead:,.0f}

Projection Summary:
{fin_summary}

Reality Engine Summary:
{reality_summary}
"""

    try:
        response = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[
                {
                    "role": "system",
                    "content": "You are a practical startup strategist who creates investor-facing business plans and pitch decks."
                },
                {"role": "user", "content": prompt},
            ],
            temperature=0.7,
        )
        return response.choices[0].message.content
    except Exception as e:
        return f"Business plan generation unavailable.\n\nError: {e}"


def extract_business_plan_section(text: str) -> str:
    cleaned = clean_ai_text(text)
    if not cleaned:
        return ""
    if "Pitch Deck Content" in cleaned:
        return cleaned.split("Pitch Deck Content", 1)[0].strip()
    return cleaned


def extract_pitch_deck_section(text: str):
    cleaned = clean_ai_text(text)
    if not cleaned:
        return [("Startup Overview", ["No pitch deck content available yet."])]

    section_match = re.search(r"(?im)^\s*pitch\s+deck\s+content\s*:?\s*$", cleaned)
    if not section_match:
        return [("Startup Overview", ["No pitch deck content available yet."])]

    section = cleaned[section_match.end():].strip()
    lines = [line.strip() for line in section.splitlines() if line.strip()]

    slides = []
    current_title = None
    current_bullets = []

    for line in lines:
        title_match = re.match(r"^(?:slide\s*)?\d+[\.\):\-]\s*(.+)$", line, flags=re.IGNORECASE)
        if title_match:
            if current_title:
                slides.append((current_title, current_bullets))
            current_title = title_match.group(1).strip()
            current_bullets = []
        else:
            current_bullets.append(re.sub(r"^[-\u2022]\s*", "", line).strip())

    if current_title:
        slides.append((current_title, current_bullets))

    if not slides:
        return [("Startup Overview", ["No pitch deck content available yet."])]

    return slides


def is_empty_pitch_deck(slides) -> bool:
    return (
        not slides
        or (
            len(slides) == 1
            and slides[0][0] == "Startup Overview"
            and slides[0][1] == ["No pitch deck content available yet."]
        )
    )


def build_standard_pitch_deck_slides(
    idea,
    industry,
    price_per_unit,
    year1_units,
    growth_y2,
    growth_y3,
    cost_per_unit,
    opex_pct,
    fixed_overhead,
    starting_cash,
    projection_df,
    reality_engine_output,
):
    idea_text = str(idea or "").strip()
    idea_summary = idea_text.splitlines()[0] if idea_text else "This startup"
    if len(idea_summary) > 130:
        idea_summary = idea_summary[:127].rstrip() + "..."

    year1 = projection_df.loc[0]
    year3 = projection_df.loc[2]
    gross_margin = ((price_per_unit - cost_per_unit) / price_per_unit) if price_per_unit else 0
    customer_segment = detect_customer_segment(idea_text, industry)
    idea_lower = idea_text.lower()
    annual_growth = ((year3["Revenue"] / year1["Revenue"]) - 1) if year1["Revenue"] else 0

    revenue_ops_terms = any(
        term in idea_lower
        for term in ["pipeline", "renewal", "forecast", "pricing", "approval", "revenue", "sales", "crm"]
    )
    if industry == "SaaS" and revenue_ops_terms:
        pain_line = "B2B revenue teams lose time and accuracy when pipeline, renewals, pricing approvals, and forecasts live in separate workflows."
        solution_line = "The product centralizes revenue workflow signals so teams can track pipeline, renewal risk, pricing approvals, and forecast changes in one place."
        product_line = "Core workflows can include pipeline status, renewal alerts, pricing approval queues, and forecast views tied to account activity."
        market_line = "The reachable starting segment is B2B teams managing recurring revenue, renewals, pipeline visibility, and forecast discipline."
        advantage_one = "SMB Focus: packages revenue workflow discipline for teams that need visibility before enterprise tooling is practical."
        advantage_two = "All-in-One Platform: connects pipeline, renewals, pricing approvals, and forecasts instead of leaving them in separate trackers."
        advantage_three = "Fast Iteration: early customer feedback can quickly sharpen alerts, dashboards, and approval workflows."
    else:
        pain_line = f"{customer_segment} customers need a clearer way to solve the problem described in the founder idea: {idea_summary}"
        solution_line = f"The product turns the founder idea into a practical {industry} workflow with measurable customer actions and outcomes."
        product_line = f"Core workflows should map directly to the startup concept: {idea_summary}"
        market_line = f"The reachable starting segment is {customer_segment} customers whose current process creates enough friction to pay for a better option."
        advantage_one = f"Segment Focus: starts with {customer_segment} customers instead of spreading effort across unrelated use cases."
        advantage_two = f"Integrated Workflow: packages the main {industry} use case into one repeatable customer experience."
        advantage_three = "Fast Iteration: uses early customer feedback to improve onboarding, activation, and retention loops."

    challenged_checks = [
        f"{label}: {item['message']}"
        for label, item in reality_engine_output["checks"].items()
        if item["status"] != "Green"
    ]
    credible_checks = [
        f"{label}: {item['message']}"
        for label, item in reality_engine_output["checks"].items()
        if item["status"] == "Green"
    ]
    investor_focus = challenged_checks or credible_checks or [reality_engine_output["summary"]]
    risk_signal = investor_focus[0]
    evidence_signal = investor_focus[1] if len(investor_focus) > 1 else "Evidence needed: prove demand quality, conversion, retention, and sales efficiency with early traction data."

    return [
        ("Problem", [
            pain_line,
            "Fragmented tracking creates visibility gaps, slower decisions, and weaker accountability across the operating workflow.",
            "The investor question is whether this pain is frequent, budgeted, and painful enough to drive repeatable adoption.",
        ]),
        ("Solution", [
            solution_line,
            f"The model assumes a ${price_per_unit:,.2f} price point, so the value proposition must be clear enough to support paid usage.",
            f"Reality Engine signal: {reality_engine_output['summary']}",
        ]),
        ("Product", [
            product_line,
            "The first product proof should show that users return to the workflow and trust the data enough to act on it.",
            "The roadmap should prioritize the few actions that most directly improve visibility, conversion, retention, or decision speed.",
        ]),
        ("Market Opportunity", [
            market_line,
            f"The initial plan models {year1_units:,} Year 1 units, making reach and conversion quality more important than broad market claims.",
            f"Growth from Year 1 to Year 3 implies {annual_growth:.0%} revenue expansion, which needs a repeatable acquisition motion.",
        ]),
        ("Business Model", [
            f"Revenue is built from ${price_per_unit:,.2f} per unit multiplied by modeled unit volume.",
            f"Unit economics assume ${cost_per_unit:,.2f} COGS per unit and a {gross_margin:.1%} gross margin.",
            f"Operating spend is modeled at {opex_pct:.0%} of revenue plus ${fixed_overhead:,.0f} in fixed annual overhead.",
        ]),
        ("Competitive Advantage", [
            advantage_one,
            advantage_two,
            advantage_three,
        ]),
        ("Go-To-Market Strategy", [
            f"Start with {customer_segment} buyers who already feel the workflow pain and can validate willingness to pay.",
            f"Acquisition needs to support {year1_units:,} Year 1 units without relying on unsupported channel scale.",
            "Retention proof should come from repeat workflow usage, renewal intent, and clear before-and-after operating metrics.",
        ]),
        ("Financial Highlights", [
            f"Year 1 revenue is ${year1['Revenue']:,.0f}; Year 3 revenue is ${year3['Revenue']:,.0f}.",
            f"Year 3 net income is ${year3['Net Income']:,.0f}.",
            f"Year 3 ending cash is ${year3['Ending Cash']:,.0f}.",
        ]),
        ("Projection Assumptions & Investor Interpretation", [
            f"Credibility: the model depends on {growth_y2:.0%} Year 2 growth, {growth_y3:.0%} Year 3 growth, and consistent unit demand.",
            f"Risk: {risk_signal}",
            f"Evidence needed: connect ${starting_cash:,.0f} starting cash to customer acquisition, product delivery, and measurable milestones.",
        ]),
        ("Funding Ask", [
            "Use of funds should prioritize product validation, customer acquisition tests, and the operating capacity needed to deliver the workflow.",
            f"The plan should show how starting cash of ${starting_cash:,.0f} supports the path to Year 1 revenue of ${year1['Revenue']:,.0f}.",
            evidence_signal,
        ]),
    ]


# ==================================================
# MAIN IDEA INPUT
# ==================================================
if st.session_state.tp_step == 1:
    st.subheader("Describe Your Startup")

    idea_input = st.text_area(
        "Describe Your Startup",
        value=st.session_state.get("idea", ""),
        height=140,
        placeholder="Explain what the business does, who it serves, and why it matters..."
    )
    st.session_state.idea = idea_input
# ==================================================
# MAIN WORKFLOW: FINANCIAL ASSUMPTIONS
# ==================================================
def render_financial_assumption_step(startup_idea: str = "") -> dict:
    st.markdown("## Financial Model")
    st.caption("Build a three-year operating model with clear, investor-facing outputs.")

    if startup_idea:
        st.success(f"**Your startup idea:** {startup_idea}")
    else:
        st.warning("⚠️ No startup idea detected. Go back to Step 1 and enter your idea before generating assumptions.")

    st.markdown("Choose how you want to build your model:")
    mode_options = ["Use Benchmarks", "Enter My Own Numbers"]
    if st.session_state.get("finance_assumption_mode") not in mode_options:
        st.session_state["finance_assumption_mode"] = mode_options[0]
    mode = st.radio(
        "Assumption mode",
        mode_options,
        horizontal=True,
        key="finance_assumption_mode",
        label_visibility="collapsed",
    )
    mode_help = {
        "Use Benchmarks": "Best if you are early and do not know pricing, customer volume, COGS, or OpEx yet.",
        "Enter My Own Numbers": "Best if you already know your pricing, volume, growth, and cost assumptions.",
    }
    st.caption(mode_help[mode])

    a, b, c = st.columns(3)
    with a: industry = st.selectbox("Startup industry", ["Vertical SaaS / AI Software", "Marketplace / Local Services", "Consumer Product", "Service Business", "FinTech", "Healthcare", "Food / Delivery", "General Business"], key="finance_industry")
    with b: business_model = st.selectbox("Business model", ["Subscription", "Usage-Based", "One-Time Purchase", "Marketplace Take Rate", "Service / Retainer", "Hybrid"], key="finance_business_model")
    revenue_driver_options = ["Customers", "Users", "Units Sold", "Subscriptions", "Transactions", "Locations", "Locations / Accounts", "Projects"]
    revenue_driver_signature = f"{industry}|{business_model}"
    if industry == "Vertical SaaS / AI Software" and business_model == "Subscription" and st.session_state.get("finance_revenue_driver_signature") != revenue_driver_signature:
        st.session_state["finance_revenue_driver"] = "Locations / Accounts"
        st.session_state["finance_revenue_driver_signature"] = revenue_driver_signature
    with c: revenue_driver = st.selectbox("Primary revenue driver", revenue_driver_options, key="finance_revenue_driver")
    benchmark = get_benchmark_assumptions(industry, business_model)
    use_benchmarks = mode == "Use Benchmarks"
    defaults = benchmark if use_benchmarks else st.session_state.get("finance_assumptions", benchmark)
    pricing_period_options = ["Monthly", "Annual", "One-Time"]
    if "finance_pricing_period" not in st.session_state:
        st.session_state["finance_pricing_period"] = defaults.get(
            "pricing_period", "Monthly" if business_model == "Subscription" else "One-Time"
        )
    benchmark_signature = f"{industry}|{business_model}|{mode}"
    widget_defaults = {
        "finance_price_per_unit": "price_per_unit",
        "finance_year_1_units": "year_1_units",
        "finance_year_2_growth_rate": "year_2_growth_rate",
        "finance_year_3_growth_rate": "year_3_growth_rate",
        "finance_cost_per_unit": "cost_per_unit",
        "finance_opex_pct_revenue": "opex_pct_revenue",
        "finance_fixed_annual_overhead": "fixed_annual_overhead",
        "finance_starting_cash": "starting_cash",
    }
    if use_benchmarks and st.session_state.get("finance_benchmark_signature") != benchmark_signature:
        for widget_key, assumption_key in widget_defaults.items():
            st.session_state[widget_key] = benchmark[assumption_key]
        st.session_state["finance_pricing_period"] = benchmark.get("pricing_period", "Monthly" if business_model == "Subscription" else "One-Time")
        st.session_state["finance_benchmark_signature"] = benchmark_signature
    else:
        for widget_key, assumption_key in widget_defaults.items():
            if widget_key not in st.session_state:
                st.session_state[widget_key] = defaults.get(assumption_key, benchmark[assumption_key])
    if use_benchmarks:
        st.info("AI-generated estimate: TurboPitch has applied industry and business-model benchmarks. Review each input before presenting this model to investors.")
    st.markdown("### User-editable assumptions")
    left, right = st.columns(2)
    with left:
        pricing_period = st.selectbox("Pricing period", pricing_period_options, key="finance_pricing_period")
        price_labels = {
            "Monthly": "Monthly price per customer / location ($)",
            "Annual": "Annual price per customer / location ($)",
            "One-Time": "One-time price per customer / unit ($)",
        }
        price = st.number_input(price_labels[pricing_period], 0.0, step=1.0, key="finance_price_per_unit")
        units = st.number_input(f"Year 1 {revenue_driver.lower()}", 0.0, step=10.0, key="finance_year_1_units")
        cost_labels = {
            "Monthly": "Monthly cost per customer / location ($)",
            "Annual": "Annual cost per customer / location ($)",
            "One-Time": "One-time cost per customer / unit ($)",
        }
        cost = st.number_input(cost_labels[pricing_period], 0.0, step=1.0, key="finance_cost_per_unit")
        cash = st.number_input("Starting cash ($)", 0.0, step=10000.0, key="finance_starting_cash")
    with right:
        growth_y2 = st.slider("Year 2 growth rate", 0.0, 3.0, step=0.05, key="finance_year_2_growth_rate")
        growth_y3 = st.slider("Year 3 growth rate", 0.0, 3.0, step=0.05, key="finance_year_3_growth_rate")
        opex = st.slider("Operating expense % of revenue", 0.0, 1.0, step=0.01, key="finance_opex_pct_revenue")
        overhead = st.number_input("Fixed annual overhead ($)", 0.0, step=10000.0, key="finance_fixed_annual_overhead")
    assumptions = {"startup_idea": startup_idea, "industry": industry, "business_model": business_model, "revenue_driver": revenue_driver, "pricing_period": pricing_period, "price_per_unit": price, "year_1_units": units, "year_2_growth_rate": growth_y2, "year_3_growth_rate": growth_y3, "cost_per_unit": cost, "opex_pct_revenue": opex, "fixed_annual_overhead": overhead, "starting_cash": cash}
    st.session_state["finance_assumptions"] = assumptions
    st.session_state.update({"industry": industry if industry in INDUSTRY_BENCHMARKS else "SaaS", "assumption_mode": "Help Me Generate Them" if use_benchmarks else "Manual", "pricing_period": pricing_period, "price_per_unit": price, "year1_units": units, "growth_y2": growth_y2, "growth_y3": growth_y3, "cost_per_unit": cost, "opex_pct": opex, "fixed_overhead": overhead, "starting_cash": cash})
    summary = pd.DataFrame([{"Assumption": "Industry", "Value": industry}, {"Assumption": "Business Model", "Value": business_model}, {"Assumption": "Revenue Driver", "Value": revenue_driver}, {"Assumption": "Pricing Period", "Value": pricing_period}, {"Assumption": "Price per Customer / Unit", "Value": f"${price:,.2f}"}, {"Assumption": "Year 1 Customers / Locations / Units", "Value": f"{units:,.0f}"}, {"Assumption": "Year 2 Growth", "Value": f"{growth_y2:.0%}"}, {"Assumption": "Year 3 Growth", "Value": f"{growth_y3:.0%}"}, {"Assumption": "Cost per Customer / Unit", "Value": f"${cost:,.2f}"}, {"Assumption": "Operating Expense % of Revenue", "Value": f"{opex:.0%}"}, {"Assumption": "Fixed Annual Overhead", "Value": f"${overhead:,.0f}"}, {"Assumption": "Starting Cash", "Value": f"${cash:,.0f}"}])
    st.markdown("### Assumption Summary"); st.dataframe(summary, use_container_width=True, hide_index=True)
    projection = calculate_projection(assumptions); st.session_state["projection_df"] = projection
    display_columns = ["Year", "Customers / Units", "Revenue", "COGS", "Gross Profit", "Gross Margin", "Operating Expenses", "Net Income", "Ending Cash"]
    formats = {"Customers / Units": "{:,.0f}", "Revenue": "${:,.0f}", "COGS": "${:,.0f}", "Gross Profit": "${:,.0f}", "Gross Margin": "{:.1%}", "Operating Expenses": "${:,.0f}", "Net Income": "${:,.0f}", "Ending Cash": "${:,.0f}"}
    year_1, year_3 = projection.iloc[0], projection.iloc[-1]
    net_income_class = "kpi-green" if year_1["Net Income"] >= 0 else "kpi-red"
    ending_cash_class = "kpi-green" if year_3["Ending Cash"] >= 0 else "kpi-red"
    annualization_multiplier = 12 if pricing_period == "Monthly" else 1
    contribution_per_unit = (price * annualization_multiplier) - (cost * annualization_multiplier) - (price * annualization_multiplier * opex)
    break_even_units = overhead / contribution_per_unit if contribution_per_unit > 0 else None
    funding_need = max(0, -float(projection["Ending Cash"].min()))
    if projection.iloc[0]["Net Income"] < 0:
        funding_need = max(funding_need, abs(float(projection.iloc[0]["Net Income"])) + (overhead / 4))
    st.session_state["estimated_funding_need"] = funding_need
    st.session_state["break_even_units"] = break_even_units

    st.markdown("### Investor-facing output")
    kpi1, kpi2, kpi3, kpi4 = st.columns(4)
    with kpi1:
        st.markdown(f'<div class="kpi-card kpi-blue"><div class="kpi-title">Year 1 Revenue</div><div class="kpi-value">${year_1["Revenue"]:,.0f}</div></div>', unsafe_allow_html=True)
    with kpi2:
        st.markdown(f'<div class="kpi-card kpi-gold"><div class="kpi-title">Year 1 Gross Margin</div><div class="kpi-value">{year_1["Gross Margin"]:.1%}</div></div>', unsafe_allow_html=True)
    with kpi3:
        st.markdown(f'<div class="kpi-card {net_income_class}"><div class="kpi-title">Year 1 Net Income</div><div class="kpi-value">${year_1["Net Income"]:,.0f}</div></div>', unsafe_allow_html=True)
    with kpi4:
        funding_display = f"${funding_need:,.0f}" if funding_need else "Funded"
        st.markdown(f'<div class="kpi-card {ending_cash_class}"><div class="kpi-title">Estimated Funding Need</div><div class="kpi-value">{funding_display}</div></div>', unsafe_allow_html=True)
    finance_left, finance_right = st.columns(2)
    with finance_left:
        break_even_text = (
            f"{break_even_units:,.0f} annual {revenue_driver.lower()} needed to cover fixed overhead."
            if break_even_units is not None else
            "Current contribution margin does not cover fixed overhead; revise price, direct cost, or operating expense assumptions."
        )
        st.markdown(f'<div class="finance-card"><strong>Break-even estimate</strong><br>{break_even_text}</div>', unsafe_allow_html=True)
    with finance_right:
        st.markdown(f'<div class="finance-card"><strong>EBITDA / profitability view</strong><br>Year 1 EBITDA proxy: ${year_1["Operating Income"]:,.0f}. Year 3 EBITDA proxy: ${year_3["Operating Income"]:,.0f}.</div>', unsafe_allow_html=True)
    st.markdown("### Revenue, expense, and profitability forecast")
    st.dataframe(projection[display_columns].style.format(formats), use_container_width=True, hide_index=True)

    assessment_status = "supports a credible first-pass model" if year_3["Operating Income"] >= 0 and year_1["Gross Margin"] >= 0.4 else "requires validation before investor use"
    st.markdown("### AI Assessment")
    st.markdown(f'<div class="finance-card"><strong>AI-generated estimate</strong><br>The model {assessment_status}. Year 1 revenue is ${year_1["Revenue"]:,.0f}, gross margin is {year_1["Gross Margin"]:.1%}, and Year 3 operating income is ${year_3["Operating Income"]:,.0f}. Treat this as decision support and validate it with customer, pricing, and cost evidence.</div>', unsafe_allow_html=True)
    suggested_adjustment = (
        "Reduce Year 1 volume or increase starting cash to cover the modeled operating loss."
        if funding_need > 0 else
        "Validate the unit-volume ramp with bottom-up pipeline, conversion, and retention evidence."
    )
    st.markdown("### Suggested Adjustment")
    st.markdown(f'<div class="finance-card">{suggested_adjustment}</div>', unsafe_allow_html=True)

    template_path = os.path.join(os.path.dirname(__file__), "templates", "financial_model_professional.xlsx")
    if not os.path.exists(template_path):
        st.warning("The professional Excel template is unavailable. TurboPitch will generate a compatible fallback workbook instead.")
    if st.button("Generate Financial Model", type="primary", key="generate_financial_model"):
        st.session_state["financial_model_generated"] = True
        st.success("Financial model prepared. Download the Excel workbook from the Export step.")
    st.markdown("### Investor Pushback Scenario")
    pushback, notes_col = st.columns([1, 2])
    with pushback: reduction = st.number_input("Reduce revenue assumptions by %", 0, 75, int(st.session_state.get("pushback_pct", 25)), 5, key="finance_revenue_reduction_pct")
    with notes_col: notes = st.text_area("Investor feedback notes", value=st.session_state.get("investor_feedback", "Investor may view early revenue assumptions as aggressive. Apply investor pushback to test whether the business still looks fundable."), height=100, key="finance_investor_notes")
    st.session_state.pushback_pct, st.session_state.investor_feedback = reduction, notes
    if st.button("Apply Investor Pushback", key="apply_investor_haircut"):
        haircut = apply_investor_haircut(assumptions, reduction); haircut_df = calculate_projection(haircut)
        st.session_state["investor_haircut_assumptions"], st.session_state["investor_haircut_projection"] = haircut, haircut_df
        st.warning("Investor pushback scenario applied below."); st.dataframe(haircut_df[display_columns].style.format(formats), use_container_width=True, hide_index=True)
    return assumptions


if st.session_state.tp_step == 2:
    finance_assumptions = render_financial_assumption_step(st.session_state.get("idea", ""))
    projection_df = st.session_state.get("projection_df")

# ==================================================
# BUILD MODEL
# ==================================================
model_industry = classify_industry(st.session_state.idea, st.session_state.industry)
projection_df = st.session_state.get("projection_df")
if projection_df is None:
    projection_df = calculate_projection({
        "price_per_unit": st.session_state.price_per_unit,
        "year_1_units": st.session_state.year1_units,
        "year_2_growth_rate": st.session_state.growth_y2,
        "year_3_growth_rate": st.session_state.growth_y3,
        "cost_per_unit": st.session_state.cost_per_unit,
        "opex_pct_revenue": st.session_state.opex_pct,
        "fixed_annual_overhead": st.session_state.fixed_overhead,
        "starting_cash": st.session_state.starting_cash,
        "pricing_period": st.session_state.get("pricing_period", "One-Time"),
    })
st.session_state["projection_df"] = projection_df

pnl_df = build_pnl_view(projection_df)
display_pnl = build_display_pnl(pnl_df)

reality_engine_output = run_reality_engine(
    st.session_state.idea,
    model_industry,
    st.session_state.price_per_unit,
    st.session_state.year1_units,
    st.session_state.growth_y2,
    st.session_state.growth_y3,
    st.session_state.opex_pct,
    st.session_state.fixed_overhead,
    st.session_state.cost_per_unit,
    projection_df["Ending Cash"].iloc[-1],
)

scorecard = build_scorecard(
    st.session_state.idea,
    model_industry,
    projection_df,
    st.session_state.price_per_unit,
    st.session_state.cost_per_unit,
    st.session_state.year1_units,
    st.session_state.growth_y2,
    st.session_state.growth_y3,
    reality_engine_output,
)

warnings = run_rule_based_sanity_check(
    st.session_state.price_per_unit,
    st.session_state.year1_units,
    st.session_state.growth_y2,
    st.session_state.growth_y3,
    st.session_state.cost_per_unit,
    st.session_state.opex_pct,
    projection_df["Ending Cash"].iloc[-1],
)

warning_summary = build_warning_summary(warnings)

benchmark_feedback = build_benchmark_feedback(
    model_industry,
    st.session_state.price_per_unit,
    st.session_state.cost_per_unit,
    st.session_state.year1_units,
    st.session_state.growth_y2,
    st.session_state.growth_y3,
    st.session_state.opex_pct,
    st.session_state.fixed_overhead,
)

if st.session_state.tp_step < 3:
    st.info("Click Next to generate your financial dashboard, investor review, and export files.")
    st.divider()
    render_nav("bottom")
    st.stop()

# ==================================================
# TOP ACTION BUTTONS
# ==================================================

import math

def safe_float(value, default=0.0):
    try:
        if value is None or value == "":
            return default
        return float(value)
    except (TypeError, ValueError):
        return default


def get_industry_benchmarks(industry: str):
    """
    Simple benchmark ranges by industry.
    You can expand this over time.
    """
    industry = (industry or "").strip().lower()

    benchmark_map = {
        "saas": {
            "price_low": 15,
            "price_high": 99,
            "growth_low": 0.10,
            "growth_high": 0.30,
            "margin_low": 0.70,
            "margin_high": 0.90,
            "cogs_low": 0.10,
            "cogs_high": 0.30,
            "churn_low": 0.02,
            "churn_high": 0.08,
        },
        "e-commerce": {
            "price_low": 20,
            "price_high": 120,
            "growth_low": 0.05,
            "growth_high": 0.20,
            "margin_low": 0.30,
            "margin_high": 0.60,
            "cogs_low": 0.40,
            "cogs_high": 0.70,
            "churn_low": 0.05,
            "churn_high": 0.15,
        },
        "service": {
            "price_low": 100,
            "price_high": 5000,
            "growth_low": 0.05,
            "growth_high": 0.25,
            "margin_low": 0.40,
            "margin_high": 0.80,
            "cogs_low": 0.20,
            "cogs_high": 0.60,
            "churn_low": 0.01,
            "churn_high": 0.10,
        },
        "consulting": {
            "price_low": 500,
            "price_high": 10000,
            "growth_low": 0.05,
            "growth_high": 0.25,
            "margin_low": 0.50,
            "margin_high": 0.85,
            "cogs_low": 0.15,
            "cogs_high": 0.50,
            "churn_low": 0.01,
            "churn_high": 0.08,
        },
        "marketplace": {
            "price_low": 10,
            "price_high": 100,
            "growth_low": 0.08,
            "growth_high": 0.25,
            "margin_low": 0.50,
            "margin_high": 0.85,
            "cogs_low": 0.15,
            "cogs_high": 0.50,
            "churn_low": 0.03,
            "churn_high": 0.10,
        },
    }

    return benchmark_map.get(industry, {
        "price_low": 20,
        "price_high": 100,
        "growth_low": 0.05,
        "growth_high": 0.20,
        "margin_low": 0.40,
        "margin_high": 0.80,
        "cogs_low": 0.20,
        "cogs_high": 0.60,
        "churn_low": 0.02,
        "churn_high": 0.10,
    })


def clamp(value, low, high):
    return max(low, min(value, high))


def generate_ai_explanations(
    industry,
    business_model,
    startup_price,
    monthly_growth_rate,
    cogs_percent,
    churn_rate=None
):
    """
    Creates dynamic explanation cards and suggested fixes.
    Returns:
      explanations: list of dicts
      suggested_fixes: dict
      severity_score: int
    """

    benchmarks = get_industry_benchmarks(industry)

    price = safe_float(startup_price)
    growth = safe_float(monthly_growth_rate)
    cogs = safe_float(cogs_percent)
    churn = safe_float(churn_rate, 0.05)

    explanations = []
    suggested_fixes = {}
    severity_score = 0

    # ----- Pricing check -----
    if price > benchmarks["price_high"]:
        severity_score += 25
        suggested_price = round((benchmarks["price_low"] + benchmarks["price_high"]) / 2, 2)
        suggested_fixes["startup_price"] = suggested_price

        explanations.append({
            "title": "Pricing looks too high for this market",
            "severity": "High",
            "reason": (
                f"Your entered price is {format_currency_ui(price)}, which is above the typical range "
                f"for {industry} businesses ({format_currency_ui(benchmarks['price_low'])} - "
                f"{format_currency_ui(benchmarks['price_high'])})."
            ),
            "why_it_matters": (
                "If pricing is too high, early customers may hesitate to buy, which lowers conversion "
                "and makes revenue assumptions less believable."
            ),
            "suggestion": (
                f"A more realistic starting range may be {format_currency_ui(benchmarks['price_low'])} - "
                f"{format_currency_ui(benchmarks['price_high'])}, with a suggested test price of about "
                f"{format_currency_ui(suggested_price)}."
            )
        })

    elif price < benchmarks["price_low"]:
        severity_score += 10
        suggested_price = round((benchmarks["price_low"] + price) / 2, 2)
        suggested_fixes["startup_price"] = suggested_price

        explanations.append({
            "title": "Pricing may be too low",
            "severity": "Medium",
            "reason": (
                f"Your price of {format_currency_ui(price)} is below the usual range for {industry} "
                f"({format_currency_ui(benchmarks['price_low'])} - {format_currency_ui(benchmarks['price_high'])})."
            ),
            "why_it_matters": (
                "Underpricing can make growth look easier on paper, but it may hurt margins and make it harder "
                "to cover operating costs."
            ),
            "suggestion": (
                f"Consider testing a higher entry point. A more balanced starting point could be around "
                f"{format_currency_ui(suggested_price)}."
            )
        })

    # ----- Growth check -----
    if growth > benchmarks["growth_high"]:
        severity_score += 25
        suggested_growth = round((benchmarks["growth_low"] + benchmarks["growth_high"]) / 2, 4)
        suggested_fixes["monthly_growth_rate"] = suggested_growth

        explanations.append({
            "title": "Growth assumption looks aggressive",
            "severity": "High",
            "reason": (
                f"You entered a monthly growth rate of {growth*100:.1f}%, which is above a more typical range "
                f"for {industry} ({benchmarks['growth_low']*100:.0f}%–{benchmarks['growth_high']*100:.0f}% per month)."
            ),
            "why_it_matters": (
                "Aggressive growth assumptions can quickly inflate revenue projections and make the business "
                "look stronger than it may be in reality."
            ),
            "suggestion": (
                f"A more realistic planning assumption may be around {suggested_growth*100:.1f}% monthly growth."
            )
        })

    elif growth < 0:
        severity_score += 20
        suggested_fixes["monthly_growth_rate"] = 0.03

        explanations.append({
            "title": "Growth rate is invalid",
            "severity": "High",
            "reason": "Your growth rate is negative.",
            "why_it_matters": (
                "A negative growth assumption can distort the model unless you are intentionally modeling decline."
            ),
            "suggestion": "Reset growth to a modest positive assumption like 3% and test from there."
        })

    # ----- COGS check -----
    if cogs > benchmarks["cogs_high"]:
        severity_score += 20
        suggested_cogs = round((benchmarks["cogs_low"] + benchmarks["cogs_high"]) / 2, 4)
        suggested_fixes["cogs_percent"] = suggested_cogs

        explanations.append({
            "title": "Cost of goods looks high",
            "severity": "High",
            "reason": (
                f"Your cost of goods is {cogs*100:.1f}% of revenue, which is above the typical range "
                f"for {industry} ({benchmarks['cogs_low']*100:.0f}%–{benchmarks['cogs_high']*100:.0f}%)."
            ),
            "why_it_matters": (
                "High delivery costs reduce gross margin, which makes profitability harder even if revenue grows."
            ),
            "suggestion": (
                f"Try stress-testing the business with COGS closer to {suggested_cogs*100:.1f}%."
            )
        })

    if cogs <= 0 or cogs < benchmarks["cogs_low"] * 0.35:
        severity_score += 30
        suggested_cogs = round((benchmarks["cogs_low"] + benchmarks["cogs_high"]) / 2, 4)
        suggested_fixes["cogs_percent"] = suggested_cogs

        explanations.append({
            "title": "COGS assumption looks unrealistically low",
            "severity": "High",
            "reason": (
                f"Your cost of goods is {cogs*100:.1f}% of revenue, well below the typical range "
                f"for {industry} ({benchmarks['cogs_low']*100:.0f}% to {benchmarks['cogs_high']*100:.0f}%)."
            ),
            "why_it_matters": (
                "Very low COGS can make gross margin look artificially strong. Investors will ask whether "
                "hosting, support, fulfillment, labor, payment fees, returns, and delivery costs are fully included."
            ),
            "suggestion": (
                f"Stress-test the model with COGS closer to {suggested_cogs*100:.1f}% unless you can defend the lower cost with operating data."
            )
        })

    # ----- Margin inference check -----
    gross_margin = 1 - cogs
    if gross_margin > benchmarks["margin_high"] + 0.10 or gross_margin > 0.95:
        severity_score += 25

        explanations.append({
            "title": "Gross margin may be too high to defend",
            "severity": "High",
            "reason": (
                f"Your estimated gross margin is {gross_margin*100:.1f}%, above the typical "
                f"{industry} benchmark ({benchmarks['margin_low']*100:.0f}% to {benchmarks['margin_high']*100:.0f}%)."
            ),
            "why_it_matters": (
                "A very high margin can be a red flag if it comes from excluding real delivery costs rather than proven pricing power."
            ),
            "suggestion": (
                "Show evidence for the margin or rerun the projection with a higher COGS assumption."
            )
        })

    if gross_margin < benchmarks["margin_low"]:
        severity_score += 15

        explanations.append({
            "title": "Gross margin may be too weak",
            "severity": "Medium",
            "reason": (
                f"Your estimated gross margin is {gross_margin*100:.1f}%, below the lower end of the typical "
                f"{industry} benchmark ({benchmarks['margin_low']*100:.0f}%–{benchmarks['margin_high']*100:.0f}%)."
            ),
            "why_it_matters": (
                "A weak margin means a lot of revenue gets consumed before covering payroll, marketing, and overhead."
            ),
            "suggestion": (
                "Look at pricing, supplier costs, packaging, delivery, or service efficiency to improve margin."
            )
        })

    # ----- Churn check -----
    if churn > benchmarks["churn_high"]:
        severity_score += 15
        suggested_churn = round((benchmarks["churn_low"] + benchmarks["churn_high"]) / 2, 4)
        suggested_fixes["churn_rate"] = suggested_churn

        explanations.append({
            "title": "Churn assumption looks too high",
            "severity": "Medium",
            "reason": (
                f"Your churn rate is {churn*100:.1f}%, above the normal range for {industry} "
                f"({benchmarks['churn_low']*100:.0f}%–{benchmarks['churn_high']*100:.0f}%)."
            ),
            "why_it_matters": (
                "High churn means you lose customers faster, so growth becomes much harder to sustain."
            ),
            "suggestion": (
                f"Try modeling churn closer to {suggested_churn*100:.1f}% unless you have strong evidence otherwise."
            )
        })

    if not explanations:
        explanations.append({
            "title": "Your assumptions look reasonably aligned",
            "severity": "Low",
            "reason": "Your main inputs are generally within a realistic planning range.",
            "why_it_matters": (
                "That does not guarantee success, but it does mean your model starts from assumptions "
                "that are easier to defend."
            ),
            "suggestion": "You can still test conservative and aggressive scenarios to compare outcomes."
        })

    severity_score = min(severity_score, 100)
    return explanations, suggested_fixes, severity_score

import math

def safe_float(value, default=0.0):
    try:
        if value is None or value == "":
            return default
        return float(value)
    except (TypeError, ValueError):
        return default

def get_industry_benchmarks(industry: str):
    industry = (industry or "").strip().lower()
    benchmark_map = {
        "saas": {"price_low": 15, "price_high": 99, "growth_low": 0.10, "growth_high": 0.30, "margin_low": 0.70, "margin_high": 0.90, "cogs_low": 0.10, "cogs_high": 0.30, "churn_low": 0.02, "churn_high": 0.08},
        "e-commerce": {"price_low": 20, "price_high": 120, "growth_low": 0.05, "growth_high": 0.20, "margin_low": 0.30, "margin_high": 0.60, "cogs_low": 0.40, "cogs_high": 0.70, "churn_low": 0.05, "churn_high": 0.15},
        "service": {"price_low": 100, "price_high": 5000, "growth_low": 0.05, "growth_high": 0.25, "margin_low": 0.40, "margin_high": 0.80, "cogs_low": 0.20, "cogs_high": 0.60, "churn_low": 0.01, "churn_high": 0.10},
        "consulting": {"price_low": 500, "price_high": 10000, "growth_low": 0.05, "growth_high": 0.25, "margin_low": 0.50, "margin_high": 0.85, "cogs_low": 0.15, "cogs_high": 0.50, "churn_low": 0.01, "churn_high": 0.08},
        "marketplace": {"price_low": 10, "price_high": 100, "growth_low": 0.08, "growth_high": 0.25, "margin_low": 0.50, "margin_high": 0.85, "cogs_low": 0.15, "cogs_high": 0.50, "churn_low": 0.03, "churn_high": 0.10},
    }
    return benchmark_map.get(industry, {"price_low": 20, "price_high": 100, "growth_low": 0.05, "growth_high": 0.20, "margin_low": 0.40, "margin_high": 0.80, "cogs_low": 0.20, "cogs_high": 0.60, "churn_low": 0.02, "churn_high": 0.10})

def clamp(value, low, high):
    return max(low, min(value, high))

# ==================================================
# DASHBOARD
# ==================================================

# ==================================================
# STEP 3 — AI REVIEW
# ==================================================
if st.session_state.tp_step == 3:
    st.markdown("## AI Pitch")
    st.caption("Use the model and startup input to prepare the investor narrative and pressure-test its claims.")
    if st.button("Generate AI Pitch", type="primary", key="generate_ai_pitch"):
        with st.spinner("Preparing investor narrative..."):
            st.session_state.business_plan_output = generate_business_plan_and_deck(
                st.session_state.idea, st.session_state.industry,
                st.session_state.price_per_unit, st.session_state.year1_units,
                st.session_state.growth_y2, st.session_state.growth_y3,
                st.session_state.cost_per_unit, st.session_state.opex_pct,
                st.session_state.fixed_overhead, projection_df, reality_engine_output,
            )
        st.success("AI pitch and supporting investor materials generated.")

    if st.session_state.get("business_plan_output"):
        st.markdown("### Investor-facing pitch preview")
        for title, bullets in extract_pitch_deck_section(st.session_state.business_plan_output)[:3]:
            st.markdown(f"**{title}**")
            for bullet in bullets[:3]:
                st.write(f"• {bullet}")

    st.markdown("### Readiness Scorecard")
    for label, status in scorecard.items():
        st.write(f"{score_metric(status)} **{label}:** {status}")

    st.markdown("---")
    st.markdown("### Risk Summary")
    st.info(
        f"{warning_summary['icon']} {warning_summary['overall']}  |  "
        f"Red Flags: {warning_summary['red_count']}  |  "
        f"Watch Items: {warning_summary['amber_count']}"
    )

    st.markdown("---")
    st.markdown("### Reality Check")
    st.write(f"{score_metric(reality_engine_output['overall'])} **Overall:** {reality_engine_output['overall']}")
    st.write(reality_engine_output["summary"])
    for label, item in reality_engine_output["checks"].items():
        st.write(f"{reality_status_icon(item['status'])} **{label}:** {item['message']}")

    st.markdown("---")
    st.markdown("### Assumption Warnings")
    for warning in warnings:
        st.write(warning)

# ==================================================
# STEP 4 — GENERATE MATERIALS
# ==================================================
elif st.session_state.tp_step == 4:

    st.markdown("## Business Plan")
    st.caption("Generate and review the investor-ready business plan and pitch-deck narrative.")

    industry_value = st.session_state.get("industry", "saas")
    startup_price_value = st.session_state.get("price_per_unit", 29.0)
    price = safe_float(startup_price_value, 1)
    cost = safe_float(st.session_state.get("cost_per_unit", 0.2))
    cogs_percent_value = cost / price if price > 0 else 0.2
    annual_growth_y2 = safe_float(st.session_state.get("growth_y2", 0.10))
    monthly_growth_rate_value = (1 + annual_growth_y2) ** (1/12) - 1
    churn_rate_value = 0.05

    explanations, suggested_fixes, severity_score = generate_ai_explanations(
        industry=industry_value,
        business_model="general",
        startup_price=startup_price_value,
        monthly_growth_rate=monthly_growth_rate_value,
        cogs_percent=cogs_percent_value,
        churn_rate=churn_rate_value
    )

    st.session_state["ai_explanations"] = explanations
    st.session_state["suggested_fixes"] = suggested_fixes
    st.session_state["assumption_severity_score"] = severity_score

    col_top1, col_top2, col_top3, col_top4, _ = st.columns([1.2, 1.6, 1.4, 1.4, 2.8])

    with col_top1:
        if st.button("Get Investor Verdict", key="run_vc_sanity_top"):
            with st.spinner("Analyzing..."):
                sanity_text = run_ai_sanity_check(
                    st.session_state.idea, st.session_state.industry,
                    st.session_state.price_per_unit, st.session_state.year1_units,
                    st.session_state.growth_y2, st.session_state.growth_y3,
                    st.session_state.cost_per_unit, st.session_state.opex_pct,
                    st.session_state.fixed_overhead, projection_df, reality_engine_output,
                )
                st.session_state.sanity_output = sanity_text
            st.success("Investor verdict generated.")

    with col_top2:
        if st.button("Generate Investor Materials", key="generate_full_plan_top"):
            with st.spinner("Generating materials..."):
                plan_text = generate_business_plan_and_deck(
                    st.session_state.idea, st.session_state.industry,
                    st.session_state.price_per_unit, st.session_state.year1_units,
                    st.session_state.growth_y2, st.session_state.growth_y3,
                    st.session_state.cost_per_unit, st.session_state.opex_pct,
                    st.session_state.fixed_overhead, projection_df, reality_engine_output,
                )
                st.session_state.business_plan_output = plan_text
            st.success("Materials generated.")

    with col_top3:
        if st.button("Generate VC Questions", key="generate_vc_questions_top"):
            with st.spinner("Generating questions..."):
                interrogation_text = run_ai_investor_interrogation(
                    st.session_state.idea, st.session_state.industry,
                    st.session_state.price_per_unit, st.session_state.year1_units,
                    st.session_state.growth_y2, st.session_state.growth_y3,
                    st.session_state.cost_per_unit, st.session_state.opex_pct,
                    st.session_state.fixed_overhead, st.session_state.starting_cash,
                    projection_df, reality_engine_output,
                )
                st.session_state.interrogation_output = interrogation_text
            st.success("Questions generated.")

    with col_top4:
        if st.button("Build Founder Answers", key="build_founder_answers_top"):
            with st.spinner("Building answers..."):
                answer_text = run_ai_founder_answer_builder(
                    st.session_state.idea, st.session_state.industry,
                    st.session_state.price_per_unit, st.session_state.year1_units,
                    st.session_state.growth_y2, st.session_state.growth_y3,
                    st.session_state.cost_per_unit, st.session_state.opex_pct,
                    st.session_state.fixed_overhead, st.session_state.starting_cash,
                    projection_df, reality_engine_output,
                )
                st.session_state.answer_builder_output = answer_text
            st.success("Founder answers built.")

    st.markdown("---")
    st.subheader("AI Explanation Engine")

    issue_severities = {str(item.get("severity", "")).strip().lower() for item in explanations}
    if issue_severities & {"high", "red", "critical"}:
        st.error(f"Your assumptions need attention. Risk level: {severity_score}/100")
    elif issue_severities & {"medium", "yellow", "orange"}:
        st.warning(f"Some assumptions need investor validation. Risk level: {severity_score}/100")
    else:
        st.success(f"Your core assumptions look reasonably grounded. Risk level: {severity_score}/100")

    for i, item in enumerate(explanations, 1):
        with st.expander(f"{i}. {item['title']} ({item['severity']})", expanded=(i == 1)):
            st.markdown(f"**What we saw:** {escape_streamlit_markdown(item['reason'])}")
            st.markdown(f"**Why it matters:** {escape_streamlit_markdown(item['why_it_matters'])}")
            st.markdown(f"**Suggested adjustment:** {escape_streamlit_markdown(item['suggestion'])}")

    st.subheader("Optimize My Model")
    if suggested_fixes and len(suggested_fixes) > 0:
        label_map = {
            "startup_price": "Price per Unit",
            "monthly_growth_rate": "Monthly Growth Rate",
            "cogs_percent": "COGS as % of Revenue",
            "churn_rate": "Churn Rate"
        }
        preview_rows = []
        for key, value in suggested_fixes.items():
            if key == "startup_price":
                current_val = st.session_state.get("price_per_unit", "—")
            elif key == "monthly_growth_rate":
                current_val = monthly_growth_rate_value
            elif key == "cogs_percent":
                current_val = cogs_percent_value
            else:
                current_val = st.session_state.get(key, "—")

            if key == "startup_price":
                current_display = format_currency_ui(current_val)
                suggested_display = format_currency_ui(value)
            elif key in {"monthly_growth_rate", "cogs_percent", "churn_rate"}:
                current_display = f"{safe_float(current_val):.1%}"
                suggested_display = f"{safe_float(value):.1%}"
            else:
                current_display = str(current_val)
                suggested_display = str(value)

            preview_rows.append({
                "Field": label_map.get(key, key),
                "Current Value": current_display,
                "Suggested Value": suggested_display,
            })
        optimize_display = pd.DataFrame(preview_rows, columns=["Field", "Current Value", "Suggested Value"])
        st.dataframe(optimize_display, use_container_width=True, hide_index=True)
        if st.button("Apply Suggested Model Fixes"):
            if "startup_price" in suggested_fixes:
                st.session_state.price_per_unit = suggested_fixes["startup_price"]
            if "monthly_growth_rate" in suggested_fixes:
                st.session_state.growth_y2 = suggested_fixes["monthly_growth_rate"]
            if "cogs_percent" in suggested_fixes:
                new_price = safe_float(st.session_state.get("price_per_unit", 1), 1)
                st.session_state.cost_per_unit = round(new_price * suggested_fixes["cogs_percent"], 2)
            if "churn_rate" in suggested_fixes:
                st.session_state.churn_rate = suggested_fixes["churn_rate"]
            st.success("Suggested fixes applied.")
            st.rerun()
    else:
        st.info("No major fixes needed based on the current assumptions.")

    st.markdown("---")
    st.subheader("Financial Dashboard")

    year1_revenue = projection_df.loc[0, "Revenue"]
    year3_revenue = projection_df.loc[2, "Revenue"]
    year3_gross_profit = projection_df.loc[2, "Gross Profit"]
    year3_gross_margin = projection_df.loc[2, "Gross Margin %"]
    year3_net_income = projection_df.loc[2, "Net Income"]
    year3_cash = projection_df.loc[2, "Ending Cash"]

    net_income_class = "kpi-green" if year3_net_income >= 0 else "kpi-red"
    cash_class = "kpi-green" if year3_cash >= 0 else "kpi-red"

    kpi1, kpi2, kpi3 = st.columns(3)
    kpi4, kpi5, kpi6 = st.columns(3)

    with kpi1:
        st.markdown(f'<div class="kpi-card kpi-blue"><div class="kpi-title">Year 1 Revenue</div><div class="kpi-value">${year1_revenue:,.0f}</div></div>', unsafe_allow_html=True)
    with kpi2:
        st.markdown(f'<div class="kpi-card kpi-green"><div class="kpi-title">Year 3 Revenue</div><div class="kpi-value">${year3_revenue:,.0f}</div></div>', unsafe_allow_html=True)
    with kpi3:
        st.markdown(f'<div class="kpi-card kpi-blue"><div class="kpi-title">Gross Profit</div><div class="kpi-value">${year3_gross_profit:,.0f}</div></div>', unsafe_allow_html=True)
    with kpi4:
        st.markdown(f'<div class="kpi-card kpi-gold"><div class="kpi-title">Gross Margin</div><div class="kpi-value">{year3_gross_margin:.1%}</div></div>', unsafe_allow_html=True)
    with kpi5:
        st.markdown(f'<div class="kpi-card {net_income_class}"><div class="kpi-title">Net Income</div><div class="kpi-value">${year3_net_income:,.0f}</div></div>', unsafe_allow_html=True)
    with kpi6:
        st.markdown(f'<div class="kpi-card {cash_class}"><div class="kpi-title">Ending Cash</div><div class="kpi-value">${year3_cash:,.0f}</div></div>', unsafe_allow_html=True)

    chart_col1, chart_col2 = st.columns(2)
    with chart_col1:
        st.markdown("#### Revenue / Net Income / Ending Cash")
        dashboard_chart_df = projection_df.set_index("Year")[["Revenue", "Net Income", "Ending Cash"]]
        st.line_chart(dashboard_chart_df, use_container_width=True)
    with chart_col2:
        st.markdown("#### Revenue vs COGS vs Operating Expenses")
        compare_df = projection_df.set_index("Year")[["Revenue", "COGS", "Operating Expenses"]]
        st.bar_chart(compare_df, use_container_width=True)

    summary_col1, summary_col2 = st.columns([2, 1])
    with summary_col1:
        st.markdown("#### Profit & Loss Summary")
        st.dataframe(display_pnl, use_container_width=True, hide_index=True)
    with summary_col2:
        st.markdown("#### Snapshot")
        st.markdown(f"""
        <div class="snapshot-card">
            <p><strong>Industry:</strong> {st.session_state.industry}</p>
            <p><strong>Customer Segment:</strong> {detect_customer_segment(st.session_state.idea, st.session_state.industry)}</p>
            <p><strong>Overall Readiness:</strong> {scorecard['Overall Investor Readiness']}</p>
            <p><strong>Pricing Realism:</strong> {scorecard['Pricing Realism']}</p>
            <p><strong>Sales Volume:</strong> {scorecard['Sales Volume']}</p>
            <p><strong>Growth Assumptions:</strong> {scorecard['Growth Assumptions']}</p>
            <p><strong>Margin Quality:</strong> {scorecard['Margin Quality']}</p>
            <p><strong>Cash Viability:</strong> {scorecard['Cash Viability']}</p>
        </div>
        """, unsafe_allow_html=True)

    if st.session_state.get("assumption_helper_output"):
        st.markdown("#### Suggested Assumption Explanation")
        st.markdown(clean_ai_text(st.session_state.assumption_helper_output))

    tab1, tab2, tab3, tab4, tab5, tab6, tab7, tab8, tab9, tab10 = st.tabs([
        "Investor Review", "Benchmarking", "Investor Q&A", "Answer Builder",
        "Business Plan & Deck", "Financial Model", "Charts", "Downloads",
        "How TurboPitch Works", "Assumption Builder",
    ])

    with tab1:
        st.markdown("## Investor Readiness Scorecard")
        for label, status in scorecard.items():
            st.write(f"{score_metric(status)} **{label}:** {status}")
        st.markdown("---")
        st.markdown("### Risk Summary")
        st.info(f"{warning_summary['icon']} {warning_summary['overall']}  |  Red Flags: {warning_summary['red_count']}  |  Watch Items: {warning_summary['amber_count']}")
        st.markdown("---")
        st.markdown("### Reality Check")
        st.write(f"{score_metric(reality_engine_output['overall'])} **Overall:** {reality_engine_output['overall']}")
        st.write(reality_engine_output["summary"])
        for label, item in reality_engine_output["checks"].items():
            st.write(f"{reality_status_icon(item['status'])} **{label}:** {item['message']}")
        st.markdown("---")
        st.markdown("### Assumption Warnings")
        for warning in warnings:
            st.write(warning)
        st.markdown("---")
        st.markdown("### AI Investor Verdict")
        if st.button("Run VC Analysis", key="run_ai_sanity_tab"):
            with st.spinner("Reviewing assumptions..."):
                sanity_text = run_ai_sanity_check(
                    st.session_state.idea, st.session_state.industry,
                    st.session_state.price_per_unit, st.session_state.year1_units,
                    st.session_state.growth_y2, st.session_state.growth_y3,
                    st.session_state.cost_per_unit, st.session_state.opex_pct,
                    st.session_state.fixed_overhead, projection_df, reality_engine_output,
                )
                st.session_state.sanity_output = sanity_text
        if st.session_state.get("sanity_output"):
            st.markdown(clean_ai_text(st.session_state.sanity_output))
            st.markdown("---")
            render_ai_methodology_note()
        else:
            st.info("No investor review yet. Click 'Get Investor Verdict' or 'Run VC Analysis'.")

    with tab2:
        st.markdown("## Industry Benchmark Feedback")
        st.write(f"Selected industry: **{st.session_state.industry}**")
        for item in benchmark_feedback:
            if item == "":
                st.markdown("---")
            elif item == "Suggested Adjustments":
                st.markdown("### Suggested Adjustments")
            else:
                st.write(item)

    with tab3:
        st.markdown("## Investor Questions")
        if st.button("Generate VC Questions", key="generate_vc_questions_tab"):
            with st.spinner("Generating questions..."):
                interrogation_text = run_ai_investor_interrogation(
                    st.session_state.idea, st.session_state.industry,
                    st.session_state.price_per_unit, st.session_state.year1_units,
                    st.session_state.growth_y2, st.session_state.growth_y3,
                    st.session_state.cost_per_unit, st.session_state.opex_pct,
                    st.session_state.fixed_overhead, st.session_state.starting_cash,
                    projection_df, reality_engine_output,
                )
                st.session_state.interrogation_output = interrogation_text
        if st.session_state.get("interrogation_output"):
            st.markdown(clean_ai_text(st.session_state.interrogation_output))
        else:
            st.info("No questions yet. Click 'Generate VC Questions'.")

    with tab4:
        st.markdown("## Founder Answer Prep")
        if st.button("Build Founder Answers", key="build_founder_answers_tab"):
            with st.spinner("Building answers..."):
                answer_text = run_ai_founder_answer_builder(
                    st.session_state.idea, st.session_state.industry,
                    st.session_state.price_per_unit, st.session_state.year1_units,
                    st.session_state.growth_y2, st.session_state.growth_y3,
                    st.session_state.cost_per_unit, st.session_state.opex_pct,
                    st.session_state.fixed_overhead, st.session_state.starting_cash,
                    projection_df, reality_engine_output,
                )
                st.session_state.answer_builder_output = answer_text
        if st.session_state.get("answer_builder_output"):
            st.markdown(clean_ai_text(st.session_state.answer_builder_output))
        else:
            st.info("No answers yet. Click 'Build Founder Answers'.")

    with tab5:
        st.markdown("### Business Plan & Pitch Deck")
        if st.button("Build Business Plan + Deck", key="generate_plan_tab"):
            with st.spinner("Building materials..."):
                plan_text = generate_business_plan_and_deck(
                    st.session_state.idea, st.session_state.industry,
                    st.session_state.price_per_unit, st.session_state.year1_units,
                    st.session_state.growth_y2, st.session_state.growth_y3,
                    st.session_state.cost_per_unit, st.session_state.opex_pct,
                    st.session_state.fixed_overhead, projection_df, reality_engine_output,
                )
                st.session_state.business_plan_output = plan_text
        if st.session_state.get("business_plan_output"):
            st.markdown(clean_ai_text(st.session_state.business_plan_output))
        else:
            st.info("No business plan yet. Click 'Generate Investor Materials' or 'Build Business Plan + Deck'.")

    with tab6:
        st.dataframe(display_pnl, use_container_width=True, hide_index=True)

    with tab7:
        chart_df = projection_df.set_index("Year")[["Revenue", "Net Income", "Ending Cash"]]
        st.line_chart(chart_df, use_container_width=True)

    with tab8:
        st.markdown("### Export Your Materials")
        st.write("Download your files below.")

    with tab9:
        render_trust_center()

    with tab10:
        st.markdown("## Assumption Builder")
        if st.session_state.assumption_mode == "Help Me Generate Them":
            st.write(f"**Price per unit:** ${st.session_state.price_per_unit:,.2f}")
            st.write(f"**Year 1 units sold:** {st.session_state.year1_units:,}")
            st.write(f"**Year 2 growth:** {st.session_state.growth_y2:.0%}")
            st.write(f"**Year 3 growth:** {st.session_state.growth_y3:.0%}")
            st.write(f"**Cost per unit:** ${st.session_state.cost_per_unit:,.2f}")
            st.write(f"**Operating expense %:** {st.session_state.opex_pct:.0%}")
            st.write(f"**Fixed overhead:** ${st.session_state.fixed_overhead:,.0f}")
            st.write(f"**Starting cash:** ${st.session_state.starting_cash:,.0f}")
            st.write(f"**Detected segment:** {detect_customer_segment(st.session_state.idea, st.session_state.industry)}")
            st.markdown("---")
        for label, item in reality_engine_output["checks"].items():
            st.write(f"{reality_status_icon(item['status'])} **{label}:** {item['message']}")
        if st.session_state.get("assumption_helper_output"):
            st.markdown(clean_ai_text(st.session_state.assumption_helper_output))
        else:
            st.info("No suggested assumptions yet.")

# ==================================================
# STEP 5 — DOWNLOADS
# ==================================================
elif st.session_state.tp_step == 5:
    st.markdown("### Download Your Investor Materials")
    st.write("Everything is ready. Download your files below.")

# ==================================================
# DOWNLOAD MATERIALS
# ==================================================
plan_raw = st.session_state.get("business_plan_output", "")
business_plan_text = extract_business_plan_section(plan_raw)
deck_slides = extract_pitch_deck_section(plan_raw)
if is_empty_pitch_deck(deck_slides):
    deck_slides = build_standard_pitch_deck_slides(
        st.session_state.idea,
        st.session_state.industry,
        st.session_state.price_per_unit,
        st.session_state.year1_units,
        st.session_state.growth_y2,
        st.session_state.growth_y3,
        st.session_state.cost_per_unit,
        st.session_state.opex_pct,
        st.session_state.fixed_overhead,
        st.session_state.starting_cash,
        projection_df,
        reality_engine_output,
    )

def clean_company_name_candidate(value):
    if value is None:
        return ""
    text = re.sub(r"\s+", " ", str(value).strip())

    stop_phrases = [
        "Executive Summary", "Problem", "Solution", "Market Opportunity",
        "Product Overview", "Go-To-Market Strategy", "Financial Overview",
        "Projection Assumptions", "Funding Ask", "Key Risks", "Target Customer",
        "Revenue Model",
    ]
    for stop in stop_phrases:
        index = text.lower().find(stop.lower())
        if index > 0:
            text = text[:index].strip()

    return text.strip(" .:-|")


def resolve_business_plan_company_name(company_name=None, business_plan_text=None, startup_idea=None):
    """Use a supplied company name, or safely recover one from plan content."""
    placeholder_values = {
        "", "your company", "company", "the company", "startup", "startup company",
    }

    def valid_name(value):
        cleaned = clean_company_name_candidate(value)
        if not cleaned or cleaned.lower() in placeholder_values or len(cleaned) > 40:
            return ""
        return cleaned

    direct = valid_name(company_name)
    if direct:
        return direct

    for source in (business_plan_text, startup_idea):
        if not source:
            continue

        for line in str(source).splitlines():
            match = re.search(r"Business Plan for\s+(.+)$", line.strip())
            if match:
                candidate = valid_name(match.group(1))
                if candidate:
                    return candidate

        patterns = [
            r"\b([A-Z][A-Za-z0-9&.'\- ]{1,40})\s+is a vertical SaaS platform\b",
            r"\b([A-Z][A-Za-z0-9&.'\- ]{1,40})\s+is an AI-powered\b",
            r"\b([A-Z][A-Za-z0-9&.'\- ]{1,40})\s+seeks\s+\$",
        ]
        for pattern in patterns:
            match = re.search(pattern, str(source))
            if match:
                candidate = valid_name(match.group(1))
                if candidate:
                    return candidate

    return "the company"


def is_car_wash_business(text):
    if text is None:
        return False
    text = str(text).lower()
    return "car wash" in text or "car washes" in text or "wash operators" in text


doc = Document()
raw_company_name = next(
    (clean_doc_text(st.session_state.get(key)) for key in ("company_name", "startup_name", "business_name")
     if clean_doc_text(st.session_state.get(key))),
    None,
)
doc_idea = clean_doc_text(st.session_state.get("idea", ""))
resolved_company_name = resolve_business_plan_company_name(
    company_name=raw_company_name,
    business_plan_text=business_plan_text,
    startup_idea=doc_idea,
)
doc_company_name = resolved_company_name
doc_industry = clean_doc_text(st.session_state.get("industry", "")) or "target industry"
doc_customer_segment = detect_customer_segment(doc_idea, doc_industry)
doc_is_car_wash = is_car_wash_business(f"{business_plan_text}\n{doc_idea}\n{doc_industry}")
doc_navy = "17365D"
doc_light_gray = "D9E1F2"
docx_footer_text = "TurboPitch Business Plan | Confidential"
docx_investor_readiness_heading = "Investor Readiness Assessment"

for section in doc.sections:
    section.top_margin = Inches(0.7)
    section.bottom_margin = Inches(0.7)
    section.left_margin = Inches(0.7)
    section.right_margin = Inches(0.7)
    section.footer_distance = Inches(0.3)
    footer = section.footer.paragraphs[0]
    footer.alignment = WD_ALIGN_PARAGRAPH.CENTER
    footer.paragraph_format.space_before = DocxPt(0)
    footer.paragraph_format.space_after = DocxPt(0)
    footer.paragraph_format.line_spacing = 1.0
    # The cover is a first page; keep its presentation clean while retaining a
    # discreet confidentiality footer on every subsequent page.
    section.different_first_page_header_footer = True
    footer_run = footer.add_run(docx_footer_text)
    footer_run.font.name = "Aptos"
    footer_run.font.size = DocxPt(8)
    footer_run.font.color.rgb = DocxRGBColor(102, 102, 102)

normal_style = doc.styles["Normal"]
normal_style.font.name = "Aptos"
normal_style.font.size = DocxPt(10.5)
normal_style.paragraph_format.space_before = DocxPt(0)
normal_style.paragraph_format.space_after = DocxPt(4)
normal_style.paragraph_format.line_spacing = 1.0

heading_style = doc.styles["Heading 1"]
heading_style.font.name = "Aptos"
heading_style.font.size = DocxPt(15)
heading_style.font.bold = True
heading_style.font.color.rgb = DocxRGBColor(23, 54, 93)
heading_style.paragraph_format.space_before = DocxPt(8)
heading_style.paragraph_format.space_after = DocxPt(3)
heading_style.paragraph_format.keep_with_next = True

def add_doc_heading(text, level=1):
    paragraph = doc.add_paragraph(style="Heading 1" if level == 1 else None)
    paragraph.paragraph_format.space_before = DocxPt(8 if level == 1 else 6)
    paragraph.paragraph_format.space_after = DocxPt(3)
    paragraph.paragraph_format.keep_with_next = True
    run = paragraph.add_run(clean_doc_text(text))
    run.bold = True
    run.font.name = "Aptos"
    run.font.size = DocxPt(15 if level == 1 else 12)
    run.font.color.rgb = DocxRGBColor(23, 54, 93)
    return paragraph

def set_doc_cell_shading(cell, fill):
    cell_properties = cell._tc.get_or_add_tcPr()
    shading = OxmlElement("w:shd")
    shading.set(qn("w:fill"), fill)
    cell_properties.append(shading)

def set_doc_cell_text(cell, value, bold=False, alignment=None, color=None):
    cell.text = ""
    paragraph = cell.paragraphs[0]
    if alignment is not None:
        paragraph.alignment = alignment
    run = paragraph.add_run(clean_doc_text(value))
    run.bold = bold
    run.font.name = "Aptos"
    run.font.size = DocxPt(9.5)
    if color:
        run.font.color.rgb = DocxRGBColor(*color)

def set_doc_cell_margins(cell, top=40, start=70, bottom=40, end=70):
    """Keep DOCX table padding compact without compromising readability."""
    cell_properties = cell._tc.get_or_add_tcPr()
    margins = cell_properties.first_child_found_in("w:tcMar")
    if margins is None:
        margins = OxmlElement("w:tcMar")
        cell_properties.append(margins)
    for side, value in (("top", top), ("start", start), ("bottom", bottom), ("end", end)):
        node = margins.find(qn(f"w:{side}"))
        if node is None:
            node = OxmlElement(f"w:{side}")
            margins.append(node)
        node.set(qn("w:w"), str(value))
        node.set(qn("w:type"), "dxa")

def prevent_doc_row_split(row):
    """Keep an individual table row intact when Word lays out a page."""
    properties = row._tr.get_or_add_trPr()
    if properties.find(qn("w:cantSplit")) is None:
        properties.append(OxmlElement("w:cantSplit"))

def add_doc_table(headers, rows, financial=False, keep_rows_together=False):
    # Keep the immediately preceding heading or introductory paragraph with the
    # table where possible, without forcing a new page for oversized tables.
    if doc.paragraphs:
        doc.paragraphs[-1].paragraph_format.keep_with_next = True
    table = doc.add_table(rows=1, cols=len(headers))
    table.style = "Table Grid"
    table.autofit = True
    table_properties = table._tbl.tblPr
    borders = OxmlElement("w:tblBorders")
    for edge in ("top", "left", "bottom", "right", "insideH", "insideV"):
        border = OxmlElement(f"w:{edge}")
        border.set(qn("w:val"), "single")
        border.set(qn("w:sz"), "4")
        border.set(qn("w:color"), doc_light_gray)
        borders.append(border)
    table_properties.append(borders)
    for index, header in enumerate(headers):
        cell = table.rows[0].cells[index]
        set_doc_cell_margins(cell)
        set_doc_cell_shading(cell, doc_navy)
        set_doc_cell_text(
            cell, header, bold=True,
            alignment=WD_ALIGN_PARAGRAPH.CENTER if index else WD_ALIGN_PARAGRAPH.LEFT,
            color=(255, 255, 255),
        )
    # A header must not be stranded at the bottom of a page. Keep every row
    # together only for small, presentation-oriented tables.
    prevent_doc_row_split(table.rows[0])
    for row in rows:
        table_row = table.add_row()
        if keep_rows_together:
            prevent_doc_row_split(table_row)
        cells = table_row.cells
        for index, value in enumerate(row):
            set_doc_cell_margins(cells[index])
            alignment = WD_ALIGN_PARAGRAPH.RIGHT if financial and index else WD_ALIGN_PARAGRAPH.LEFT
            set_doc_cell_text(cells[index], value, bold=(financial and index == 0), alignment=alignment)
    return table

def add_doc_body_paragraph(text):
    """Add generated copy, turning inline dash lists into compact Word bullets."""
    parts = [part.strip() for part in re.split(r"\s+-\s+", clean_doc_text(text)) if part.strip()]
    if len(parts) <= 1:
        return doc.add_paragraph(clean_doc_text(text))

    intro = doc.add_paragraph(parts[0])
    for item in parts[1:]:
        bullet = doc.add_paragraph(item, style="List Bullet")
        bullet.paragraph_format.space_before = DocxPt(0)
        bullet.paragraph_format.space_after = DocxPt(2)
        bullet.paragraph_format.line_spacing = 1.0
    return intro

def add_doc_bullet(text):
    """Add a compact, native Word bullet rather than dash-prefixed body copy."""
    bullet = doc.add_paragraph(clean_doc_text(text), style="List Bullet")
    bullet.paragraph_format.space_before = DocxPt(0)
    bullet.paragraph_format.space_after = DocxPt(2)
    bullet.paragraph_format.line_spacing = 1.0
    return bullet

def add_docx_product_overview(features_text):
    """Render the product introduction as body copy and features as bullets."""
    intro = doc.add_paragraph(
        f"{resolved_company_name} features a user-friendly software platform that integrates seamlessly into existing "
        "car wash operations. Key functionalities include:"
    )
    intro.paragraph_format.keep_with_next = True
    for item in (
        "Demand prediction algorithms that analyze real-time data",
        "Staffing optimization tools that adjust labor needs dynamically",
        "Resource tracking for chemicals and water usage",
        "Automated pricing recommendations based on predictive analytics",
    ):
        add_doc_bullet(item)
    doc.add_paragraph(
        "The platform is designed for easy integration into existing operations, enabling rapid deployment and a clear path to customer value."
    )

def add_docx_solution():
    """Render the solution as an introductory paragraph followed by native bullets."""
    if doc_is_car_wash:
        intro = doc.add_paragraph(
            f"{resolved_company_name} addresses these challenges by offering a comprehensive AI-driven platform that:"
        )
        items = (
            "Predicts demand using historical and real-time data, enabling better staffing decisions.",
            "Tracks chemical and water usage to minimize waste and optimize resource allocation.",
            "Delivers automated pricing recommendations based on dynamic factors such as weather and local events.",
            "Enhances customer experience through targeted promotions and pricing.",
        )
    else:
        intro = doc.add_paragraph(
            f"{resolved_company_name} addresses these challenges through a focused {doc_industry} platform that:"
        )
        items = (
            "Connects the core workflow data needed for faster operating decisions.",
            "Identifies demand, capacity, and resource signals that affect day-to-day execution.",
            "Recommends practical actions based on current business conditions.",
            "Improves customer outcomes through more consistent and responsive operations.",
        )
    intro.paragraph_format.keep_with_next = True
    for item in items:
        add_doc_bullet(item)

def add_docx_go_to_market_strategy():
    """Render go-to-market copy as an introduction followed by native bullets."""
    if doc_is_car_wash:
        intro = doc.add_paragraph("Our strategy focuses on targeting independent car washes through:")
        items = (
            "Direct sales outreach to operators with tailored demonstrations of product benefits.",
            "Partnerships with car wash equipment suppliers to support bundled offerings.",
            "Digital marketing campaigns aimed at educating potential customers on AI-driven operating improvements.",
            "Participation in industry trade shows and conferences to build brand awareness.",
        )
    else:
        intro = doc.add_paragraph(
            f"Our strategy focuses on reaching {doc_customer_segment} customers through:"
        )
        items = (
            "Direct outreach with tailored demonstrations of measurable product benefits.",
            "Partnerships with relevant industry providers and implementation channels.",
            "Digital marketing that educates prospective customers about operating improvements.",
            "Industry events and communities that build credibility and early pipeline.",
        )
    intro.paragraph_format.keep_with_next = True
    for item in items:
        add_doc_bullet(item)

def compact_paragraph(paragraph, is_heading=False):
    fmt = paragraph.paragraph_format
    fmt.line_spacing = 1.0
    if is_page_break_paragraph(paragraph):
        fmt.space_before = DocxPt(0)
        fmt.space_after = DocxPt(0)
        return
    if is_heading:
        fmt.space_before = DocxPt(8)
        fmt.space_after = DocxPt(3)
        fmt.keep_with_next = True
    elif paragraph.style and paragraph.style.name == "List Bullet":
        fmt.space_before = DocxPt(0)
        fmt.space_after = DocxPt(2)
    else:
        fmt.space_before = DocxPt(0)
        fmt.space_after = DocxPt(4)

def compact_table(table):
    """Apply the standard compact treatment without forcing small-table sizing."""
    table.autofit = True
    for row in table.rows:
        row.height = None
        for cell in row.cells:
            for paragraph in cell.paragraphs:
                paragraph.paragraph_format.space_before = DocxPt(0)
                paragraph.paragraph_format.space_after = DocxPt(1)
                paragraph.paragraph_format.line_spacing = 1.0

def compact_small_table(table, font_size=8):
    """Make short presentation tables compact enough to stay on one page."""
    table.autofit = True
    for row in table.rows:
        row.height = None
        for cell in row.cells:
            for paragraph in cell.paragraphs:
                paragraph.paragraph_format.space_before = DocxPt(0)
                paragraph.paragraph_format.space_after = DocxPt(0)
                paragraph.paragraph_format.line_spacing = 1.0
                for run in paragraph.runs:
                    run.font.size = DocxPt(font_size)

def compact_investor_readiness_table(table, font_size=7.5):
    """Use a denser, readable treatment for the investor-readiness table."""
    table.autofit = True
    for row in table.rows:
        row.height = None
        for cell in row.cells:
            for paragraph in cell.paragraphs:
                paragraph.paragraph_format.space_before = DocxPt(0)
                paragraph.paragraph_format.space_after = DocxPt(0)
                paragraph.paragraph_format.line_spacing = 1.0
                for run in paragraph.runs:
                    run.font.size = DocxPt(font_size)

def prevent_table_row_splits(table):
    """Keep rows intact for tables small enough to fit naturally on a page."""
    for row in table.rows:
        prevent_doc_row_split(row)

def is_page_break_paragraph(paragraph):
    return any(
        br.get(qn("w:type")) == "page"
        for br in paragraph._p.findall(".//" + qn("w:br"))
    )

def remove_excess_blank_paragraphs(document):
    """Remove body-only blank paragraphs that can create unintended blank pages."""
    paragraphs = list(document.paragraphs)
    for index, paragraph in enumerate(paragraphs):
        if paragraph.text.strip() or is_page_break_paragraph(paragraph):
            continue
        previous = paragraphs[index - 1] if index else None
        following = paragraphs[index + 1] if index + 1 < len(paragraphs) else None
        if (
            previous is None
            or following is None
            or (previous and not previous.text.strip())
            or (following and not following.text.strip())
            or (following and is_page_break_paragraph(following))
            or (previous and is_page_break_paragraph(previous))
        ):
            paragraph._element.getparent().remove(paragraph._element)

def remove_blank_page_artifacts(document):
    """Final DOCX guard for generated blank paragraphs and duplicate page breaks."""
    remove_excess_blank_paragraphs(document)
    page_break_seen = False
    for paragraph in list(document.paragraphs):
        if not is_page_break_paragraph(paragraph):
            if paragraph.text.strip():
                page_break_seen = False
            continue
        if page_break_seen:
            paragraph._element.getparent().remove(paragraph._element)
        else:
            page_break_seen = True
    # Catch blanks exposed by removing duplicate page-break paragraphs.
    remove_excess_blank_paragraphs(document)

def remove_empty_paragraphs_before_page_break(document):
    """Avoid carrying empty body paragraphs into an intentional page break."""
    for paragraph in list(document.paragraphs):
        if not is_page_break_paragraph(paragraph):
            continue
        previous = paragraph._p.getprevious()
        while previous is not None and previous.tag == qn("w:p"):
            if "".join(previous.itertext()).strip():
                break
            prior = previous.getprevious()
            previous.getparent().remove(previous)
            previous = prior

def apply_compact_business_plan_layout(document):
    for section in document.sections:
        section.top_margin = Inches(0.7)
        section.bottom_margin = Inches(0.7)
        section.left_margin = Inches(0.7)
        section.right_margin = Inches(0.7)

    remove_blank_page_artifacts(document)
    for paragraph in document.paragraphs:
        style_name = paragraph.style.name if paragraph.style else ""
        is_heading = style_name.startswith("Heading")
        compact_paragraph(paragraph, is_heading=is_heading)
        for run in paragraph.runs:
            if is_heading:
                run.font.size = DocxPt(14)
                run.font.bold = True
            else:
                run.font.size = DocxPt(10)

    for table in document.tables:
        compact_table(table)

# Cover page
cover_title = doc.add_paragraph()
cover_title.alignment = WD_ALIGN_PARAGRAPH.CENTER
cover_title.paragraph_format.space_before = DocxPt(90)
title_run = cover_title.add_run("Investor Business Plan")
title_run.bold = True
title_run.font.name = "Aptos"
title_run.font.size = DocxPt(24)
title_run.font.color.rgb = DocxRGBColor(23, 54, 93)
cover_subtitle = doc.add_paragraph()
cover_subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
cover_subtitle.paragraph_format.space_after = DocxPt(16)
cover_subtitle_text = f"Business Plan for {doc_company_name}" if doc_company_name != "the company" else "Business Plan"
subtitle_run = cover_subtitle.add_run(cover_subtitle_text)
subtitle_run.font.name = "Aptos"
subtitle_run.font.size = DocxPt(14)
subtitle_run.font.color.rgb = DocxRGBColor(75, 75, 75)
cover_tagline = doc.add_paragraph("AI-powered operating intelligence for independent car wash operators.")
cover_tagline.alignment = WD_ALIGN_PARAGRAPH.CENTER
cover_tagline.paragraph_format.space_after = DocxPt(30)
cover_tagline.runs[0].italic = True
cover_date = doc.add_paragraph("Generated by TurboPitch")
cover_date.alignment = WD_ALIGN_PARAGRAPH.CENTER
cover_date.runs[0].font.size = DocxPt(10)
doc.add_page_break()

known_plan_headings = {
    "executive summary", "problem", "solution", "market opportunity", "business model",
    "competitive advantage", "go-to-market strategy", "financial highlights", "funding ask",
    "projection assumptions & investor interpretation", "product overview", "competitive landscape",
    "financial overview", "key risks", "pitch deck content",
}

def add_financial_overview():
    """Add a deterministic, properly structured financial-summary section."""
    year1, year2, year3 = projection_df.iloc[0], projection_df.iloc[1], projection_df.iloc[2]
    intro = doc.add_paragraph("Initial projections indicate:")
    intro.paragraph_format.keep_with_next = True
    bullets = (
        f"Year 1 Revenue: {format_currency_doc(year1['Revenue'])}",
        f"Year 1 Gross Margin: {format_percent_doc(year1['Gross Margin %'])}",
        f"Year 2 Revenue Growth Rate: {st.session_state.growth_y2:.0%}",
        f"Year 3 Revenue Growth Rate: {st.session_state.growth_y3:.0%}",
    )
    for text in bullets:
        bullet = doc.add_paragraph(text, style="List Bullet")
        bullet.paragraph_format.space_before = DocxPt(0)
        bullet.paragraph_format.space_after = DocxPt(2)
        bullet.paragraph_format.line_spacing = 1.0
        bullet.paragraph_format.keep_with_next = True
    profitability_year = "Year 2" if year2["Net Income"] >= 0 else "a later projected year"
    profitability = doc.add_paragraph(
        f"Despite initial operating losses, the model reaches profitability in {profitability_year}, "
        "driven by increased adoption and operating leverage."
    )
    profitability.paragraph_format.keep_together = True

def add_docx_executive_summary():
    doc.add_paragraph(
        f"{doc_company_name} is a vertical SaaS platform designed specifically for independent car wash operators. "
        "The platform uses AI-powered forecasting to help operators optimize staffing, resource usage, and pricing "
        "decisions based on external demand signals such as weather, traffic, and local events."
    )

def add_docx_market_opportunity():
    if doc_is_car_wash:
        text = (
            "The car wash market is large and fragmented, with many independent and regional operators still relying "
            "on manual scheduling, static pricing, and limited operational analytics. This creates an opportunity for "
            "vertical software that improves demand forecasting, resource planning, and operating decisions."
        )
    else:
        text = (
            "The target market is large and fragmented, with many independent operators still relying on manual "
            "workflows, static pricing, and limited operational analytics. This creates an opportunity for vertical "
            "software that improves forecasting, resource planning, and operating decisions."
        )
    doc.add_paragraph(text)

def normalize_docx_plan_text(text):
    """Remove unsupported canned market claims and correct generated grammar."""
    text = text.replace("Reality Engine Summary", docx_investor_readiness_heading)
    text = re.sub(
        r"investor\s+readiness\s+assessment",
        docx_investor_readiness_heading,
        text,
        flags=re.IGNORECASE,
    )
    text = text.replace("as indicated by Reality Engine concerns", "based on investor-readiness review")
    text = text.replace("innovative software platform", "vertical SaaS platform")
    text = text.replace("innovative SaaS platform", "vertical SaaS platform")
    text = re.sub(r"innovative\s+vertical\s+SaaS\s+platform", "vertical SaaS platform", text, flags=re.IGNORECASE)
    text = re.sub(r"revolutionizing", "modernizing", text, flags=re.IGNORECASE)
    text = re.sub(r"capture a substantial share", "build an initial position", text, flags=re.IGNORECASE)
    text = text.replace("an vertical SaaS platform", "a vertical SaaS platform")
    text = text.replace("multi-billion dollar market", "large and fragmented market")
    text = text.replace("Your Company", doc_company_name)
    text = re.sub(
        r"(?:The global\s+)?The target market is large and fragmented[^.]*\.",
        "The target market is large and fragmented, with many independent operators still relying on manual workflows, static pricing, and limited operational analytics. This creates an opportunity for vertical software that improves forecasting, resource planning, and operating decisions.",
        text,
        flags=re.IGNORECASE,
    )
    text = re.sub(
        r"(?:The )?car wash industry[^.]*\$\s*[\d.,]+[^.]*\.?(?:\s*)",
        "The target market is large and fragmented, with many independent operators still relying on manual workflows, static pricing, and limited operational analytics. ",
        text,
        flags=re.IGNORECASE,
    )
    text = re.sub(
        r"(?:(?:there (?:are|is)|with) )?over [\d,]+ independent car wash operators[^.]*\.?(?:\s*)",
        "",
        text,
        flags=re.IGNORECASE,
    )
    return clean_doc_text(text)

if business_plan_text:
    suppressed_plan_sections = {
        "business model", "competitive landscape", "executive summary", "financial overview", "market opportunity",
        "projection assumptions & investor interpretation", "funding ask", "key risks", "pitch deck content",
    }
    suppress_current_section = False
    current_plan_section = None
    heading_pattern = "|".join(re.escape(heading) for heading in sorted(known_plan_headings, key=len, reverse=True))
    # Generators may put a heading and its text on adjacent lines instead of
    # separating them with a blank line. Split at those heading lines before
    # normalizing whitespace so headings cannot run into body copy.
    plan_blocks = re.split(
        rf"\n\s*\n|\n(?=(?:\d+[.):\-]\s*)?(?:{heading_pattern})\s*:?\s*(?:\n|$))",
        business_plan_text,
        flags=re.IGNORECASE,
    )
    for paragraph in plan_blocks:
        text = normalize_docx_plan_text(paragraph)
        if not text:
            continue
        if re.match(rf"^(?:Investor )?Business Plan(?: for {re.escape(doc_company_name)})?$", text, re.IGNORECASE) or re.match(rf"^Business Plan for {re.escape(doc_company_name)}$", text, re.IGNORECASE):
            continue

        normalized_text = re.sub(r"^\d+[.):\-]\s*", "", text)
        heading_match = None
        for heading in sorted(known_plan_headings, key=len, reverse=True):
            match = re.match(rf"^{re.escape(heading)}\s*:?\s*(.*)$", normalized_text, re.IGNORECASE)
            if match:
                heading_match = (heading, match.group(1).strip())
                break
        if heading_match:
            heading_key, remainder = heading_match
            suppress_current_section = heading_key in suppressed_plan_sections
            current_plan_section = heading_key
            if heading_key == "executive summary":
                add_doc_heading("Executive Summary")
                add_docx_executive_summary()
            elif heading_key == "market opportunity":
                add_doc_heading("Market Opportunity")
                add_docx_market_opportunity()
            elif heading_key == "financial overview":
                add_doc_heading("Financial Overview")
                add_financial_overview()
            elif heading_key == "solution":
                add_doc_heading("Solution")
                add_docx_solution()
            elif heading_key == "go-to-market strategy":
                add_doc_heading("Go-To-Market Strategy")
                add_docx_go_to_market_strategy()
            if not suppress_current_section:
                if heading_key in {"solution", "go-to-market strategy"}:
                    continue
                add_doc_heading(heading_key.title() if heading_key != "go-to-market strategy" else "Go-To-Market Strategy")
                if remainder:
                    if heading_key == "product overview":
                        add_docx_product_overview(remainder)
                    else:
                        add_doc_body_paragraph(remainder)
        elif not suppress_current_section and current_plan_section not in {"solution", "go-to-market strategy"}:
            if current_plan_section == "product overview":
                add_docx_product_overview(text)
            else:
                add_doc_body_paragraph(text)
else:
    doc.add_paragraph("No business plan available yet. Generate the full business plan and deck first.")

year1, year2, year3 = projection_df.iloc[0], projection_df.iloc[1], projection_df.iloc[2]
profitability_year = next(
    (f"Year {index}" for index, net_income in enumerate(projection_df["Net Income"], start=1) if net_income >= 0),
    "Year 2",
)
add_doc_heading("Projection Assumptions & Investor Interpretation")
projection_bullets = [
    f"Revenue build: Revenue is projected at {format_currency_doc(year1['Revenue'])} in Year 1, {format_currency_doc(year2['Revenue'])} in Year 2, and {format_currency_doc(year3['Revenue'])} in Year 3.",
    f"Customer/unit assumptions: The Year 1 target of {float(year1['Units']):,.0f} units is aggressive and should be tied to customer count, locations served, transaction volume, conversion rate, and sales capacity.",
    f"Pricing logic: The {format_price_doc(st.session_state.price_per_unit)} per-unit assumption may understate value and should be validated against customer willingness to pay.",
    f"Gross margin logic: Projected gross margin of {format_percent_doc(year1['Gross Margin %'])} is directionally attractive for a software model, but COGS assumptions should be validated.",
    f"Operating expense logic: Operating expenses equal {year1['Operating Expenses'] / year1['Revenue']:.0%} of Year 1 revenue, which may create investor concern unless supported by a clear hiring and spend plan.",
    "Cash runway: Ending cash balances suggest no immediate funding gap, but runway should be tied to milestone execution.",
    "Funding need vs. milestones: Funding should be linked to product development, pilot customers, conversion, and go-to-market execution.",
    "Likely investor pushback: Investors will likely challenge Year 1 adoption, pricing clarity, CAC, sales cycle, and operating expense discipline.",
]
for bullet in projection_bullets:
    add_doc_bullet(bullet)

add_doc_heading("Funding Ask")
doc_funding_ask = st.session_state.get("funding_ask", 1_000_000)
doc.add_paragraph(
    f"{resolved_company_name} is seeking {format_funding_ask_doc(doc_funding_ask)} in seed funding to support product development, pilot customer acquisition, data integrations, and go-to-market execution. Based on the current model, {resolved_company_name} reaches profitability in Year 2, but investors will expect the funding plan to be tied to measurable milestones including pilot launches, paid customer conversion, retention, and sales capacity."
)

add_doc_heading("Key Risks")
for risk in (
    "Market adoption: Independent operators may adopt new software more slowly than projected.",
    "Competitive pressure: Existing software providers or POS vendors may add similar optimization features.",
    f"Pricing risk: The initial {format_price_doc(st.session_state.price_per_unit)} per-unit assumption may not reflect the value delivered or may require a clearer pricing unit.",
    "Execution risk: Scaling customer acquisition, onboarding, and support may require more resources than initially modeled.",
    "Technical risk: Forecasting accuracy depends on the quality of data integrations and continuous model improvement.",
):
    add_doc_bullet(risk)

add_doc_heading("Target Customer")
if doc_is_car_wash:
    doc.add_paragraph(
        f"{doc_company_name} targets independent and regional car wash operators, especially businesses with 1 to 20 locations that rely on manual scheduling, static pricing, and limited operational analytics."
    )
else:
    doc.add_paragraph(
        f"{doc_company_name} targets the selected customer segment, especially customers that rely on manual workflows, fragmented tools, static pricing, and limited operational analytics."
    )
revenue_model_heading = add_doc_heading("Revenue Model")
revenue_model_heading.paragraph_format.keep_with_next = True
revenue_model_intro = doc.add_paragraph(
    f"{doc_company_name} currently models revenue using a {format_price_doc(st.session_state.price_per_unit)} per-unit {st.session_state.get('pricing_period', 'One-Time').lower()} pricing assumption. Before investor presentation, {doc_company_name} should translate this assumption into customer economics, including whether the unit represents a paid wash transaction, managed location, active customer account, or vehicle processed. A tiered SaaS pricing model may better reflect value across single-location and multi-location operators."
)
revenue_model_intro.paragraph_format.keep_with_next = True
revenue_model_table = add_doc_table(
    ["Plan", "Target Customer", "Example Pricing Logic"],
    [
        ["Starter", "Single-location operator", "Entry monthly SaaS plan"],
        ["Growth", "Multi-location operator", "Higher monthly plan with forecasting and analytics"],
        ["Regional / Enterprise", "Larger operator groups", "Custom pricing and integrations"],
    ],
    keep_rows_together=True,
)
add_doc_heading("Customer Economics")
doc.add_paragraph(
    f"Investors will expect {doc_company_name} to validate customer acquisition cost, sales cycle, pilot-to-paid conversion rate, retention, gross margin, payback period, and revenue per customer, location, or account. The {float(year1['Units']):,.0f}-unit Year 1 assumption should be translated into a clear operating build, including expected customer count, locations served, transaction volume, and sales capacity."
)
add_doc_heading("Market Sizing Framework")
market_sizing_table = add_doc_table(
    ["Metric", "Definition", "Validation Needed"],
    (
        [
            ["TAM", "Total addressable market for car wash operating software", "Validate with industry location count and software spend assumptions"],
            ["SAM", "Independent and regional car wash operators reachable by initial go-to-market", "Validate target operator count and location profile"],
            ["SOM", "Realistic 3-5 year capture based on sales capacity and adoption", "Validate with pipeline, conversion rate, and churn assumptions"],
        ] if doc_is_car_wash else [
            ["TAM", "Total addressable market for the selected vertical software category", "Validate with customer count and software spend assumptions"],
            ["SAM", "Target customer segment reachable by the initial go-to-market motion", "Validate target customer count and customer profile"],
            ["SOM", "Realistic 3-5 year capture based on sales capacity and adoption", "Validate with pipeline, conversion rate, and churn assumptions"],
        ]
    ),
    keep_rows_together=True,
)
add_doc_heading("Competitive Landscape")
doc.add_paragraph(
    f"{doc_company_name} competes against manual workflows, spreadsheets, incumbent operational systems, and generic business intelligence tools. Its differentiation should be demonstrated through a focused vertical workflow and measurable customer outcomes."
)
add_doc_table(
    ["Alternative", "Limitation", f"{doc_company_name} Advantage"],
    [
        ["Manual scheduling / spreadsheets", "Reactive and labor intensive", "Predictive staffing and demand planning"],
        ["POS dashboards", "Shows historical data but limited forecasting", "Forward-looking demand and pricing recommendations"],
        ["Generic BI tools", "Not built for the target workflow", "Vertical SaaS workflow and operator-specific insights"],
        ["Legacy management software", "Operational system of record but often static", "AI-driven optimization layer"],
    ],
)
add_doc_heading("Use of Funds")
use_of_funds_table = add_doc_table(
    ["Use of Funds", "Allocation"],
    [["Product Development", "35%"], ["Sales & Marketing", "30%"], ["Customer Success / Onboarding", "15%"], ["Data Integrations", "10%"], ["Working Capital", "10%"]],
    keep_rows_together=True,
)
milestones_heading = add_doc_heading("Milestones")
milestones_heading.paragraph_format.keep_with_next = True
milestones_table = add_doc_table(
    ["Milestone", "Target"],
    [
        ["Product Launch / MVP", "3–6 months"],
        ["Pilot Operators", "10–25 operators"],
        ["Paid Conversion", "Convert pilots to paid accounts"],
        ["Data Integrations", "Weather, traffic, POS, usage data"],
        ["Investor Metrics", "Retention, CAC, payback, margin"],
    ],
    keep_rows_together=True,
)

# Let the financial section follow Milestones naturally. A forced page break
# here can strand the last Milestones row on an otherwise blank page.
add_doc_heading("Financial Projections")
financial_rows = []
currency_line_items = {"Revenue", "COGS", "Gross Profit", "Operating Expenses", "Operating Income", "Taxes", "Net Income", "Ending Cash"}
for _, row in pnl_df.iterrows():
    line_item = clean_doc_text(row.get("Line Item", ""))
    values = [line_item]
    for column in pnl_df.columns[1:]:
        value = row[column]
        if line_item == "Gross Margin %":
            values.append(format_percent_doc(value))
        elif line_item == "Units":
            values.append(f"{float(value):,.0f}")
        elif line_item in currency_line_items:
            values.append(format_currency_doc(value))
        else:
            values.append(clean_doc_text(value))
    financial_rows.append(values)
add_doc_table([clean_doc_text(column) for column in pnl_df.columns], financial_rows, financial=True)

methodology_heading = add_doc_heading("Methodology Note")
methodology_note = doc.add_paragraph(
    "TurboPitch combines founder inputs, internal benchmark ranges, rule-based business heuristics, "
    "investor-readiness review, and structured financial modeling logic to create investor-readiness analysis and starter assumptions."
)
methodology_note.paragraph_format.space_before = DocxPt(4)
methodology_note.paragraph_format.space_after = DocxPt(2)
methodology_note.paragraph_format.line_spacing = 1.0
methodology_note.paragraph_format.keep_with_next = True
for run in methodology_note.runs:
    run.font.size = DocxPt(9)

investor_readiness_heading = add_doc_heading(docx_investor_readiness_heading)
investor_readiness_heading.paragraph_format.keep_with_next = True
investor_readiness_table = add_doc_table(
    ["Assessment Area", "Status", "Investor View", "Suggested Action"],
    [
        ["Customer Segment Detection", "Green", "Segment is identifiable, but buyer profile needs narrowing.", "Define buyer title, customer size, and beachhead."],
        ["Pricing Market Fit", "Yellow", "Pricing may understate value and needs validation.", "Test willingness to pay and tiered pricing."],
        ["Adoption Realism", "Red", "Year 1 adoption is the largest credibility risk.", "Tie units to pipeline, conversion, sales capacity, and churn."],
        ["Growth Realism", "Green", "Growth path is reasonable if early adoption is proven.", "Support growth with funnel metrics and retention."],
        ["Operating Model Reality", "Orange", "OpEx may pressure early profitability.", "Break out overhead, sales, product, and support spend."],
        ["Financial Assumption Reality", "Green", "Gross margin and runway appear workable.", "Validate COGS, runway, and working capital needs."],
    ],
)

if st.session_state.get("assumption_helper_output"):
    add_doc_heading("Suggested Assumption Rationale")
    for paragraph in clean_ai_text(st.session_state.assumption_helper_output).split("\n\n"):
        paragraph = clean_doc_text(paragraph)
        if paragraph:
            doc.add_paragraph(paragraph)

doc_buffer = io.BytesIO()
apply_compact_business_plan_layout(doc)

# Keep only the short presentation tables compact and prevent their rows from
# splitting across pages. The financial table intentionally remains excluded.
for small_table, font_size in (
    (revenue_model_table, 8.5),
    (market_sizing_table, 8.5),
    (use_of_funds_table, 8),
    (milestones_table, 8),
):
    compact_small_table(small_table, font_size=font_size)
    prevent_table_row_splits(small_table)

# This table follows the financial projections, so it gets a more compact
# treatment to keep the entire assessment intentional on the same page.
compact_investor_readiness_table(investor_readiness_table, font_size=7.5)
prevent_table_row_splits(investor_readiness_table)

# The global compaction pass normalizes body copy, so restore the deliberately
# compact methodology note afterwards.
methodology_note.paragraph_format.space_before = DocxPt(4)
methodology_note.paragraph_format.space_after = DocxPt(2)
methodology_note.paragraph_format.line_spacing = 1.0
methodology_note.paragraph_format.keep_with_next = True
for run in methodology_note.runs:
    run.font.size = DocxPt(9)

remove_empty_paragraphs_before_page_break(doc)
remove_blank_page_artifacts(doc)

# Restore the intentionally distinct cover after the global compacting pass.
cover_title.paragraph_format.space_before = DocxPt(90)
cover_title.paragraph_format.space_after = DocxPt(4)
title_run.font.size = DocxPt(24)
cover_subtitle.paragraph_format.space_after = DocxPt(10)
subtitle_run.font.size = DocxPt(14)
cover_tagline.paragraph_format.space_after = DocxPt(18)
cover_date.paragraph_format.space_after = DocxPt(0)
cover_date.runs[0].font.size = DocxPt(10)
doc.save(doc_buffer)
doc_buffer.seek(0)

wb = Workbook()

# The export is intentionally built in code. The optional template is a visual
# reference only, so downloads remain available if it is not packaged at runtime.
navy_fill = PatternFill("solid", fgColor="17365D")
header_fill = PatternFill("solid", fgColor="1F4E78")
section_fill = PatternFill("solid", fgColor="D9EAF7")
input_fill = PatternFill("solid", fgColor="FFF2CC")
status_fills = {
    "green": PatternFill("solid", fgColor="C6EFCE"),
    "yellow": PatternFill("solid", fgColor="FFEB9C"),
    "orange": PatternFill("solid", fgColor="FCE4D6"),
    "red": PatternFill("solid", fgColor="FFC7CE"),
    "action": PatternFill("solid", fgColor="D9EAF7"),
}
header_font = Font(color="FFFFFF", bold=True)
title_font = Font(size=18, bold=True, color="FFFFFF")
section_font = Font(size=12, bold=True, color="1F1F1F")
hardcode_font = Font(color="0000FF")
thin_border = Border(*(Side(style="thin", color="B7C9D6") for _ in range(4)))
currency_format = '$#,##0;[Red]($#,##0);-'
currency_per_unit_format = '$#,##0.00 " / unit";[Red]($#,##0.00 " / unit");-'
percent_format = '0.0%;[Red](0.0%);-'
units_format = '#,##0;[Red](#,##0);-'


def excel_status_fill(status):
    return status_fills.get(str(status).strip().lower(), status_fills["action"])


def apply_excel_status_style(cell, status):
    """Apply the shared professional status treatment to a status cell only."""
    status_key = str(status).strip().lower()
    status_fonts = {
        "green": "006100", "yellow": "7F6000", "orange": "C65911",
        "red": "9C0006", "action": "1F4E78",
    }
    cell.fill = excel_status_fill(status_key)
    cell.font = Font(name="Aptos", size=10, bold=True,
                     color=status_fonts.get(status_key, "1F4E78"))


def format_table(ws, start_row, end_row, start_col, end_col):
    for row in ws.iter_rows(min_row=start_row, max_row=end_row, min_col=start_col, max_col=end_col):
        for cell in row:
            cell.border = thin_border
            cell.alignment = Alignment(vertical="top", wrap_text=True)


def section_title(ws, row, title, end_col=4):
    ws.merge_cells(start_row=row, start_column=1, end_row=row, end_column=end_col)
    cell = ws.cell(row=row, column=1, value=clean_excel_text(title))
    cell.fill = navy_fill
    cell.font = header_font
    cell.alignment = Alignment(vertical="center")


# ---------------- Dashboard Sheet ----------------
ws_dash = wb.active
ws_dash.title = "Dashboard"
ws_dash.sheet_view.showGridLines = False
ws_dash.merge_cells("A1:L1")
ws_dash["A1"] = "TurboPitch Financial Dashboard"
ws_dash["A1"].fill = navy_fill
ws_dash["A1"].font = title_font
ws_dash["A1"].alignment = Alignment(horizontal="center", vertical="center")
ws_dash.merge_cells("A2:L2")
ws_dash["A2"] = "Investor-ready projection summary with model-linked KPIs and readiness flags."
ws_dash["A2"].font = Font(size=10, italic=True, color="666666")
ws_dash["A2"].alignment = Alignment(horizontal="center")
ws_dash.row_dimensions[1].height = 28

section_title(ws_dash, 4, "Key Metrics", 4)
for col, value in enumerate(["Metric", "Year 1", "Year 2", "Year 3"], 1):
    cell = ws_dash.cell(5, col, value)
    cell.fill, cell.font, cell.border = section_fill, Font(bold=True), thin_border
    cell.alignment = Alignment(horizontal="center")
dashboard_rows = [("Revenue", 8, currency_format), ("Gross Margin %", 15, percent_format),
                  ("Net Income", 27, currency_format), ("Ending Cash", 29, currency_format)]
for row, (label, model_row, number_format) in enumerate(dashboard_rows, 6):
    ws_dash.cell(row, 1, label).font = Font(bold=True)
    for col, model_col in enumerate(("B", "C", "D"), 2):
        cell = ws_dash.cell(row, col, f"='Financial Model'!{model_col}{model_row}")
        cell.number_format = number_format
format_table(ws_dash, 5, 9, 1, 4)

# Keep this header separate from A4:D4.  Excel repairs overlapping merged
# ranges, so dashboard sections must never share any cells.
ws_dash.merge_cells("F4:H4")
ws_dash["F4"] = "Investor Snapshot"
ws_dash["F4"].fill = navy_fill
ws_dash["F4"].font = header_font
ws_dash["F4"].alignment = Alignment(vertical="center")
for col, value in enumerate(["Assessment", "Status", "Notes"], 6):
    cell = ws_dash.cell(5, col, value)
    cell.fill, cell.font, cell.border = section_fill, Font(bold=True), thin_border
    cell.alignment = Alignment(horizontal="center")
snapshot_items = [("Overall Readiness", scorecard["Overall Investor Readiness"], reality_engine_output["summary"]),
                  ("Pricing Realism", scorecard["Pricing Realism"], reality_engine_output["checks"]["Pricing Market Fit"]["message"]),
                  ("Adoption Realism", scorecard["Sales Volume"], reality_engine_output["checks"]["Adoption Realism"]["message"]),
                  ("Growth Realism", scorecard["Growth Assumptions"], reality_engine_output["checks"]["Growth Realism"]["message"]),
                  ("Cash Viability", scorecard["Cash Viability"], reality_engine_output["checks"]["Financial Assumption Reality"]["message"])]
for row, (assessment, status, notes) in enumerate(snapshot_items, 6):
    ws_dash.cell(row, 6, assessment).font = Font(bold=True)
    ws_dash.cell(row, 7, status).fill = excel_status_fill(status)
    ws_dash.cell(row, 8, notes)
format_table(ws_dash, 5, 10, 6, 8)

section_title(ws_dash, 12, "P&L Summary", 4)
for col, value in enumerate(["Line Item", "Year 1", "Year 2", "Year 3"], 1):
    cell = ws_dash.cell(13, col, value)
    cell.fill, cell.font, cell.border = section_fill, Font(bold=True), thin_border
    cell.alignment = Alignment(horizontal="center")
pnl_dashboard = [("Revenue", 8, currency_format), ("Gross Profit", 14, currency_format),
                 ("Operating Income", 22, currency_format), ("Net Income", 27, currency_format),
                 ("Ending Cash", 29, currency_format)]
for row, (label, model_row, number_format) in enumerate(pnl_dashboard, 14):
    ws_dash.cell(row, 1, label).font = Font(bold=True)
    for col, model_col in enumerate(("B", "C", "D"), 2):
        cell = ws_dash.cell(row, col, f"='Financial Model'!{model_col}{model_row}")
        cell.number_format = number_format
format_table(ws_dash, 13, 18, 1, 4)
# Keep this header separate from A12:D12 for the same merged-cell integrity
# requirement as the Investor Snapshot section above.
ws_dash.merge_cells("F12:L12")
ws_dash["F12"] = "Projection Chart"
ws_dash["F12"].fill = navy_fill
ws_dash["F12"].font = header_font
ws_dash["F12"].alignment = Alignment(vertical="center")
projection_chart_stream = create_excel_projection_chart_image(projection_df)
projection_chart_img = XLImage(io.BytesIO(projection_chart_stream.getvalue()))
projection_chart_img.width, projection_chart_img.height = 680, 320
ws_dash.add_image(projection_chart_img, "F13")
for col, width in {"A": 24, "B": 16, "C": 16, "D": 16, "E": 3, "F": 23, "G": 18, "H": 58, "I": 12, "J": 12, "K": 12, "L": 12}.items():
    ws_dash.column_dimensions[col].width = width
ws_dash.freeze_panes = "A5"

# ---------------- Assumptions Sheet ----------------
ws_assump = wb.create_sheet("Assumptions")
ws_assump.sheet_view.showGridLines = False
ws_assump.merge_cells("A1:D1")
ws_assump["A1"] = "TurboPitch Model Assumptions"
ws_assump["A1"].fill, ws_assump["A1"].font = navy_fill, title_font
ws_assump["A1"].alignment = Alignment(horizontal="center")
for col, value in enumerate(["Assumption", "Value", "Type", "Notes"], 1):
    cell = ws_assump.cell(3, col, value)
    cell.fill, cell.font, cell.border = header_fill, header_font, thin_border
assumptions_data = [
    ("Startup idea / company name", st.session_state.idea or "No startup idea provided.", "Hardcoded", "Founder input"),
    ("Industry", st.session_state.industry, "Hardcoded", "Founder input"),
    ("Customer segment", detect_customer_segment(st.session_state.idea, st.session_state.industry), "Calculated", "Detected from founder input"),
    ("Assumption mode", st.session_state.get("assumption_mode", "Manual"), "Hardcoded", "Model setup"),
    ("Model version", "Professional", "Calculated", "Export profile"),
    ("Price per unit", st.session_state.price_per_unit, "Hardcoded", "Revenue driver"),
    ("Year 1 units", st.session_state.year1_units, "Hardcoded", "Initial sales volume"),
    ("Year 2 growth", st.session_state.growth_y2, "Hardcoded", "Unit growth rate"),
    ("Year 3 growth", st.session_state.growth_y3, "Hardcoded", "Unit growth rate"),
    ("COGS per unit", st.session_state.cost_per_unit, "Hardcoded", "Direct cost assumption"),
    ("Gross margin", (st.session_state.price_per_unit - st.session_state.cost_per_unit) / st.session_state.price_per_unit if st.session_state.price_per_unit else 0, "Calculated", "Derived from price and COGS"),
    ("Operating expense %", st.session_state.opex_pct, "Hardcoded", "Variable operating expense ratio"),
    ("Fixed annual overhead", st.session_state.fixed_overhead, "Hardcoded", "Annual fixed operating costs"),
    ("Tax rate", 0.21, "Hardcoded", "Applied only to positive operating income"),
    ("Starting cash", st.session_state.starting_cash, "Hardcoded", "Opening cash balance"),
    ("Funding ask", st.session_state.get("funding_ask", "Not provided"), "Hardcoded", "Included when captured by the app"),
    ("Investor pushback %", st.session_state.pushback_pct / 100, "Hardcoded", "Revenue reduction for the conservative investor case"),
    ("Pricing period", st.session_state.get("pricing_period", "One-Time"), "Hardcoded", "Determines whether price is annualized"),
]
for row, (label, value, kind, notes) in enumerate(assumptions_data, 4):
    for col, item in enumerate((label, value, kind, notes), 1):
        cell = ws_assump.cell(row, col, clean_excel_text(item) if isinstance(item, str) else item)
        cell.border, cell.alignment = thin_border, Alignment(vertical="top", wrap_text=True)
    ws_assump.cell(row, 1).font = Font(bold=True)
    if kind == "Hardcoded":
        ws_assump.cell(row, 2).font, ws_assump.cell(row, 2).fill = hardcode_font, input_fill
    if label == "Price per unit": ws_assump.cell(row, 2).number_format = currency_format
    elif label == "COGS per unit": ws_assump.cell(row, 2).number_format = currency_per_unit_format
    elif label in {"Year 1 units"}: ws_assump.cell(row, 2).number_format = units_format
    elif label in {"Year 2 growth", "Year 3 growth", "Gross margin", "Operating expense %", "Tax rate", "Investor pushback %"}: ws_assump.cell(row, 2).number_format = percent_format
    elif label in {"Fixed annual overhead", "Starting cash", "Funding ask"} and isinstance(value, (int, float)): ws_assump.cell(row, 2).number_format = currency_format
for col, width in {"A": 28, "B": 35, "C": 16, "D": 42}.items(): ws_assump.column_dimensions[col].width = width
ws_assump.freeze_panes = "A4"

# ---------------- Financial Model Sheet ----------------
ws_model = wb.create_sheet("Financial Model")
ws_model.sheet_view.showGridLines = False
ws_model.merge_cells("A1:D1")
ws_model["A1"] = "TurboPitch Financial Model"
ws_model["A1"].fill, ws_model["A1"].font = navy_fill, title_font
ws_model["A1"].alignment = Alignment(horizontal="center")
for col, value in enumerate(["Line Item", "Year 1", "Year 2", "Year 3"], 1):
    cell = ws_model.cell(3, col, value)
    cell.fill, cell.font, cell.border = header_fill, header_font, thin_border
    cell.alignment = Alignment(horizontal="center")
model_rows = [
    (4, "Revenue Build", None), (6, "Units", ("='Assumptions'!B10", "=B6*(1+'Assumptions'!B11)", "=C6*(1+'Assumptions'!B12)")),
    (7, "Average Selling Price", ("='Assumptions'!B9",) * 3), (8, "Revenue", ("=B6*B7*IF('Assumptions'!B21=\"Monthly\",12,1)", "=C6*C7*IF('Assumptions'!B21=\"Monthly\",12,1)", "=D6*D7*IF('Assumptions'!B21=\"Monthly\",12,1)")),
    (9, "Revenue Growth", ("=0", "=C8/B8-1", "=D8/C8-1")), (11, "Cost of Sales", None),
    (12, "Cost per Unit", ("='Assumptions'!B13",) * 3), (13, "COGS", ("=B6*B12", "=C6*C12", "=D6*D12")),
    (14, "Gross Profit", ("=B8-B13", "=C8-C13", "=D8-D13")), (15, "Gross Margin %", ("=IFERROR(B14/B8,0)", "=IFERROR(C14/C8,0)", "=IFERROR(D14/D8,0)")),
    (17, "Operating Expenses", None), (18, "Variable OpEx", ("=B8*'Assumptions'!B14", "=C8*'Assumptions'!B14", "=D8*'Assumptions'!B14")),
    (19, "Fixed Annual Overhead", ("='Assumptions'!B15",) * 3), (20, "Total Operating Expenses", ("=B18+B19", "=C18+C19", "=D18+D19")),
    (21, "OpEx % of Revenue", ("=IFERROR(B20/B8,0)", "=IFERROR(C20/C8,0)", "=IFERROR(D20/D8,0)")),
    (22, "Operating Income", ("=B14-B20", "=C14-C20", "=D14-D20")), (24, "Cash Flow / Profitability", None),
    (25, "Beginning Cash", ("='Assumptions'!B16", "=B29", "=C29")), (26, "Taxes", ("=MAX(0,B22*'Assumptions'!B17)", "=MAX(0,C22*'Assumptions'!B17)", "=MAX(0,D22*'Assumptions'!B17)")),
    (27, "Net Income", ("=B22-B26", "=C22-C26", "=D22-D26")), (28, "Net Margin %", ("=IFERROR(B27/B8,0)", "=IFERROR(C27/C8,0)", "=IFERROR(D27/D8,0)")), (29, "Ending Cash", ("=B25+B27", "=C25+C27", "=D25+D27")), (34, "Investor Pushback Case", None),
    (35, "Revenue @ Investor Pushback", ("=B8*(1-'Assumptions'!B18)", "=C8*(1-'Assumptions'!B18)", "=D8*(1-'Assumptions'!B18)")),
    (36, "Net Income @ Investor Pushback", ("=B35-(B13*(1-'Assumptions'!B18))-(B18*(1-'Assumptions'!B18))-B19-MAX(0,(B35-(B13*(1-'Assumptions'!B18))-(B18*(1-'Assumptions'!B18))-B19)*'Assumptions'!B17)", "=C35-(C13*(1-'Assumptions'!B18))-(C18*(1-'Assumptions'!B18))-C19-MAX(0,(C35-(C13*(1-'Assumptions'!B18))-(C18*(1-'Assumptions'!B18))-C19)*'Assumptions'!B17)", "=D35-(D13*(1-'Assumptions'!B18))-(D18*(1-'Assumptions'!B18))-D19-MAX(0,(D35-(D13*(1-'Assumptions'!B18))-(D18*(1-'Assumptions'!B18))-D19)*'Assumptions'!B17)")),
    (37, "Ending Cash @ Investor Pushback", ("='Assumptions'!B16+B36", "=B37+C36", "=C37+D36")),
]
percent_rows, unit_rows = {9, 15, 21, 28}, {6}
for row, label, formulas in model_rows:
    ws_model.cell(row, 1, label)
    if formulas is None:
        for col in range(1, 5):
            cell = ws_model.cell(row, col)
            cell.fill, cell.font = section_fill, Font(bold=True)
    else:
        ws_model.cell(row, 1).font = Font(bold=True)
        for col, formula in enumerate(formulas, 2): ws_model.cell(row, col, formula)
        for col in range(2, 5):
            ws_model.cell(row, col).number_format = percent_format if row in percent_rows else units_format if row in unit_rows else currency_format
format_table(ws_model, 3, 37, 1, 4)
for col, width in {"A": 34, "B": 17, "C": 17, "D": 17}.items(): ws_model.column_dimensions[col].width = width
ws_model.freeze_panes = "B4"

# ---------------- Investor Review Sheet ----------------
ws_review = wb.create_sheet("Investor Review")
ws_review.sheet_view.showGridLines = False
ws_review.merge_cells("A1:C1")
ws_review["A1"] = "Investor Review"
ws_review["A1"].fill, ws_review["A1"].font = navy_fill, title_font
for col, value in enumerate(["Section", "Assessment", "Notes"], 1):
    cell = ws_review.cell(3, col, value); cell.fill, cell.font, cell.border = header_fill, header_font, thin_border
credible = [item["message"] for item in reality_engine_output["checks"].values() if "green" in item["status"].lower()]
challenges = [item["message"] for item in reality_engine_output["checks"].values() if "green" not in item["status"].lower()]
review_rows = [("Overall readiness", reality_engine_output["overall"], reality_engine_output["summary"]),
               ("What looks credible", "Positive", "\n".join(credible) or "No checks are currently marked positive."),
               ("What investors may challenge", "Review", "\n".join(challenges) or "No material challenges identified by the Reality Engine."),
               ("Recommended fixes", "Action", "\n".join(item for item in benchmark_feedback if item and item != "Suggested Adjustments"))]
for row, values in enumerate(review_rows, 4):
    for col, value in enumerate(values, 1): ws_review.cell(row, col, clean_excel_text(value))
    ws_review.cell(row, 1).font = Font(bold=True); ws_review.cell(row, 2).fill = excel_status_fill(values[1])
format_table(ws_review, 3, 7, 1, 3)
for row in range(4, 8): ws_review.row_dimensions[row].height = 75
for col, width in {"A": 30, "B": 18, "C": 105}.items(): ws_review.column_dimensions[col].width = width
ws_review.freeze_panes = "A4"

# ---------------- Benchmark Feedback Sheet ----------------
ws_bench = wb.create_sheet("Benchmark Feedback")
ws_bench.sheet_view.showGridLines = False
ws_bench.merge_cells("A1:C1")
ws_bench["A1"] = "Industry Benchmark Feedback"
ws_bench["A1"].fill, ws_bench["A1"].font = navy_fill, title_font
for col, value in enumerate(["Category", "Status", "Feedback"], 1):
    cell = ws_bench.cell(3, col, value); cell.fill, cell.font, cell.border = header_fill, header_font, thin_border
categories = ["Pricing realism", "Adoption realism", "Growth realism", "Margin quality", "Cash viability", "Operating model realism"]
checks = reality_engine_output["checks"]
benchmark_rows = [
    (categories[0], checks["Pricing Market Fit"]), (categories[1], checks["Adoption Realism"]),
    (categories[2], checks["Growth Realism"]),
    (categories[3], {"status": scorecard["Margin Quality"], "message": checks["Financial Assumption Reality"]["message"]}),
    (categories[4], {"status": scorecard["Cash Viability"], "message": checks["Financial Assumption Reality"]["message"]}),
    (categories[5], checks["Operating Model Reality"]),
]
for row, (category, item) in enumerate(benchmark_rows, 4):
    ws_bench.cell(row, 1, clean_excel_text(category)); ws_bench.cell(row, 1).font = Font(bold=True)
    ws_bench.cell(row, 2, clean_excel_text(item["status"])); ws_bench.cell(row, 2).fill = excel_status_fill(item["status"])
    ws_bench.cell(row, 3, clean_excel_text(item["message"]))
format_table(ws_bench, 3, 9, 1, 3)
for col, width in {"A": 28, "B": 18, "C": 105}.items(): ws_bench.column_dimensions[col].width = width
ws_bench.freeze_panes = "A4"

# ---------------- Reality Engine Sheet ----------------
ws_reality = wb.create_sheet("Reality Engine")
ws_reality.sheet_view.showGridLines = False
ws_reality.merge_cells("A1:D1")
ws_reality["A1"] = "Reality Engine Review"
ws_reality["A1"].fill, ws_reality["A1"].font = navy_fill, title_font
for col, value in enumerate(["Check", "Result", "Status", "Notes"], 1):
    cell = ws_reality.cell(3, col, value); cell.fill, cell.font, cell.border = header_fill, header_font, thin_border
for row, (label, item) in enumerate(checks.items(), 4):
    ws_reality.cell(row, 1, clean_excel_text(label)); ws_reality.cell(row, 1).font = Font(bold=True)
    ws_reality.cell(row, 2, clean_excel_text(item["status"])); ws_reality.cell(row, 3, clean_excel_text(item["status"])); ws_reality.cell(row, 3).fill = excel_status_fill(item["status"])
    ws_reality.cell(row, 4, clean_excel_text(item["message"]))
format_table(ws_reality, 3, 3 + len(checks), 1, 4)
for col, width in {"A": 30, "B": 20, "C": 18, "D": 105}.items(): ws_reality.column_dimensions[col].width = width
ws_reality.freeze_panes = "A4"

# ---------------- Methodology Sheet ----------------
ws_method = wb.create_sheet("Methodology")
ws_method.sheet_view.showGridLines = False
ws_method.merge_cells("A1:B1")
ws_method["A1"] = "TurboPitch Methodology"
ws_method["A1"].fill, ws_method["A1"].font = navy_fill, title_font
for col, value in enumerate(["Topic", "Methodology"], 1):
    cell = ws_method.cell(3, col, value); cell.fill, cell.font, cell.border = header_fill, header_font, thin_border
methodology_rows = [
    ("Revenue calculation", "Units multiplied by average selling price for each projected year."),
    ("COGS calculation", "Units multiplied by cost per unit."),
    ("Gross margin", "Gross profit divided by revenue; gross profit equals revenue less COGS."),
    ("Operating expenses", "Variable operating expense equals revenue times the operating-expense percentage; fixed annual overhead is added."),
    ("Taxes", "Taxes apply at the tax rate only when operating income is positive."),
    ("Net income", "Operating income less taxes."),
    ("Ending cash", "Beginning cash plus net income; each following year begins with the prior year's ending cash."),
    ("Investor pushback case", "Revenue, COGS, and variable OpEx flex with the conservative investor case; fixed overhead remains fixed, taxes apply only to positive operating income, and net income and cash are recalculated."),
    ("Reality Engine interpretation", "Green indicates a favorable signal, yellow indicates an assumption to support, and red indicates material investor risk. Outputs are decision support, not validation."),
]
for row, values in enumerate(methodology_rows, 4):
    for col, value in enumerate(values, 1): ws_method.cell(row, col, clean_excel_text(value))
    ws_method.cell(row, 1).font = Font(bold=True)
format_table(ws_method, 3, 12, 1, 2)
for col, width in {"A": 28, "B": 110}.items(): ws_method.column_dimensions[col].width = width
ws_method.freeze_panes = "A4"

# ---------------- Save Workbook ----------------
def populate_professional_excel_template(template_path):
    """Populate the supplied finance template without changing its structure or formulas."""
    template_wb = load_workbook(template_path)
    assumptions = template_wb["Assumptions"]
    dashboard = template_wb["Dashboard"]
    model = template_wb["Financial Model"]
    review = template_wb["Investor Review"]
    benchmark = template_wb["Benchmark Feedback"]
    reality = template_wb["Reality Engine"]

    export_industry = classify_industry(st.session_state.idea, st.session_state.industry)
    customer_segment = detect_customer_segment(st.session_state.idea, export_industry)
    pricing_period = st.session_state.get("pricing_period", "One-Time")
    revenue_multiplier = 12 if pricing_period == "Monthly" else 1
    assumption_values = {
        "B5": clean_excel_text(st.session_state.idea or "No startup idea provided."),
        "B6": clean_excel_text(export_industry),
        "B7": clean_excel_text(customer_segment),
        "B8": clean_excel_text(st.session_state.get("assumption_mode", "Manual")),
        "B9": st.session_state.price_per_unit,
        "B10": st.session_state.year1_units,
        "B11": st.session_state.growth_y2,
        "B12": st.session_state.growth_y3,
        "B13": st.session_state.cost_per_unit,
        "B14": st.session_state.opex_pct,
        "B15": st.session_state.fixed_overhead,
        "B16": st.session_state.starting_cash,
        "B17": 0.21,
        "B18": st.session_state.pushback_pct / 100,
        "B19": clean_excel_text(st.session_state.get("investor_feedback", "Not provided")),
    }
    for cell, value in assumption_values.items():
        assumptions[cell] = value

    # Explicitly reset numeric formats because template input styles can carry
    # a prior percentage or currency format into a newly populated value.
    assumption_number_formats = {
        "B9": currency_format, "B10": units_format,
        "B11": percent_format, "B12": percent_format,
        "B13": currency_per_unit_format, "B14": percent_format,
        "B15": currency_format, "B16": currency_format,
        "B17": percent_format, "B18": percent_format,
    }
    for cell, number_format in assumption_number_formats.items():
        assumptions[cell].number_format = number_format
    assumption_units = {
        "C9": "$ / unit", "C10": "units", "C11": "%", "C12": "%",
        "C13": "$ / unit", "C14": "% of revenue", "C15": "$ / year",
        "C16": "$", "C17": "%", "C18": "% pushback",
    }
    for cell, unit in assumption_units.items():
        assumptions[cell] = unit
    assumptions["A18"] = "Investor Pushback %"

    # The template has no dedicated funding row. Use its blank next row only when
    # the app has a funding ask, retaining the template's existing row styling.
    funding_ask = st.session_state.get("funding_ask")
    if funding_ask not in (None, "", "Not provided"):
        from copy import copy
        for column in range(1, 6):
            source = assumptions.cell(19, column)
            target = assumptions.cell(20, column)
            target._style = copy(source._style)
            target.alignment = copy(source.alignment)
            target.number_format = source.number_format
        assumptions["A20"] = clean_excel_text("Funding Ask")
        assumptions["B20"] = clean_excel_text(funding_ask) if isinstance(funding_ask, str) else funding_ask
        assumptions["C20"] = clean_excel_text("$")
        assumptions["D20"] = clean_excel_text("Founder input")
        assumptions["E20"] = clean_excel_text("Funding requested, if captured by the app.")
        if isinstance(funding_ask, (int, float)):
            assumptions["B20"].number_format = currency_format

    checks = reality_engine_output["checks"]
    benchmark_rows, suggested_fix = build_benchmark_feedback_rows(
        export_industry, st.session_state.price_per_unit,
        st.session_state.cost_per_unit, st.session_state.year1_units,
        st.session_state.growth_y2, st.session_state.growth_y3,
        st.session_state.opex_pct, st.session_state.fixed_overhead,
    )
    # The Year 1 volume benchmark must use the same status as the adoption
    # reality check so Dashboard, Investor Review, Benchmark Feedback, and
    # Reality Engine tell one coherent investor story.
    benchmark_rows[4]["status"] = checks["Adoption Realism"]["status"]
    # Normalize the professional workbook's reporting headers.
    for cell, value in {
        "F4": "Investor Snapshot", "G4": "Assessment", "H4": "Status", "I4": "Notes",
    }.items():
        dashboard[cell] = clean_excel_text(value)
    for cell, value in {
        "A5": "Category", "B5": "Metric", "C5": "Current Value", "D5": "Investor Read", "E5": "Recommended Fix",
    }.items():
        review[cell] = clean_excel_text(value)
    for cell, value in {
        "A5": "Category", "B5": "Status", "C5": "Feedback", "D5": "Benchmark Basis",
    }.items():
        benchmark[cell] = clean_excel_text(value)

    # Restore the template's three-year calculation grid.  Assumptions B9:B18
    # are the canonical input locations; calculated formulas remain black and
    # same-workbook links remain green.
    model_formulas = {
        "B6": "='Assumptions'!B10", "C6": "=B6*(1+'Assumptions'!B11)", "D6": "=C6*(1+'Assumptions'!B12)",
        "B7": "='Assumptions'!B9", "C7": "='Assumptions'!B9", "D7": "='Assumptions'!B9",
        "B8": f"=B6*B7*{revenue_multiplier}", "C8": f"=C6*C7*{revenue_multiplier}", "D8": f"=D6*D7*{revenue_multiplier}",
        "B9": "=0", "C9": "=IFERROR(C8/B8-1,0)", "D9": "=IFERROR(D8/C8-1,0)",
        "B12": "='Assumptions'!B13", "C12": "='Assumptions'!B13", "D12": "='Assumptions'!B13",
        "B13": "=B6*B12", "C13": "=C6*C12", "D13": "=D6*D12",
        "B14": "=B8-B13", "C14": "=C8-C13", "D14": "=D8-D13",
        "B15": "=IFERROR(B14/B8,0)", "C15": "=IFERROR(C14/C8,0)", "D15": "=IFERROR(D14/D8,0)",
        "B18": "=B8*'Assumptions'!B14", "C18": "=C8*'Assumptions'!B14", "D18": "=D8*'Assumptions'!B14",
        "B19": "='Assumptions'!B15", "C19": "='Assumptions'!B15", "D19": "='Assumptions'!B15",
        "B20": "=B18+B19", "C20": "=C18+C19", "D20": "=D18+D19",
        "B21": "=IFERROR(B20/B8,0)", "C21": "=IFERROR(C20/C8,0)", "D21": "=IFERROR(D20/D8,0)",
        "B22": "=B14-B20", "C22": "=C14-C20", "D22": "=D14-D20",
        "B25": "='Assumptions'!B16", "C25": "=B29", "D25": "=C29",
        "B26": "=MAX(0,B22*'Assumptions'!B17)", "C26": "=MAX(0,C22*'Assumptions'!B17)", "D26": "=MAX(0,D22*'Assumptions'!B17)",
        "B27": "=B22-B26", "C27": "=C22-C26", "D27": "=D22-D26",
        "B28": "=IFERROR(B27/B8,0)", "C28": "=IFERROR(C27/C8,0)", "D28": "=IFERROR(D27/D8,0)",
        "B29": "=B25+B27", "C29": "=C25+C27", "D29": "=D25+D27",
        "B35": "=B8*(1-'Assumptions'!B18)", "C35": "=C8*(1-'Assumptions'!B18)", "D35": "=D8*(1-'Assumptions'!B18)",
        "B36": "=B35-(B13*(1-'Assumptions'!B18))-(B18*(1-'Assumptions'!B18))-B19-MAX(0,(B35-(B13*(1-'Assumptions'!B18))-(B18*(1-'Assumptions'!B18))-B19)*'Assumptions'!B17)",
        "C36": "=C35-(C13*(1-'Assumptions'!B18))-(C18*(1-'Assumptions'!B18))-C19-MAX(0,(C35-(C13*(1-'Assumptions'!B18))-(C18*(1-'Assumptions'!B18))-C19)*'Assumptions'!B17)",
        "D36": "=D35-(D13*(1-'Assumptions'!B18))-(D18*(1-'Assumptions'!B18))-D19-MAX(0,(D35-(D13*(1-'Assumptions'!B18))-(D18*(1-'Assumptions'!B18))-D19)*'Assumptions'!B17)",
        "B37": "='Assumptions'!B16+B36", "C37": "=B37+C36", "D37": "=C37+D36",
    }
    percentage_rows = {9, 15, 21, 28}
    for cell_ref, formula in model_formulas.items():
        cell = model[cell_ref]
        cell.value = formula
        cell.number_format = percent_format if cell.row in percentage_rows else units_format if cell.row == 6 else currency_format
        cell.font = Font(name="Aptos", size=10, color="008000" if "Assumptions" in formula and formula.startswith("='") else "000000")
    model["A34"] = "Investor Pushback Case"
    model["A35"] = "Revenue @ Investor Pushback"
    model["A36"] = "Net Income @ Investor Pushback"
    model["A37"] = "Ending Cash @ Investor Pushback"

    def assessment_from_status(status, green_value, warning_value, red_value):
        return green_value if status == "Green" else red_value if status == "Red" else warning_value

    # Keep margin interpretation tied to the margin benchmark, rather than a
    # generic status mapping.  A margin above the high end is not a low-margin
    # problem; it is an assumption that needs support.
    gross_margin = ((st.session_state.price_per_unit - st.session_state.cost_per_unit)
                    / st.session_state.price_per_unit
                    if st.session_state.price_per_unit else 0)
    margin_low, margin_high = INDUSTRY_BENCHMARKS[export_industry]["gross_margin"]
    if gross_margin > margin_high:
        margin_assessment = "Overstated"
    elif gross_margin < margin_low:
        margin_assessment = "Low Margin"
    else:
        margin_assessment = "Healthy"
    margin_status = (
        "Red" if gross_margin > 0.95 or gross_margin > margin_high + 0.10
        else benchmark_rows[0]["status"]
    )

    # Cash viability must assess liquidity, not gross-margin completeness.
    # Use the current calculated projection because the template formulas are
    # recalculated by Excel after download.
    ending_cash_values = projection_df["Ending Cash"]
    ending_cash_year3 = ending_cash_values.iloc[-1]
    year1_net_income = projection_df.loc[0, "Net Income"]
    if ending_cash_values.min() < 0:
        cash_assessment = "Funding Gap"
        cash_status = "Red"
        cash_note = "Projected cash balance turns negative and requires outside funding or cost reduction."
    elif year1_net_income < 0:
        cash_assessment = "Needs Runway Support"
        cash_status = "Yellow"
        cash_note = "Ending cash remains positive, but Year 1 losses require clear runway and milestone support."
    else:
        cash_assessment = "Positive"
        cash_status = "Green"
        cash_note = "Projected cash remains positive through the forecast period."

    segment_status = "Yellow" if customer_segment in ("Mixed / Unclear", "Mixed", "Unclear") else "Green"
    snapshot_rows = [
        ("Industry", export_industry, "Green", "Category detected from the founder input."),
        ("Customer Segment", customer_segment, segment_status, "Customer segment inferred from the startup description."),
        ("Overall Readiness", assessment_from_status(reality_engine_output["overall"], "Ready", "Needs Work", "High Risk"), reality_engine_output["overall"], reality_engine_output["summary"]),
        ("Pricing Realism", assessment_from_status(checks["Pricing Market Fit"]["status"], "Reasonable", "Needs Support", "High"), checks["Pricing Market Fit"]["status"], checks["Pricing Market Fit"]["message"]),
        ("Sales Volume", assessment_from_status(checks["Adoption Realism"]["status"], "Reasonable", "Aggressive", "Unrealistic"), checks["Adoption Realism"]["status"], checks["Adoption Realism"]["message"]),
        ("Growth Assumptions", assessment_from_status(checks["Growth Realism"]["status"], "Plausible", "Conservative", "Aggressive"), checks["Growth Realism"]["status"], checks["Growth Realism"]["message"]),
        ("Margin Quality", margin_assessment, margin_status, benchmark_rows[0]["feedback"]),
        ("Cash Viability", cash_assessment, cash_status, cash_note),
    ]
    dashboard_values = {
    }
    for row, (category, assessment, status, notes) in enumerate(snapshot_rows, 5):
        dashboard_values.update({f"F{row}": category, f"G{row}": assessment,
                                 f"H{row}": status, f"I{row}": notes})
    for cell, value in dashboard_values.items():
        dashboard[cell] = clean_excel_text(value)
    for cell in ("H5", "H6", "H7", "H8", "H9", "H10", "H11", "H12"):
        apply_excel_status_style(dashboard[cell], dashboard[cell].value)
    for row in range(5, 13):
        for col in (7, 9):
            dashboard.cell(row, col).fill = PatternFill(fill_type=None)
            dashboard.cell(row, col).font = Font(name="Aptos", size=10, color="1F1F1F")
        dashboard.cell(row, 9).alignment = Alignment(vertical="top", wrap_text=True)
        dashboard.row_dimensions[row].height = 42

    # Long founder inputs and investor-feedback text should stay readable.
    assumptions.column_dimensions["B"].width = 60
    assumptions.row_dimensions[5].height = 85
    assumptions.row_dimensions[19].height = 90
    for row in range(5, 20):
        for col in range(1, 6):
            assumptions.cell(row, col).alignment = Alignment(vertical="top", wrap_text=True)
    for cell in ("B5", "B19"):
        assumptions[cell].alignment = Alignment(vertical="top", wrap_text=True)

    # Keep the notes column readable without moving the chart section below it.
    dashboard.column_dimensions["I"].width = 65
    for row in range(5, 13):
        for col in range(6, 10):
            dashboard.cell(row, col).alignment = Alignment(vertical="top", wrap_text=True)

    # Widen the investor-facing review sheets and give wrapped content enough
    # vertical space to remain readable without changing their layout.
    for column, width in {"A": 22, "B": 28, "C": 18, "D": 18, "E": 68}.items():
        review.column_dimensions[column].width = width
    for row in range(6, 12):
        review.row_dimensions[row].height = 60

    if is_software_industry(export_industry):
        review_rows = [
            ("Traction", "Year 1 Units", "='Financial Model'!B6", checks["Adoption Realism"]["status"], "Define whether units are paid locations, subscriptions, active accounts, or transactions, then support the target with pipeline and conversion math."),
            ("Unit Economics", "Gross Margin %", "='Financial Model'!B15", benchmark_rows[0]["status"], "Validate hosting, data, onboarding, customer support, payment processing, and implementation costs."),
            ("Growth", "Year 2 Growth", "='Assumptions'!B11", benchmark_rows[1]["status"], "Tie growth to paid customer acquisition, retention, expansion revenue, and sales capacity."),
            ("Growth", "Year 3 Growth", "='Assumptions'!B12", benchmark_rows[2]["status"], "Tie growth to paid customer acquisition, retention, expansion revenue, and sales capacity."),
            ("Operating Model", "Total OpEx % of Revenue", "='Financial Model'!B21", benchmark_rows[3]["status"], "Break out Sales & Marketing, Product/R&D, G&A, implementation, and customer success."),
            ("Liquidity", "Ending Cash", "='Financial Model'!B29", checks["Financial Assumption Reality"]["status"], "Connect the cash plan to burn, runway, milestone timing, and any required funding round."),
        ]
    else:
        review_rows = [
            ("Traction", "Year 1 Units", "='Financial Model'!B6", checks["Adoption Realism"]["status"], "Provide provider-supply, booking-frequency, route-density, and conversion evidence."),
            ("Unit Economics", "Gross Margin %", "='Financial Model'!B15", benchmark_rows[0]["status"], "Validate contractor payouts, insurance, refunds, payment fees, and support costs."),
            ("Growth", "Year 2 Growth", "='Assumptions'!B11", benchmark_rows[1]["status"], "Tie growth to retention, repeat bookings, and provider capacity."),
            ("Growth", "Year 3 Growth", "='Assumptions'!B12", benchmark_rows[2]["status"], "Show how growth remains credible as the local supply base expands."),
            ("Operating Model", "Total OpEx % of Revenue", "='Financial Model'!B21", benchmark_rows[3]["status"], "Document total GTM, operations, and fixed-overhead investment."),
            ("Liquidity", "Ending Cash", "='Financial Model'!B29", checks["Financial Assumption Reality"]["status"], "Connect the cash plan to milestones, burn, and any funding requirement."),
        ]
    for row, (category, metric, current_value, status, recommended_fix) in enumerate(review_rows, 6):
        review.cell(row, 1, clean_excel_text(category))
        review.cell(row, 2, clean_excel_text(metric))
        review.cell(row, 3, current_value)
        review.cell(row, 3).font = Font(name="Aptos", size=10, color="1F1F1F")
        review.cell(row, 3).fill = PatternFill(fill_type=None)
        review.cell(row, 3).number_format = (
            units_format if metric == "Year 1 Units" else
            currency_format if metric == "Ending Cash" else percent_format
        )
        review.cell(row, 4, clean_excel_text(status))
        review.cell(row, 5, clean_excel_text(recommended_fix))
        apply_excel_status_style(review.cell(row, 4), status)
        for col in (1, 2, 5):
            review.cell(row, col).fill = PatternFill(fill_type=None)
            review.cell(row, col).font = Font(name="Aptos", size=10, bold=(col in (1, 2)), color="1F1F1F")
        review.cell(row, 5).alignment = Alignment(vertical="top", wrap_text=True)
        review.cell(row, 5).font = Font(name="Aptos", size=10, bold=False, color="000000")
        review.row_dimensions[row].height = 48

    # The investor-preparation questions use the template's lower section.
    # Apply the same wrapped, top-aligned treatment to all cells in that block.
    review.column_dimensions["B"].width = max(review.column_dimensions["B"].width or 0, 28)
    review.column_dimensions["C"].width = max(review.column_dimensions["C"].width or 0, 42)
    for row in range(15, 20):
        for col in range(1, 4):
            review.cell(row, col).alignment = Alignment(vertical="top", wrap_text=True)
    for row in range(16, 20):
        review.row_dimensions[row].height = 68

    if is_software_industry(export_industry):
        software_questions = [
            ("What exactly is one unit?", "Unit definition determines whether 80,000 units is realistic.", "Clarify whether units are paid locations, subscriptions, active accounts, monthly transactions, or monitored sites."),
            ("What costs are included in COGS?", "Gross margin depends on what is included below the line.", "Add hosting, data, implementation, onboarding, customer support, payment fees, and AI infrastructure costs."),
            ("How will you acquire 80,000 units in Year 1?", "Customer acquisition is the biggest model risk.", "Show funnel math: target accounts, outreach volume, pilot conversion, paid conversion, churn, and expansion."),
            ("Why is OpEx only 25% of revenue before fixed overhead?", "Growth-stage software companies usually need meaningful S&M, product, and support investment.", "Break out Sales & Marketing, Product/R&D, G&A, implementation, and customer success costs."),
        ]
        for row, values in enumerate(software_questions, 16):
            for col, value in enumerate(values, 1):
                review.cell(row, col, clean_excel_text(value))

    benchmark_categories = ["Gross Margin", "Year 2 Growth", "Year 3 Growth", "OpEx Ratio", "Year 1 Volume"]
    for row, category, feedback in zip(range(6, 11), benchmark_categories, benchmark_feedback):
        text = str(feedback).replace("\n", " ")
        status = "Green" if "ðŸŸ¢" in text else "Yellow" if "ðŸŸ¡" in text or "ðŸŸ " in text else "Red"
        benchmark.cell(row, 1, clean_excel_text(category))
        benchmark.cell(row, 2, clean_excel_text(status))
        benchmark.cell(row, 3, clean_excel_text(text))
        apply_excel_status_style(benchmark.cell(row, 2), status)
        benchmark.cell(row, 3).alignment = Alignment(vertical="top", wrap_text=True)
        benchmark.cell(row, 4).alignment = Alignment(vertical="top", wrap_text=True)
        benchmark.cell(row, 3).font = Font(name="Aptos", size=10, bold=False, color="000000")
        benchmark.cell(row, 4).font = Font(name="Aptos", size=10, bold=False, color="000000")
        benchmark.row_dimensions[row].height = 54

    # Overwrite the legacy text rows with structured, emoji-free benchmark output.
    for row, item in zip(range(6, 11), benchmark_rows):
        benchmark.cell(row, 1, clean_excel_text(item["category"]))
        benchmark.cell(row, 2, clean_excel_text(item["status"]))
        benchmark.cell(row, 3, clean_excel_text(item["feedback"]))
        benchmark.cell(row, 4, clean_excel_text(item["basis"]))
        apply_excel_status_style(benchmark.cell(row, 2), item["status"])
        for col in (3, 4):
            benchmark.cell(row, col).fill = PatternFill(fill_type=None)
            benchmark.cell(row, col).font = Font(name="Aptos", size=10, bold=False, color="1F1F1F")

    # C11 must contain one concise recommendation, not the entire feedback list.
    suggested_adjustments = [
        str(item).replace("\n", " ") for item in benchmark_feedback
        if str(item).startswith("Suggested fix:")
    ]
    suggested_text = next(
        (item for item in suggested_adjustments if "Year 1 units" in item),
        suggested_adjustments[0] if suggested_adjustments else
        "Suggested fix: provide stronger bottom-up demand evidence for the current assumptions.",
    )
    benchmark["A11"] = clean_excel_text("Suggested Adjustment")
    benchmark["B11"] = clean_excel_text("Action")
    benchmark["C11"] = clean_excel_text(
        "Clarify the unit definition, reduce or support Year 1 volume with bottom-up pipeline math, and validate gross margin with hosting, support, onboarding, and AI infrastructure costs."
        if is_software_industry(export_industry) else suggested_fix
    )
    benchmark["D11"] = clean_excel_text("Investor readiness adjustment")
    apply_excel_status_style(benchmark["B11"], "Action")
    for cell in ("C11", "D11"):
        benchmark[cell].alignment = Alignment(vertical="top", wrap_text=True)
        benchmark[cell].font = Font(name="Aptos", size=10, bold=False, color="000000")
    benchmark.row_dimensions[11].height = 42

    # Benchmark feedback stays plain text, while status styling remains
    # isolated to column B.
    for row in range(6, 12):
        for col in (3, 4):
            benchmark.cell(row, col).alignment = Alignment(vertical="top", wrap_text=True)
            benchmark.cell(row, col).fill = PatternFill(fill_type=None)
            benchmark.cell(row, col).font = Font(name="Aptos", size=10, bold=False, color="1F1F1F")
        benchmark.row_dimensions[row].height = 48 if row < 11 else 42
    for column, width in {"A": 24, "B": 16, "C": 72, "D": 62}.items():
        benchmark.column_dimensions[column].width = width
    for row in range(6, 12):
        benchmark.row_dimensions[row].height = 60

    for row, (label, item) in enumerate(checks.items(), 6):
        reality.cell(row, 1, clean_excel_text(label))
        reality.cell(row, 2, clean_excel_text(item["message"]))
        reality.cell(row, 3, clean_excel_text(item["status"]))
        apply_excel_status_style(reality.cell(row, 3), item["status"])
        reality.cell(row, 2).fill = PatternFill(fill_type=None)
        reality.cell(row, 2).alignment = Alignment(vertical="top", wrap_text=True)
        reality.cell(row, 2).font = Font(name="Aptos", size=10, bold=False, color="000000")
        reality.row_dimensions[row].height = 48
    reality["B12"] = clean_excel_text(reality_engine_output["summary"])
    reality["C12"] = clean_excel_text(reality_engine_output["overall"])
    apply_excel_status_style(reality["C12"], reality_engine_output["overall"])
    reality["B12"].alignment = Alignment(vertical="top", wrap_text=True)
    reality["B12"].font = Font(name="Aptos", size=10, bold=False, color="000000")
    reality.row_dimensions[12].height = 48
    price_display = f"${st.session_state.price_per_unit:,.2f}".rstrip("0").rstrip(".")
    reality["D7"] = clean_excel_text(
        f"Clarify whether {price_display} is monthly, per order, per route, "
        "per location, or per transaction."
    )
    for column, width in {"A": 32, "B": 72, "C": 16, "D": 72}.items():
        reality.column_dimensions[column].width = width
    # Keep the main Dashboard chart labels explicit in the exported workbook.
    if dashboard._charts:
        main_dashboard_chart = dashboard._charts[0]
        main_dashboard_chart.title = "Revenue, Net Income & Cash by Year"
        main_dashboard_chart.legend.position = "b"
        main_dashboard_chart.x_axis.title = "Forecast Year"
        main_dashboard_chart.y_axis.numFmt = '$#,##0'

    # Give only the long, wrapped Dashboard and Reality Engine notes extra room.
    for row in range(5, 13):
        if len(str(dashboard.cell(row, 9).value or "")) > 120:
            dashboard.row_dimensions[row].height = max(dashboard.row_dimensions[row].height or 0, 60)
    for row in range(6, 13):
        if len(str(reality.cell(row, 2).value or "")) > 120 or len(str(reality.cell(row, 4).value or "")) > 120:
            reality.row_dimensions[row].height = max(reality.row_dimensions[row].height or 0, 66)

    # Keep exported sheets presentation-ready when printed or opened in Excel.
    for sheet, print_area in {
        assumptions: "A1:E20", dashboard: "A1:I32", model: "A1:D37",
        review: "A1:E19", benchmark: "A1:D11", reality: "A1:C12",
    }.items():
        sheet.sheet_view.showGridLines = False
        sheet.print_area = print_area
        sheet.page_setup.fitToWidth = 1
        sheet.page_setup.fitToHeight = 0

    # Leave every existing formula untouched and have Excel refresh dependent
    # dashboard/model values when the workbook is opened. Some templates do
    # not include calculation properties, so this optional setting must not
    # prevent the populated workbook from being saved.
    calculation = template_wb.calculation
    if calculation is not None:
        calculation.fullCalcOnLoad = True
        calculation.forceFullCalc = True
    return template_wb


template_path = os.path.join(os.path.dirname(__file__), "templates", "financial_model_professional.xlsx")
if os.path.exists(template_path):
    wb = populate_professional_excel_template(template_path)


def verify_export_workbook(workbook):
    """Fail fast if the professional workbook receives an incomplete export."""
    model = workbook["Financial Model"]
    dashboard = workbook["Dashboard"]
    benchmark = workbook["Benchmark Feedback"]
    review = workbook["Investor Review"]
    required_model_cells = [f"{col}{row}" for row in (6, 7, 8, 9, 12, 13, 14, 15, 18, 19,
                                                         20, 21, 22, 25, 26, 27, 28, 29, 35, 36, 37)
                            for col in "BCD"]
    if any(model[cell].value in (None, "") for cell in required_model_cells):
        raise ValueError("Excel export verification failed: Financial Model contains blank projection cells.")
    color_words = {"Green", "Yellow", "Orange", "Red", "Action"}
    assessments = [dashboard[f"G{row}"].value for row in range(5, 13)]
    statuses = [dashboard[f"H{row}"].value for row in range(5, 13)]
    if assessments and all(value in color_words for value in assessments):
        raise ValueError("Excel export verification failed: Dashboard assessments are status words.")
    if any(value not in color_words for value in statuses):
        raise ValueError("Excel export verification failed: Dashboard status cells are incomplete.")
    benchmark_statuses = [benchmark[f"B{row}"].value for row in range(6, 11)]
    if benchmark_statuses and all(value == "Red" for value in benchmark_statuses):
        raise ValueError("Excel export verification failed: Benchmark statuses are all Red.")
    emoji_prefixes = ("🟢", "🟡", "🟠", "🔴")
    if any(str(benchmark.cell(row, col).value or "").startswith(emoji_prefixes)
           for row in range(6, 12) for col in (3, 4)):
        raise ValueError("Excel export verification failed: Benchmark feedback contains emoji prefixes.")
    if any(review[f"C{row}"].value in (None, "") for row in range(6, 12)):
        raise ValueError("Excel export verification failed: Investor Review current values are blank.")
    expected_review_headers = ["Category", "Metric", "Current Value", "Investor Read", "Recommended Fix"]
    if [review.cell(5, col).value for col in range(1, 6)] != expected_review_headers:
        raise ValueError("Excel export verification failed: Investor Review headers are incorrect.")
    volume_statuses = [dashboard["H9"].value, review["D6"].value, benchmark["B10"].value, workbook["Reality Engine"]["C8"].value]
    if len(set(volume_statuses)) != 1:
        raise ValueError("Excel export verification failed: Year 1 volume status differs across tabs.")
    for col in "BCD":
        downside_formula = str(model[f"{col}36"].value or "")
        if "(1-'Assumptions'!B18)" not in downside_formula or f"{col}19" not in downside_formula or "MAX(0" not in downside_formula:
            raise ValueError("Excel export verification failed: Investor pushback formulas are incomplete.")
    margin_assessment = dashboard["G11"].value
    margin_note = str(dashboard["I11"].value or "").lower()
    if "above the typical" in margin_note and margin_assessment == "Low Margin":
        raise ValueError("Excel export verification failed: Above-benchmark margin is labeled Low Margin.")
    cash_note = str(dashboard["I12"].value or "").lower()
    if "gross margin" in cash_note or "cogs" in cash_note:
        raise ValueError("Excel export verification failed: Cash Viability note contains margin commentary.")
    assumptions = workbook["Assumptions"]
    if assumptions.column_dimensions["B"].width < 58 or assumptions.row_dimensions[5].height < 85 or assumptions.row_dimensions[19].height < 90:
        raise ValueError("Excel export verification failed: Assumptions long-text fields are not readable.")
    if any(not review.cell(row, col).alignment.wrap_text or review.cell(row, col).alignment.vertical != "top"
           for row in range(15, 20) for col in range(1, 4)):
        raise ValueError("Excel export verification failed: Investor preparation text is not wrapped and top-aligned.")
    if any((review.row_dimensions[row].height or 0) < 68 for row in range(16, 20)):
        raise ValueError("Excel export verification failed: Investor preparation rows are too short.")

    export_industry = str(assumptions["B6"].value or "")
    if is_software_industry(export_industry):
        if dashboard["G5"].value != export_industry or dashboard["G6"].value != "B2B / SMB Operators":
            raise ValueError("Excel export verification failed: SaaS dashboard classification is incorrect.")
        forbidden_terms = ("restaurants", "route density", "provider supply", "delivery ops")
        dashboard_review_text = " ".join(
            str(sheet.cell(row, col).value or "")
            for sheet in (dashboard, review)
            for row in range(1, sheet.max_row + 1)
            for col in range(1, sheet.max_column + 1)
        ).lower()
        if any(term in dashboard_review_text for term in forbidden_terms):
            raise ValueError("Excel export verification failed: SaaS export contains marketplace-specific language.")
        required_question_text = (
            "paid locations, subscriptions, active accounts, monthly transactions, or monitored sites.",
            "hosting, data, implementation, onboarding, customer support, payment fees, and ai infrastructure costs.",
            "target accounts, outreach volume, pilot conversion, paid conversion, churn, and expansion.",
        )
        question_text = " ".join(str(review.cell(row, col).value or "") for row in range(16, 20) for col in range(1, 4)).lower()
        if any(text not in question_text for text in required_question_text):
            raise ValueError("Excel export verification failed: SaaS investor questions are incomplete.")
        suggestion = str(benchmark["C11"].value or "").lower()
        if not all(text in suggestion for text in ("unit definition", "year 1 volume", "gross margin")):
            raise ValueError("Excel export verification failed: SaaS benchmark adjustment is incomplete.")


verify_export_workbook(wb)
excel_buffer = io.BytesIO()
wb.save(excel_buffer)
excel_buffer.seek(0)

# ---------------- PowerPoint ----------------
prs = Presentation()

INDUSTRY_PPT_THEMES = {
    "SaaS": {
        "navy": RGBColor(11, 31, 63),
        "navy_2": RGBColor(20, 52, 101),
        "blue": RGBColor(37, 99, 235),
        "cyan": RGBColor(56, 189, 248),
        "ice": RGBColor(239, 246, 255),
        "pale_blue": RGBColor(219, 234, 254),
        "slate": RGBColor(45, 55, 72),
        "muted": RGBColor(100, 116, 139),
        "soft": RGBColor(248, 250, 252),
        "border": RGBColor(191, 219, 254),
        "white": RGBColor(255, 255, 255),
        "green": RGBColor(220, 252, 231),
        "yellow": RGBColor(254, 249, 195),
        "red": RGBColor(254, 226, 226),
        "green_text": RGBColor(22, 101, 52),
        "yellow_text": RGBColor(133, 77, 14),
        "red_text": RGBColor(153, 27, 27),
        "primary": RGBColor(11, 31, 63),
        "secondary": RGBColor(20, 52, 101),
        "accent": RGBColor(37, 99, 235),
        "background": RGBColor(248, 250, 252),
        "text": RGBColor(45, 55, 72),
        "chart": ["#2563EB", "#38BDF8", "#0B1F3F"],
        "visual_motif": "clean software dashboard geometry",
        "hero_style": "polished business software with crisp blue accents",
        "image_prompt": "modern SaaS dashboard interface, clean navy and bright blue palette, white and light gray background",
    },
    "Marketplace": {
        "navy": RGBColor(49, 46, 129),
        "navy_2": RGBColor(67, 56, 202),
        "blue": RGBColor(13, 148, 136),
        "cyan": RGBColor(45, 212, 191),
        "ice": RGBColor(240, 253, 250),
        "pale_blue": RGBColor(221, 214, 254),
        "slate": RGBColor(45, 55, 72),
        "muted": RGBColor(100, 116, 139),
        "soft": RGBColor(245, 243, 255),
        "border": RGBColor(153, 246, 228),
        "white": RGBColor(255, 255, 255),
        "green": RGBColor(220, 252, 231),
        "yellow": RGBColor(254, 249, 195),
        "red": RGBColor(254, 226, 226),
        "green_text": RGBColor(22, 101, 52),
        "yellow_text": RGBColor(133, 77, 14),
        "red_text": RGBColor(153, 27, 27),
        "primary": RGBColor(49, 46, 129),
        "secondary": RGBColor(67, 56, 202),
        "accent": RGBColor(13, 148, 136),
        "background": RGBColor(245, 243, 255),
        "text": RGBColor(45, 55, 72),
        "chart": ["#4338CA", "#0D9488", "#2DD4BF"],
        "visual_motif": "connected nodes and platform network lines",
        "hero_style": "modern marketplace platform with indigo structure and teal connection accents",
        "image_prompt": "abstract marketplace network platform, indigo and teal palette, cool light background, connected nodes",
    },
    "Consumer Product": {
        "navy": RGBColor(38, 38, 38),
        "navy_2": RGBColor(79, 70, 63),
        "blue": RGBColor(217, 119, 6),
        "cyan": RGBColor(245, 158, 11),
        "ice": RGBColor(255, 247, 237),
        "pale_blue": RGBColor(254, 215, 170),
        "slate": RGBColor(63, 55, 48),
        "muted": RGBColor(120, 113, 108),
        "soft": RGBColor(255, 251, 235),
        "border": RGBColor(231, 203, 166),
        "white": RGBColor(255, 253, 250),
        "green": RGBColor(220, 252, 231),
        "yellow": RGBColor(254, 249, 195),
        "red": RGBColor(254, 226, 226),
        "green_text": RGBColor(22, 101, 52),
        "yellow_text": RGBColor(133, 77, 14),
        "red_text": RGBColor(153, 27, 27),
        "primary": RGBColor(38, 38, 38),
        "secondary": RGBColor(79, 70, 63),
        "accent": RGBColor(217, 119, 6),
        "background": RGBColor(255, 251, 235),
        "text": RGBColor(63, 55, 48),
        "chart": ["#D97706", "#F59E0B", "#262626"],
        "visual_motif": "premium product surfaces and warm editorial accents",
        "hero_style": "premium consumer brand with cream background, charcoal typography, and gold accents",
        "image_prompt": "premium consumer product presentation, cream off-white background, charcoal text, warm orange gold accents",
    },
    "Food / Delivery": {
        "navy": RGBColor(47, 34, 30),
        "navy_2": RGBColor(91, 52, 39),
        "blue": RGBColor(220, 38, 38),
        "cyan": RGBColor(249, 115, 22),
        "ice": RGBColor(255, 237, 213),
        "pale_blue": RGBColor(254, 202, 147),
        "slate": RGBColor(74, 54, 45),
        "muted": RGBColor(133, 77, 54),
        "soft": RGBColor(255, 247, 237),
        "border": RGBColor(253, 186, 116),
        "white": RGBColor(255, 255, 255),
        "green": RGBColor(220, 252, 231),
        "yellow": RGBColor(254, 249, 195),
        "red": RGBColor(254, 226, 226),
        "green_text": RGBColor(22, 101, 52),
        "yellow_text": RGBColor(133, 77, 14),
        "red_text": RGBColor(153, 27, 27),
        "primary": RGBColor(47, 34, 30),
        "secondary": RGBColor(91, 52, 39),
        "accent": RGBColor(220, 38, 38),
        "background": RGBColor(255, 247, 237),
        "text": RGBColor(74, 54, 45),
        "chart": ["#DC2626", "#F97316", "#2F221E"],
        "visual_motif": "energetic service flow with warm delivery accents",
        "hero_style": "energetic food delivery brand with orange-red accents and soft beige surfaces",
        "image_prompt": "food delivery service visual theme, orange red accents, warm charcoal typography, beige peach background",
    },
    "AI Startup": {
        "navy": RGBColor(2, 6, 23),
        "navy_2": RGBColor(20, 16, 56),
        "blue": RGBColor(0, 178, 255),
        "cyan": RGBColor(139, 92, 246),
        "ice": RGBColor(238, 242, 255),
        "pale_blue": RGBColor(196, 181, 253),
        "slate": RGBColor(30, 41, 59),
        "muted": RGBColor(100, 116, 139),
        "soft": RGBColor(248, 250, 252),
        "border": RGBColor(165, 180, 252),
        "white": RGBColor(255, 255, 255),
        "green": RGBColor(220, 252, 231),
        "yellow": RGBColor(254, 249, 195),
        "red": RGBColor(254, 226, 226),
        "green_text": RGBColor(22, 101, 52),
        "yellow_text": RGBColor(133, 77, 14),
        "red_text": RGBColor(153, 27, 27),
        "primary": RGBColor(2, 6, 23),
        "secondary": RGBColor(20, 16, 56),
        "accent": RGBColor(0, 178, 255),
        "background": RGBColor(248, 250, 252),
        "text": RGBColor(30, 41, 59),
        "chart": ["#00B2FF", "#8B5CF6", "#020617"],
        "visual_motif": "technical grid, signal lines, and neural interface accents",
        "hero_style": "futuristic technical deck with dark navy, electric blue, and violet highlights",
        "image_prompt": "futuristic AI startup interface, dark navy, electric blue accents, violet highlights, technical grid",
    },
}

selected_ppt_industry = st.session_state.get("industry", "SaaS")
PPT_THEME = INDUSTRY_PPT_THEMES.get(selected_ppt_industry, INDUSTRY_PPT_THEMES["SaaS"])


def add_block(slide, x, y, w, h, fill_color, line_color=None):
    block = slide.shapes.add_textbox(PPTInches(x), PPTInches(y), PPTInches(w), PPTInches(h))
    block.fill.solid()
    block.fill.fore_color.rgb = fill_color
    block.line.color.rgb = line_color or fill_color
    return block


def clean_pitch_body_text(text):
    clean_text = str(text or "")
    duplicate_prefixes = [
        "Segment Focus:",
        "Integrated Workflow:",
        "Fast Iteration:",
        "SMB Focus:",
        "All-in-One Platform:",
        "Credibility:",
        "Risk:",
        "Evidence needed:",
    ]
    for prefix in duplicate_prefixes:
        if clean_text.lower().startswith(prefix.lower()):
            return clean_text[len(prefix):].lstrip()
    return clean_text


def add_textbox(slide, x, y, w, h, text, size, color, bold=False, align=None):
    box = slide.shapes.add_textbox(PPTInches(x), PPTInches(y), PPTInches(w), PPTInches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = clean_pitch_body_text(text)
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    if align:
        p.alignment = align
    return box


def add_shape_block(slide, x, y, w, h, fill_color, line_color=None, shape_type=MSO_SHAPE.RECTANGLE):
    shape = slide.shapes.add_shape(
        shape_type,
        PPTInches(x),
        PPTInches(y),
        PPTInches(w),
        PPTInches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = line_color or fill_color
    return shape


def get_ppt_hero_name():
    for key in ("company_name", "startup_name", "business_name"):
        value = st.session_state.get(key, "")
        if value and str(value).strip():
            return str(value).strip()[:42]

    idea = st.session_state.get("idea", "")
    if idea and str(idea).strip():
        first_line = str(idea).strip().splitlines()[0]
        first_sentence = re.split(r"(?<=[.!?])\s+", first_line)[0]
        clean_idea = first_sentence.strip()
        clean_idea = re.split(r"\s*[:|]\s*", clean_idea, maxsplit=1)[0]
        name_candidate = re.split(
            r"\s+(?:is\s+a|is\s+an|helps|helping|provides|delivers|enables|offers|builds|creates|for)\s+",
            clean_idea,
            maxsplit=1,
            flags=re.IGNORECASE,
        )[0]
        name_candidate = re.sub(r"^(?:the\s+startup|startup|company|business|idea)\s*[:\-]\s*", "", name_candidate, flags=re.IGNORECASE)
        name_candidate = name_candidate.strip(" -:,.")
        if 1 < len(name_candidate) <= 42 and len(name_candidate.split()) <= 4:
            return name_candidate

    return "Startup Investor Deck"


def get_ppt_subtitle():
    return {
        "SaaS": "Revenue operations software for modern teams",
        "Marketplace": "Platform connecting buyers and sellers",
        "AI Startup": "AI-powered platform for intelligent workflows",
        "Consumer Product": "Premium consumer brand built for modern customers",
        "Food / Delivery": "Food and delivery platform for modern service operations",
    }.get(selected_ppt_industry, "Revenue operations software for modern teams")


def get_industry_badge_label(industry):
    return {
        "SaaS": "SaaS",
        "Marketplace": "Marketplace Platform",
        "Consumer Product": "Consumer Product",
        "Food / Delivery": "Food / Delivery",
        "AI Startup": "AI Startup",
    }.get(industry, industry or "Startup")


def add_industry_badge(slide, x, y, label, dark=False):
    fill = PPT_THEME["navy_2"] if dark else PPT_THEME["ice"]
    line = PPT_THEME["cyan"] if dark else PPT_THEME["border"]
    text = PPT_THEME["cyan"] if dark else PPT_THEME["navy"]
    badge_width = max(1.45, min(2.35, 0.18 + (len(label) * 0.08)))
    badge = add_shape_block(slide, x, y, badge_width, 0.3, fill, line, MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = badge.text_frame
    tf.margin_left = PPTInches(0.1)
    tf.margin_right = PPTInches(0.1)
    tf.margin_top = PPTInches(0.035)
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = label.upper()
    p.font.size = Pt(7.5)
    p.font.bold = True
    p.font.color.rgb = text
    p.alignment = PP_ALIGN.CENTER
    return badge


def get_theme_image_path(industry):
    image_dir = os.path.join("assets", "industry_images")
    image_paths = {
        "SaaS": os.path.join(image_dir, "saas.jpg"),
        "Marketplace": os.path.join(image_dir, "marketplace.jpg"),
        "AI Startup": os.path.join(image_dir, "ai_startup.jpg"),
        "Consumer Product": os.path.join(image_dir, "consumer_product.jpg"),
        "Food / Delivery": os.path.join(image_dir, "food_delivery.jpg"),
        "Healthcare": os.path.join(image_dir, "healthcare.jpg"),
        "Fintech": os.path.join(image_dir, "fintech.jpg"),
        "General Business": os.path.join(image_dir, "general_business.jpg"),
    }
    image_path = image_paths.get(industry)
    return image_path if image_path and os.path.exists(image_path) else None


def add_safe_picture(slide, image_path, x, y, w, h):
    if not image_path or not os.path.exists(image_path):
        return False
    try:
        slide.shapes.add_picture(
            image_path,
            PPTInches(x),
            PPTInches(y),
            width=PPTInches(w),
            height=PPTInches(h),
        )
        return True
    except Exception:
        return False


def add_saas_image_accent(slide, x, y, w, h):
    return add_safe_picture(slide, get_theme_image_path("SaaS"), x, y, w, h)


def add_panel_label(slide, x, y, label, value=None, dark=False):
    color = PPT_THEME["white"] if dark else PPT_THEME["navy"]
    muted = PPT_THEME["pale_blue"] if dark else PPT_THEME["muted"]
    add_textbox(slide, x, y, 0.95, 0.15, label.upper(), 5.6, muted, True)
    if value:
        add_textbox(slide, x, y + 0.17, 1.2, 0.18, value, 8.5, color, True)


def add_theme_cover_visual(slide, industry):
    if industry == "SaaS":
        add_saas_image_accent(slide, 6.18, 0.64, 3.38, 6.02)

    image_path = None if industry == "SaaS" else get_theme_image_path(industry)
    if image_path:
        add_shape_block(slide, 6.38, 0.88, 3.12, 5.55, PPT_THEME["navy_2"], PPT_THEME["cyan"], MSO_SHAPE.ROUNDED_RECTANGLE)
        if add_safe_picture(slide, image_path, 6.46, 0.98, 2.96, 5.35):
            return

    dark_panel = industry in ("SaaS", "Marketplace", "AI Startup")
    panel_fill = PPT_THEME["navy_2"] if dark_panel else PPT_THEME["white"]
    panel_line = PPT_THEME["cyan"] if dark_panel else PPT_THEME["border"]
    panel_h = 5.05 if industry == "SaaS" else 5.55
    add_shape_block(slide, 6.35, 0.86, 3.12, panel_h, panel_fill, panel_line, MSO_SHAPE.ROUNDED_RECTANGLE)

    if industry == "SaaS":
        add_shape_block(slide, 6.5, 1.08, 2.82, 4.5, PPT_THEME["white"], PPT_THEME["blue"], MSO_SHAPE.ROUNDED_RECTANGLE)
        add_block(slide, 6.5, 1.08, 2.82, 0.42, PPT_THEME["navy"])
        for i, color in enumerate([PPT_THEME["cyan"], PPT_THEME["blue"], PPT_THEME["pale_blue"]]):
            add_shape_block(slide, 6.66 + (i * 0.18), 1.23, 0.08, 0.08, color, color, MSO_SHAPE.OVAL)
        add_textbox(slide, 8.02, 1.21, 1.0, 0.14, "Forecast", 6.6, PPT_THEME["cyan"], True, PP_ALIGN.RIGHT)

        add_shape_block(slide, 6.5, 1.5, 0.42, 4.08, PPT_THEME["ice"], PPT_THEME["border"], MSO_SHAPE.RECTANGLE)
        for i in range(4):
            y = 1.82 + (i * 0.58)
            add_shape_block(slide, 6.64, y, 0.14, 0.14, PPT_THEME["cyan"] if i == 0 else PPT_THEME["blue"], PPT_THEME["cyan"] if i == 0 else PPT_THEME["blue"], MSO_SHAPE.OVAL)

        kpis = [("ARR", "$2.4M"), ("NRR", "118%"), ("Pipeline", "$640K")]
        for i, (label, value) in enumerate(kpis):
            x = 7.06 + (i * 0.72)
            add_shape_block(slide, x, 1.78, 0.62, 0.58, PPT_THEME["ice"], PPT_THEME["border"], MSO_SHAPE.ROUNDED_RECTANGLE)
            add_textbox(slide, x + 0.08, 1.92, 0.48, 0.1, label, 5.0, PPT_THEME["muted"], True)
            add_textbox(slide, x + 0.08, 2.1, 0.48, 0.13, value, 6.3, PPT_THEME["navy"], True)

        add_shape_block(slide, 7.06, 2.6, 1.6, 1.86, PPT_THEME["white"], PPT_THEME["border"], MSO_SHAPE.ROUNDED_RECTANGLE)
        add_textbox(slide, 7.24, 2.82, 1.1, 0.12, "Revenue Trend", 5.9, PPT_THEME["navy"], True, PP_ALIGN.CENTER)
        add_block(slide, 7.24, 4.22, 1.16, 0.02, PPT_THEME["border"])
        for i, h in enumerate([0.34, 0.54, 0.44, 0.78, 0.66]):
            add_block(slide, 7.34 + (i * 0.2), 4.2 - h, 0.1, h, PPT_THEME["blue"] if i % 2 else PPT_THEME["cyan"])

        add_shape_block(slide, 8.78, 2.6, 0.36, 1.86, PPT_THEME["ice"], PPT_THEME["border"], MSO_SHAPE.ROUNDED_RECTANGLE)
        add_textbox(slide, 8.82, 2.84, 0.28, 0.12, "Forecast", 4.7, PPT_THEME["navy"], True, PP_ALIGN.CENTER)
        for i, width in enumerate([0.26, 0.22, 0.28, 0.18]):
            add_block(slide, 8.84, 3.24 + (i * 0.24), width, 0.04, PPT_THEME["blue"] if i % 2 else PPT_THEME["cyan"])
    elif industry == "Marketplace":
        add_textbox(slide, 6.62, 1.18, 1.8, 0.24, "Platform Network", 10, PPT_THEME["white"], True)
        nodes = [
            (7.8, 1.65, 0.52, "CORE", PPT_THEME["cyan"]),
            (6.72, 2.34, 0.42, "BUY", PPT_THEME["blue"]),
            (8.72, 2.42, 0.42, "SELL", PPT_THEME["blue"]),
            (7.16, 3.62, 0.38, "VND", PPT_THEME["pale_blue"]),
            (8.35, 3.66, 0.38, "FUL", PPT_THEME["pale_blue"]),
            (7.83, 4.72, 0.42, "PAY", PPT_THEME["cyan"]),
        ]
        for x1, y1, x2, y2 in [(8.04, 2.0, 6.94, 2.52), (8.12, 2.0, 8.92, 2.6), (6.95, 2.72, 7.36, 3.76), (8.92, 2.78, 8.54, 3.84), (7.44, 3.86, 8.08, 4.88), (8.5, 3.86, 8.1, 4.88)]:
            line = add_shape_block(slide, x1, y1, 0.84, 0.035, PPT_THEME["cyan"])
            line.rotation = 24 if y2 > y1 else -24
        for x, y, size, label, color in nodes:
            add_shape_block(slide, x, y, size, size, color, color, MSO_SHAPE.OVAL)
            add_textbox(slide, x - 0.04, y + 0.17, size + 0.08, 0.12, label, 5.2, PPT_THEME["navy"], True, PP_ALIGN.CENTER)
        add_panel_label(slide, 6.72, 5.45, "Flywheel", "Demand -> Supply", True)
        add_panel_label(slide, 8.38, 5.45, "Take Rate", "Platform Layer", True)
    elif industry == "AI Startup":
        for i in range(8):
            add_shape_block(slide, 6.58 + (i * 0.34), 1.18, 0.015, 4.65, PPT_THEME["blue"])
        for i in range(9):
            add_shape_block(slide, 6.55, 1.2 + (i * 0.5), 2.62, 0.015, PPT_THEME["cyan"])
        for x, y, size in [(7.02, 1.72, 0.18), (7.78, 2.28, 0.24), (8.58, 1.86, 0.16), (6.88, 3.38, 0.2), (8.3, 3.68, 0.24), (7.55, 4.62, 0.18)]:
            add_shape_block(slide, x, y, size, size, PPT_THEME["cyan"], PPT_THEME["cyan"], MSO_SHAPE.OVAL)
        for x, y, w, rot in [(7.14, 1.9, 0.82, 28), (7.92, 2.43, 0.8, -22), (7.05, 3.56, 1.45, 10), (7.7, 4.72, 0.82, -18)]:
            line = add_shape_block(slide, x, y, w, 0.028, PPT_THEME["blue"])
            line.rotation = rot
        add_shape_block(slide, 6.72, 5.25, 2.34, 0.52, PPT_THEME["navy"], PPT_THEME["blue"], MSO_SHAPE.ROUNDED_RECTANGLE)
        add_textbox(slide, 6.92, 5.38, 1.9, 0.18, "MODEL SIGNAL / DATA PIPELINE", 7, PPT_THEME["cyan"], True, PP_ALIGN.CENTER)
    elif industry == "Consumer Product":
        add_shape_block(slide, 6.68, 1.22, 1.0, 3.15, PPT_THEME["ice"], PPT_THEME["border"], MSO_SHAPE.ROUNDED_RECTANGLE)
        add_shape_block(slide, 6.86, 1.55, 0.64, 2.18, PPT_THEME["white"], PPT_THEME["cyan"], MSO_SHAPE.ROUNDED_RECTANGLE)
        add_shape_block(slide, 7.95, 1.52, 1.1, 1.18, PPT_THEME["white"], PPT_THEME["border"], MSO_SHAPE.ROUNDED_RECTANGLE)
        add_shape_block(slide, 7.95, 2.95, 1.1, 1.18, PPT_THEME["ice"], PPT_THEME["border"], MSO_SHAPE.ROUNDED_RECTANGLE)
        add_block(slide, 8.12, 1.82, 0.58, 0.06, PPT_THEME["cyan"])
        add_block(slide, 8.12, 2.03, 0.72, 0.04, PPT_THEME["blue"])
        add_block(slide, 8.12, 3.28, 0.52, 0.06, PPT_THEME["cyan"])
        add_block(slide, 8.12, 3.5, 0.68, 0.04, PPT_THEME["blue"])
        add_shape_block(slide, 6.62, 4.78, 2.54, 0.62, PPT_THEME["navy"], PPT_THEME["cyan"], MSO_SHAPE.ROUNDED_RECTANGLE)
        add_textbox(slide, 6.9, 4.96, 1.96, 0.2, "PREMIUM BRAND SYSTEM", 7.2, PPT_THEME["white"], True, PP_ALIGN.CENTER)
    elif industry == "Food / Delivery":
        add_textbox(slide, 6.62, 1.18, 1.85, 0.24, "Service Route", 10, PPT_THEME["navy"], True)
        route_segments = [(6.84, 2.05, 1.25, 18), (7.78, 2.66, 0.92, -22), (7.05, 3.54, 1.28, 18), (7.96, 4.16, 0.84, -18)]
        for x, y, w, rot in route_segments:
            route = add_shape_block(slide, x, y, w, 0.06, PPT_THEME["cyan"])
            route.rotation = rot
        for x, y, label in [(6.78, 1.84, "ORDER"), (8.62, 2.48, "KITCHEN"), (6.86, 3.36, "RIDER"), (8.58, 4.06, "DOOR")]:
            add_shape_block(slide, x, y, 0.54, 0.54, PPT_THEME["white"], PPT_THEME["border"], MSO_SHAPE.ROUNDED_RECTANGLE)
            add_textbox(slide, x - 0.04, y + 0.2, 0.62, 0.12, label, 5.2, PPT_THEME["navy"], True, PP_ALIGN.CENTER)
        add_shape_block(slide, 6.76, 5.0, 2.3, 0.54, PPT_THEME["ice"], PPT_THEME["cyan"], MSO_SHAPE.ROUNDED_RECTANGLE)
        add_textbox(slide, 6.96, 5.16, 1.85, 0.16, "LIVE OPERATIONS LOOP", 7, PPT_THEME["navy"], True, PP_ALIGN.CENTER)
    else:
        add_shape_block(slide, 6.62, 1.14, 2.58, 0.4, PPT_THEME["navy"], PPT_THEME["blue"], MSO_SHAPE.ROUNDED_RECTANGLE)
        add_textbox(slide, 6.82, 1.26, 1.72, 0.15, "Growth Dashboard", 7.5, PPT_THEME["cyan"], True)
        add_shape_block(slide, 6.62, 1.72, 0.52, 3.62, PPT_THEME["navy"], PPT_THEME["blue"], MSO_SHAPE.ROUNDED_RECTANGLE)
        for i in range(5):
            add_shape_block(slide, 6.8, 2.0 + (i * 0.48), 0.18, 0.18, PPT_THEME["cyan"], PPT_THEME["cyan"], MSO_SHAPE.OVAL)
        for i, (x, y) in enumerate([(7.38, 1.74), (8.28, 1.74), (7.38, 2.44), (8.28, 2.44)]):
            add_shape_block(slide, x, y, 0.72, 0.46, PPT_THEME["white"], PPT_THEME["border"], MSO_SHAPE.ROUNDED_RECTANGLE)
            add_block(slide, x + 0.12, y + 0.16, 0.36 + (i * 0.06), 0.05, PPT_THEME["blue"])
            add_block(slide, x + 0.12, y + 0.29, 0.26, 0.04, PPT_THEME["cyan"])
        add_shape_block(slide, 7.36, 3.24, 1.8, 1.1, PPT_THEME["white"], PPT_THEME["border"], MSO_SHAPE.ROUNDED_RECTANGLE)
        for i, h in enumerate([0.34, 0.56, 0.42, 0.74, 0.62]):
            add_block(slide, 7.55 + (i * 0.25), 4.04 - h, 0.12, h, PPT_THEME["blue"] if i % 2 else PPT_THEME["cyan"])
        add_shape_block(slide, 7.36, 4.6, 1.8, 0.52, PPT_THEME["ice"], PPT_THEME["cyan"], MSO_SHAPE.ROUNDED_RECTANGLE)
        add_textbox(slide, 7.68, 4.76, 1.16, 0.16, "ARR WORKFLOW", 7, PPT_THEME["navy"], True, PP_ALIGN.CENTER)


def add_theme_content_visual(slide, industry):
    try:
        add_block(slide, 0.22, 1.28, 0.08, 4.95, PPT_THEME["blue"])
        add_block(slide, 0.34, 1.28, 0.035, 4.95, PPT_THEME["cyan"])
        if industry == "Marketplace":
            for x, y in [(8.86, 1.5), (9.34, 2.08), (8.92, 2.74), (9.38, 3.34)]:
                add_shape_block(slide, x, y, 0.15, 0.15, PPT_THEME["cyan"], PPT_THEME["cyan"], MSO_SHAPE.OVAL)
            add_shape_block(slide, 8.96, 2.14, 0.42, 0.03, PPT_THEME["blue"])
            add_shape_block(slide, 9.0, 2.8, 0.36, 0.03, PPT_THEME["cyan"])
        elif industry == "Consumer Product":
            add_shape_block(slide, 8.72, 1.42, 0.72, 1.0, PPT_THEME["white"], PPT_THEME["border"], MSO_SHAPE.ROUNDED_RECTANGLE)
            add_shape_block(slide, 8.86, 1.62, 0.44, 0.58, PPT_THEME["ice"], PPT_THEME["cyan"], MSO_SHAPE.ROUNDED_RECTANGLE)
            add_block(slide, 8.72, 2.78, 0.72, 0.05, PPT_THEME["cyan"])
            add_block(slide, 8.86, 3.02, 0.5, 0.04, PPT_THEME["blue"])
        elif industry == "Food / Delivery":
            for x, y, w, rot in [(8.64, 1.5, 0.75, 18), (8.82, 2.05, 0.62, -24), (8.58, 2.72, 0.78, 18)]:
                route = add_shape_block(slide, x, y, w, 0.05, PPT_THEME["cyan"])
                route.rotation = rot
            add_shape_block(slide, 9.2, 3.14, 0.22, 0.22, PPT_THEME["blue"], PPT_THEME["blue"], MSO_SHAPE.OVAL)
        elif industry == "AI Startup":
            for i in range(4):
                add_shape_block(slide, 8.62 + (i * 0.22), 1.34, 0.015, 1.3, PPT_THEME["blue"])
                add_shape_block(slide, 8.56, 1.48 + (i * 0.26), 1.0, 0.015, PPT_THEME["cyan"])
            add_shape_block(slide, 8.9, 2.22, 0.18, 0.18, PPT_THEME["cyan"], PPT_THEME["cyan"], MSO_SHAPE.OVAL)
        else:
            for i in range(3):
                add_shape_block(slide, 8.66 + (i * 0.28), 1.42, 0.2, 0.2, PPT_THEME["ice"], PPT_THEME["border"], MSO_SHAPE.ROUNDED_RECTANGLE)
            add_block(slide, 8.66, 1.92, 0.88, 0.05, PPT_THEME["cyan"])
    except Exception:
        pass


def add_industry_motif(slide, industry, surface="title"):
    if surface == "title":
        add_theme_cover_visual(slide, industry)
    else:
        add_theme_content_visual(slide, industry)


def add_logo(slide, x=8.65, y=6.58, width=0.62):
    logo_path = "turbopitch_logo.png"
    if os.path.exists(logo_path):
        try:
            slide.shapes.add_picture(logo_path, PPTInches(x), PPTInches(y), width=PPTInches(width))
        except Exception:
            pass


def add_dark_header(slide, title, eyebrow=None):
    add_block(slide, 0, 0, 10, 1.02, PPT_THEME["navy"])
    add_block(slide, 0, 1.02, 10, 0.08, PPT_THEME["blue"])
    header = slide.shapes.add_textbox(PPTInches(0.55), PPTInches(0.18), PPTInches(8.35), PPTInches(0.65))
    tf = header.text_frame
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0

    if eyebrow:
        p = tf.paragraphs[0]
        p.text = eyebrow.upper()
        p.font.size = Pt(8)
        p.font.bold = True
        p.font.color.rgb = PPT_THEME["cyan"]
        p.space_after = Pt(2)
        title_p = tf.add_paragraph()
    else:
        title_p = tf.paragraphs[0]

    title_p.text = title
    title_p.font.size = Pt(22)
    title_p.font.bold = True
    title_p.font.color.rgb = PPT_THEME["white"]


def add_footer(slide, dark=False, branded=False):
    if not branded:
        return
    footer = slide.shapes.add_textbox(PPTInches(0.55), PPTInches(6.86), PPTInches(4.8), PPTInches(0.22))
    tf = footer.text_frame
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = "TurboPitch"
    p.font.size = Pt(7)
    p.font.bold = False
    p.font.color.rgb = PPT_THEME["pale_blue"] if dark else PPT_THEME["muted"]


def add_metric_card(slide, x, y, w, h, year_label, metrics, accent_color=None):
    card = add_shape_block(slide, x, y, w, h, PPT_THEME["white"], PPT_THEME["border"], MSO_SHAPE.ROUNDED_RECTANGLE)
    add_block(slide, x, y, 0.08, h, accent_color or PPT_THEME["blue"])
    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = PPTInches(0.2)
    tf.margin_right = PPTInches(0.14)
    tf.margin_top = PPTInches(0.15)
    tf.margin_bottom = PPTInches(0.12)

    p = tf.paragraphs[0]
    p.text = year_label
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = PPT_THEME["navy"]
    p.space_after = Pt(8)

    for label, value in metrics:
        metric_p = tf.add_paragraph()
        metric_p.text = label
        metric_p.font.size = Pt(7.5)
        metric_p.font.bold = True
        metric_p.font.color.rgb = PPT_THEME["muted"]
        metric_p.space_after = Pt(0)

        value_p = tf.add_paragraph()
        value_p.text = value
        value_p.font.size = Pt(13)
        value_p.font.bold = True
        value_p.font.color.rgb = PPT_THEME["navy"]
        value_p.space_after = Pt(5)


def add_bullet_panel(slide, x, y, w, h, title, bullets, fill_color=None, title_color=None):
    panel = add_shape_block(slide, x, y, w, h, fill_color or PPT_THEME["white"], PPT_THEME["border"], MSO_SHAPE.ROUNDED_RECTANGLE)
    add_block(slide, x, y, w, 0.08, PPT_THEME["cyan"])
    tf = panel.text_frame
    tf.word_wrap = True
    tf.margin_left = PPTInches(0.2)
    tf.margin_right = PPTInches(0.18)
    tf.margin_top = PPTInches(0.18)
    tf.margin_bottom = PPTInches(0.1)

    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = title_color or PPT_THEME["navy"]
    p.space_after = Pt(4)

    for item in bullets[:5]:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(8.8)
        p.font.color.rgb = PPT_THEME["slate"]
        p.level = 0
        p.space_after = Pt(3)


def add_bullet_card(slide, bullets, x=0.62, y=1.42, w=4.15, h=4.95):
    row_items = (bullets[:6] if bullets else ["No content available."])
    add_shape_block(slide, x, y, w, h, PPT_THEME["white"], PPT_THEME["border"], MSO_SHAPE.ROUNDED_RECTANGLE)
    add_block(slide, x, y, w, 0.1, PPT_THEME["blue"])
    add_block(slide, x + 0.22, y + 0.34, 0.62, 0.05, PPT_THEME["cyan"])
    add_textbox(slide, x + 0.22, y + 0.48, w - 0.44, 0.22, "Key points", 8.2, PPT_THEME["muted"], True)

    content = slide.shapes.add_textbox(
        PPTInches(x + 0.32),
        PPTInches(y + 0.84),
        PPTInches(w - 0.64),
        PPTInches(h - 1.08),
    )
    tf = content.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0

    for i, item in enumerate(row_items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(9.4 if len(row_items) <= 4 else 8.8)
        p.font.color.rgb = PPT_THEME["slate"]
        p.space_after = Pt(6 if len(row_items) <= 4 else 4)
        p.level = 0

    add_shape_block(slide, x + w - 0.62, y + h - 0.58, 0.34, 0.34, PPT_THEME["ice"], PPT_THEME["border"], MSO_SHAPE.OVAL)


def get_slide_kind(title):
    title_text = str(title or "").lower()
    if "problem" in title_text or "pain" in title_text:
        return "problem"
    if "solution" in title_text:
        return "solution"
    if "product" in title_text:
        return "product"
    if "go-to-market" in title_text or "go to market" in title_text or "gtm" in title_text or "market strategy" in title_text:
        return "gtm"
    if "market" in title_text or "opportunity" in title_text:
        return "market"
    if "business model" in title_text or "revenue" in title_text or "monetization" in title_text:
        return "business_model"
    if "competitive" in title_text or "advantage" in title_text or "moat" in title_text:
        return "competitive"
    if "financial" in title_text:
        return "financial"
    if "assumption" in title_text or "investor interpretation" in title_text or "interpretation" in title_text:
        return "investor_lens"
    if "funding" in title_text or "ask" in title_text or "use of funds" in title_text:
        return "funding"
    return "generic"


def get_bullet_keywords(bullets, max_items=3):
    stop_words = {
        "about", "above", "across", "after", "again", "against", "also", "because", "being",
        "build", "built", "business", "company", "could", "customer", "customers", "early",
        "from", "have", "into", "market", "more", "need", "needs", "offer", "platform",
        "product", "solution", "startup", "than", "that", "their", "there", "these", "this",
        "through", "with", "without", "would", "while", "will", "using",
    }
    keywords = []
    for bullet in bullets or []:
        text = re.sub(r"^[\s\-*•]+", "", str(bullet or "")).strip()
        text = re.sub(r"[^A-Za-z0-9%$ ]+", " ", text)
        words = [word for word in text.split() if len(word) > 2 and word.lower() not in stop_words]
        phrase = " ".join(words[:3]).strip()
        if phrase and phrase.lower() not in [item.lower() for item in keywords]:
            keywords.append(phrase[:32])
        if len(keywords) >= max_items:
            break
    fallback = ["Signal", "Risk", "Evidence"]
    for item in fallback:
        if len(keywords) >= max_items:
            break
        keywords.append(item)
    return keywords[:max_items]


def shorten_chip_text(text, max_words=6, max_chars=36):
    clean_text = re.sub(r"\s+", " ", str(text or "")).strip()
    clean_text = re.sub(r"^[\-*â€¢\s]+", "", clean_text)
    clean_text = re.sub(r"\bhampers\b", "slows", clean_text, flags=re.IGNORECASE)
    clean_text = clean_text.rstrip(" .,:;")
    if len(clean_text) <= max_chars and len(clean_text.split()) <= max_words:
        return clean_text
    words = clean_text.split()
    target_words = min(max(max_words, 4), 7)
    shortened_words = words[:target_words]
    shortened = " ".join(shortened_words).strip(" .,:;")
    if len(shortened) > max_chars:
        shortened = shortened[:max_chars].rsplit(" ", 1)[0].strip(" .,:;")
    while shortened and len(shortened.split()) < 4 and len(shortened_words) < len(words):
        shortened_words.append(words[len(shortened_words)])
        shortened = " ".join(shortened_words).strip(" .,:;")
        if len(shortened) > max_chars:
            shortened = " ".join(shortened_words[:-1]).strip(" .,:;")
            break
    return f"{shortened}..." if shortened else clean_text[: max_chars - 3].strip(" .,:;") + "..."


def add_panel_chip(slide, x, y, w, h, text, fill=None):
    chip = add_shape_block(slide, x, y, w, h, fill or PPT_THEME["ice"], PPT_THEME["border"], MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = chip.text_frame
    tf.word_wrap = True
    tf.margin_left = PPTInches(0.1)
    tf.margin_right = PPTInches(0.08)
    tf.margin_top = PPTInches(0.055)
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = shorten_chip_text(text, 6, 34)
    p.font.size = Pt(8.5)
    p.font.bold = True
    p.font.color.rgb = PPT_THEME["navy"]
    p.alignment = PP_ALIGN.CENTER
    return chip


def add_insight_panel_shell(slide, x, y, w, h, heading):
    panel = add_shape_block(slide, x, y, w, h, PPT_THEME["white"], PPT_THEME["border"], MSO_SHAPE.ROUNDED_RECTANGLE)
    add_block(slide, x, y, w, 0.1, PPT_THEME["cyan"])
    add_textbox(slide, x + 0.24, y + 0.24, w - 0.48, 0.28, heading, 12, PPT_THEME["navy"], True)
    add_block(slide, x + 0.24, y + 0.62, 0.72, 0.05, PPT_THEME["blue"])
    return panel


def add_industry_panel_texture(slide, x, y, w, h, industry):
    if industry == "AI Startup":
        for i in range(4):
            add_block(slide, x + 0.28 + (i * 0.38), y + h - 0.92, 0.012, 0.54, PPT_THEME["border"])
            add_block(slide, x + 0.22, y + h - 0.82 + (i * 0.12), 1.62, 0.012, PPT_THEME["border"])
        for nx, ny in [(x + 2.75, y + 3.82), (x + 3.24, y + 4.12), (x + 3.62, y + 3.76)]:
            add_shape_block(slide, nx, ny, 0.12, 0.12, PPT_THEME["cyan"], PPT_THEME["cyan"], MSO_SHAPE.OVAL)
    elif industry == "Marketplace":
        for nx, ny in [(x + 2.72, y + 3.78), (x + 3.32, y + 4.1), (x + 3.72, y + 3.74)]:
            add_shape_block(slide, nx, ny, 0.16, 0.16, PPT_THEME["cyan"], PPT_THEME["cyan"], MSO_SHAPE.OVAL)
        add_block(slide, x + 2.84, y + 3.88, 0.56, 0.025, PPT_THEME["blue"])
        add_block(slide, x + 3.42, y + 3.98, 0.36, 0.025, PPT_THEME["cyan"])
    elif industry == "Food / Delivery":
        route = add_block(slide, x + 2.55, y + 4.0, 1.0, 0.05, PPT_THEME["cyan"])
        route.rotation = -16
        add_shape_block(slide, x + 3.58, y + 3.78, 0.18, 0.18, PPT_THEME["blue"], PPT_THEME["blue"], MSO_SHAPE.OVAL)
    elif industry == "Consumer Product":
        add_shape_block(slide, x + 2.72, y + 3.58, 0.54, 0.72, PPT_THEME["ice"], PPT_THEME["border"], MSO_SHAPE.ROUNDED_RECTANGLE)
        add_shape_block(slide, x + 3.4, y + 3.72, 0.42, 0.44, PPT_THEME["white"], PPT_THEME["cyan"], MSO_SHAPE.ROUNDED_RECTANGLE)
    else:
        for i in range(3):
            add_shape_block(slide, x + 2.58 + (i * 0.42), y + 3.78, 0.32, 0.32, PPT_THEME["ice"], PPT_THEME["border"], MSO_SHAPE.ROUNDED_RECTANGLE)
        add_block(slide, x + 2.62, y + 4.28, 1.24, 0.05, PPT_THEME["cyan"])


def add_slide_insight_panel(slide, title, bullets, industry):
    x, y, w, h = 5.05, 1.42, 4.25, 4.95
    kind = get_slide_kind(title)
    keywords = get_bullet_keywords(bullets, 3)

    headings = {
        "problem": "Core friction",
        "solution": "How it works",
        "product": "Product modules",
        "market": "Market lens",
        "business_model": "Revenue engine",
        "competitive": "Differentiation stack",
        "gtm": "GTM motion",
        "financial": "Financial signal",
        "investor_lens": "Investor lens",
        "funding": "Use of funds",
        "generic": "Investor lens",
    }
    add_insight_panel_shell(slide, x, y, w, h, headings.get(kind, "Investor lens"))

    if kind == "problem":
        warning = add_shape_block(slide, x + 0.28, y + 0.92, 0.58, 0.52, PPT_THEME["yellow"], PPT_THEME["cyan"], MSO_SHAPE.ISOSCELES_TRIANGLE)
        warning.rotation = 0
        add_textbox(slide, x + 0.48, y + 1.08, 0.16, 0.16, "!", 10, PPT_THEME["navy"], True, PP_ALIGN.CENTER)
        for i, label in enumerate(keywords):
            add_panel_chip(slide, x + 1.02, y + 0.88 + (i * 0.58), 2.82, 0.38, label)
    elif kind == "solution":
        for i, label in enumerate(keywords):
            step_x = x + 0.32 + (i * 1.25)
            add_shape_block(slide, step_x, y + 1.12, 0.92, 1.05, PPT_THEME["ice"], PPT_THEME["border"], MSO_SHAPE.ROUNDED_RECTANGLE)
            add_shape_block(slide, step_x + 0.29, y + 1.28, 0.34, 0.34, PPT_THEME["blue"], PPT_THEME["blue"], MSO_SHAPE.OVAL)
            add_textbox(slide, step_x + 0.38, y + 1.38, 0.16, 0.12, str(i + 1), 6.8, PPT_THEME["white"], True, PP_ALIGN.CENTER)
            add_textbox(slide, step_x + 0.12, y + 1.76, 0.68, 0.24, label, 6.5, PPT_THEME["navy"], True, PP_ALIGN.CENTER)
            if i < 2:
                add_block(slide, step_x + 0.96, y + 1.62, 0.27, 0.035, PPT_THEME["cyan"])
    elif kind == "product":
        for i, label in enumerate(keywords):
            card_y = y + 0.98 + (i * 0.9)
            add_shape_block(slide, x + 0.3, card_y, 3.62, 0.68, PPT_THEME["ice"] if i % 2 == 0 else PPT_THEME["white"], PPT_THEME["border"], MSO_SHAPE.ROUNDED_RECTANGLE)
            add_block(slide, x + 0.52, card_y + 0.2, 0.48, 0.06, PPT_THEME["blue"])
            add_block(slide, x + 0.52, card_y + 0.36, 0.72, 0.04, PPT_THEME["cyan"])
            add_textbox(slide, x + 1.42, card_y + 0.22, 2.0, 0.22, label, 7.4, PPT_THEME["navy"], True)
    elif kind == "market":
        market_labels = ["TAM", "SAM", "SOM"]
        for i, label in enumerate(market_labels):
            card_x = x + 0.32 + (i * 1.25)
            add_shape_block(slide, card_x, y + 1.08, 0.95, 1.42, PPT_THEME["ice"], PPT_THEME["border"], MSO_SHAPE.ROUNDED_RECTANGLE)
            add_textbox(slide, card_x + 0.14, y + 1.34, 0.66, 0.22, label, 12, PPT_THEME["blue"], True, PP_ALIGN.CENTER)
            add_textbox(slide, card_x + 0.12, y + 1.78, 0.7, 0.3, keywords[i], 6.2, PPT_THEME["muted"], True, PP_ALIGN.CENTER)
    elif kind == "business_model":
        revenue_terms = []
        content = " ".join(str(item) for item in bullets or [])
        for term in ["pricing", "subscription", "recurring", "take rate", "margin", "revenue"]:
            if term in content.lower():
                revenue_terms.append(term.title())
        labels = (revenue_terms + keywords)[:3]
        for i, label in enumerate(labels):
            add_panel_chip(slide, x + 0.36, y + 1.0 + (i * 0.68), 3.46, 0.44, label, PPT_THEME["ice"])
    elif kind == "competitive":
        for i, label in enumerate(keywords):
            stack_y = y + 1.08 + (i * 0.78)
            add_shape_block(slide, x + 0.44 + (i * 0.12), stack_y, 3.28 - (i * 0.24), 0.58, PPT_THEME["white"], PPT_THEME["border"], MSO_SHAPE.ROUNDED_RECTANGLE)
            add_block(slide, x + 0.58 + (i * 0.12), stack_y + 0.17, 0.42, 0.06, PPT_THEME["blue"])
            add_textbox(slide, x + 1.16 + (i * 0.12), stack_y + 0.18, 1.92, 0.2, label, 7.2, PPT_THEME["navy"], True)
    elif kind == "gtm":
        for i, label in enumerate(["Acquire", "Convert", "Retain"]):
            step_y = y + 1.02 + (i * 0.74)
            add_shape_block(slide, x + 0.36, step_y, 3.4, 0.54, PPT_THEME["ice"], PPT_THEME["border"], MSO_SHAPE.ROUNDED_RECTANGLE)
            add_shape_block(slide, x + 0.56, step_y + 0.16, 0.22, 0.22, PPT_THEME["cyan"], PPT_THEME["cyan"], MSO_SHAPE.OVAL)
            add_textbox(slide, x + 0.96, step_y + 0.15, 0.86, 0.18, label, 7.6, PPT_THEME["navy"], True)
            add_textbox(slide, x + 1.86, step_y + 0.15, 1.36, 0.18, keywords[i], 6.4, PPT_THEME["muted"], True)
    elif kind == "financial":
        final_projection = projection_df.iloc[-1]
        metric_labels = [
            ("Revenue", final_projection["Revenue"]),
            ("Net Income", final_projection["Net Income"]),
            ("Ending Cash", final_projection["Ending Cash"]),
        ]
        for i, (label, value) in enumerate(metric_labels):
            card_y = y + 0.98 + (i * 0.78)
            add_shape_block(slide, x + 0.34, card_y, 3.5, 0.58, PPT_THEME["ice"], PPT_THEME["border"], MSO_SHAPE.ROUNDED_RECTANGLE)
            add_textbox(slide, x + 0.56, card_y + 0.12, 1.15, 0.16, label.upper(), 6.2, PPT_THEME["muted"], True)
            add_textbox(slide, x + 2.14, card_y + 0.11, 1.2, 0.2, f"${value / 1000:,.0f}K", 10.4, PPT_THEME["navy"], True, PP_ALIGN.RIGHT)
    elif kind == "investor_lens":
        for i, label in enumerate(["Credibility", "Risk", "Evidence needed"]):
            add_panel_chip(slide, x + 0.36, y + 1.0 + (i * 0.64), 3.46, 0.42, label)
    elif kind == "funding":
        for i, label in enumerate(keywords):
            card_x = x + 0.32 + (i * 1.25)
            add_shape_block(slide, card_x, y + 1.08, 0.95, 1.28, PPT_THEME["ice"], PPT_THEME["border"], MSO_SHAPE.ROUNDED_RECTANGLE)
            add_textbox(slide, card_x + 0.14, y + 1.3, 0.66, 0.18, f"Phase {i + 1}", 6.4, PPT_THEME["muted"], True, PP_ALIGN.CENTER)
            add_textbox(slide, card_x + 0.1, y + 1.7, 0.74, 0.3, label, 6.2, PPT_THEME["navy"], True, PP_ALIGN.CENTER)
    else:
        for i, label in enumerate(keywords):
            add_panel_chip(slide, x + 0.36, y + 1.0 + (i * 0.64), 3.46, 0.42, label)

    add_industry_panel_texture(slide, x, y, w, h, industry)


def add_content_slide_frame(slide, title, eyebrow):
    add_block(slide, 0, 0, 10, 7.5, PPT_THEME["soft"])
    add_theme_content_visual(slide, selected_ppt_industry)
    add_block(slide, 0, 0, 10, 0.16, PPT_THEME["navy"])
    add_textbox(slide, 0.62, 0.26, 1.45, 0.24, eyebrow.upper(), 7.5, PPT_THEME["blue"], True)
    add_industry_badge(slide, 7.36, 0.28, get_industry_badge_label(selected_ppt_industry))
    add_textbox(slide, 0.62, 0.52, 8.25, 0.44, title, 20, PPT_THEME["navy"], True)
    add_block(slide, 0.62, 1.0, 0.95, 0.05, PPT_THEME["blue"])
    add_block(slide, 1.64, 1.0, 0.36, 0.05, PPT_THEME["cyan"])


def add_pitch_text_card(slide, x, y, w, h, text, accent_color=None, fill_color=None, size=8.0, bold=False):
    card = add_shape_block(
        slide,
        x,
        y,
        w,
        h,
        fill_color or PPT_THEME["white"],
        PPT_THEME["border"],
        MSO_SHAPE.ROUNDED_RECTANGLE,
    )
    add_block(slide, x, y, 0.07, h, accent_color or PPT_THEME["blue"])
    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = PPTInches(0.18)
    tf.margin_right = PPTInches(0.12)
    tf.margin_top = PPTInches(0.12)
    tf.margin_bottom = PPTInches(0.08)
    p = tf.paragraphs[0]
    p.text = str(text or "")
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = PPT_THEME["slate"]
    return card


def add_pitch_visual_panel(slide, x, y, w, h, industry, dark=False):
    if industry == "SaaS":
        add_shape_block(slide, x, y, w, h, PPT_THEME["navy_2"], PPT_THEME["cyan"], MSO_SHAPE.ROUNDED_RECTANGLE)
        add_saas_image_accent(slide, x + 0.08, y + 0.08, w - 0.16, h - 0.16)
        add_shape_block(slide, x + 0.22, y + 0.38, w - 0.88, h - 0.82, PPT_THEME["white"], PPT_THEME["border"], MSO_SHAPE.ROUNDED_RECTANGLE)
        add_block(slide, x + 0.22, y + 0.38, w - 0.88, 0.28, PPT_THEME["navy"])
        for i, color in enumerate([PPT_THEME["cyan"], PPT_THEME["blue"], PPT_THEME["pale_blue"]]):
            add_shape_block(slide, x + 0.38 + (i * 0.16), y + 0.48, 0.07, 0.07, color, color, MSO_SHAPE.OVAL)
        for i, hgt in enumerate([0.38, 0.64, 0.52, 0.88, 0.74]):
            add_block(slide, x + 0.58 + (i * 0.26), y + h - 0.62 - hgt, 0.12, hgt, PPT_THEME["blue"] if i % 2 else PPT_THEME["cyan"])
        add_shape_block(slide, x + w - 1.1, y + 0.92, 0.56, 1.42, PPT_THEME["ice"], PPT_THEME["border"], MSO_SHAPE.ROUNDED_RECTANGLE)
        for i in range(3):
            add_block(slide, x + w - 0.96, y + 1.18 + (i * 0.3), 0.32 + (i * 0.06), 0.04, PPT_THEME["cyan"] if i != 1 else PPT_THEME["blue"])
        return True

    image_path = None if industry == "SaaS" else get_theme_image_path(industry)
    if image_path:
        add_shape_block(slide, x, y, w, h, PPT_THEME["white"], PPT_THEME["border"], MSO_SHAPE.ROUNDED_RECTANGLE)
        if add_safe_picture(slide, image_path, x + 0.08, y + 0.08, w - 0.16, h - 0.16):
            return True

    fill = PPT_THEME["navy_2"] if dark else PPT_THEME["white"]
    line = PPT_THEME["cyan"] if dark else PPT_THEME["border"]
    add_shape_block(slide, x, y, w, h, fill, line, MSO_SHAPE.ROUNDED_RECTANGLE)
    return False


def add_hero_insight_layout(slide, title, bullets, industry):
    kind = get_slide_kind(title)
    row_items = bullets[:4] if bullets else ["No content available."]
    chip_items = bullets[:3] if bullets else row_items[:3]
    section_label = "Market Pain" if kind == "problem" else "Market Signal" if kind == "market" else "Investor Signal"

    add_textbox(slide, 0.74, 1.34, 2.2, 0.26, section_label, 12, PPT_THEME["muted"], True)
    add_block(slide, 0.74, 1.68, 0.86, 0.05, PPT_THEME["blue"])
    add_block(slide, 1.7, 1.68, 0.32, 0.05, PPT_THEME["cyan"])

    content = slide.shapes.add_textbox(PPTInches(0.76), PPTInches(1.94), PPTInches(4.35), PPTInches(2.72))
    tf = content.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, item in enumerate(row_items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = str(item or "")
        p.font.size = Pt(11.4 if len(row_items) <= 3 else 10.5)
        p.font.color.rgb = PPT_THEME["slate"]
        p.space_after = Pt(7)

    image_used = add_pitch_visual_panel(slide, 5.48, 1.38, 3.85, 3.18, industry, dark=industry in ("SaaS", "Marketplace", "AI Startup"))
    if not image_used:
        if industry == "Marketplace":
            nodes = [(6.42, 2.48), (7.34, 2.04), (8.16, 2.64), (7.12, 3.34), (8.04, 3.62)]
            for x, y in nodes:
                add_shape_block(slide, x, y, 0.34, 0.34, PPT_THEME["cyan"], PPT_THEME["cyan"], MSO_SHAPE.OVAL)
            for x, y, w, rot in [(6.72, 2.6, 0.78, -22), (7.52, 2.22, 0.8, 26), (7.28, 3.48, 0.88, 12)]:
                line = add_block(slide, x, y, w, 0.035, PPT_THEME["blue"])
                line.rotation = rot
        elif industry == "Food / Delivery":
            for i, label in enumerate(["ORDER", "PREP", "DELIVER"]):
                step_x = 5.92 + (i * 1.05)
                add_shape_block(slide, step_x, 2.54 + (i % 2) * 0.34, 0.72, 0.48, PPT_THEME["ice"], PPT_THEME["border"], MSO_SHAPE.ROUNDED_RECTANGLE)
                add_textbox(slide, step_x + 0.08, 2.72 + (i % 2) * 0.34, 0.56, 0.12, label, 5.5, PPT_THEME["navy"], True, PP_ALIGN.CENTER)
                if i < 2:
                    connector = add_block(slide, step_x + 0.76, 2.78 + (i % 2) * 0.34, 0.36, 0.04, PPT_THEME["cyan"])
                    connector.rotation = 16 if i == 0 else -16
        else:
            for i, height in enumerate([1.05, 1.48, 0.86, 1.72, 1.2]):
                add_block(slide, 6.04 + (i * 0.48), 3.78 - height, 0.24, height, PPT_THEME["cyan"] if i % 2 else PPT_THEME["blue"])
            add_shape_block(slide, 7.98, 2.02, 0.8, 0.8, PPT_THEME["ice"], PPT_THEME["cyan"], MSO_SHAPE.OVAL)
            add_shape_block(slide, 8.18, 2.22, 0.4, 0.4, PPT_THEME["blue"], PPT_THEME["blue"], MSO_SHAPE.OVAL)

    for i, item in enumerate(chip_items[:3]):
        add_pitch_text_card(
            slide,
            0.74 + (i * 2.86),
            5.08,
            2.58,
            1.02,
            shorten_chip_text(item, 7, 46),
            PPT_THEME["cyan"] if i == 1 else PPT_THEME["blue"],
            PPT_THEME["ice"] if i == 1 else PPT_THEME["white"],
            8.8,
            True,
        )


def add_process_flow_layout(slide, title, bullets, industry):
    kind = get_slide_kind(title)
    labels = ["Acquire", "Convert", "Retain"] if kind == "gtm" else ["Input", "Platform", "Outcome"]
    row_items = bullets[:3] if bullets else ["No content available."]
    while len(row_items) < 3:
        row_items.append("")

    section_label = "GTM Motion" if kind == "gtm" else "Solution Workflow"
    add_textbox(slide, 0.74, 1.32, 5.65, 0.28, section_label, 12, PPT_THEME["muted"], True)

    for i, label in enumerate(labels):
        x = 0.78 + (i * 3.0)
        y = 2.06 + (0.22 if i == 1 else 0)
        add_shape_block(slide, x, y, 2.2, 2.7, PPT_THEME["white"], PPT_THEME["border"], MSO_SHAPE.ROUNDED_RECTANGLE)
        add_block(slide, x, y, 2.2, 0.1, PPT_THEME["blue"] if i != 1 else PPT_THEME["cyan"])
        add_shape_block(slide, x + 0.18, y + 0.28, 0.44, 0.44, PPT_THEME["ice"], PPT_THEME["cyan"], MSO_SHAPE.OVAL)
        add_textbox(slide, x + 0.32, y + 0.42, 0.16, 0.12, str(i + 1), 7, PPT_THEME["navy"], True, PP_ALIGN.CENTER)
        add_textbox(slide, x + 0.72, y + 0.3, 1.18, 0.28, label, 13, PPT_THEME["navy"], True)
        add_textbox(slide, x + 0.24, y + 0.92, 1.72, 1.18, row_items[i], 9.2, PPT_THEME["slate"])
        add_block(slide, x + 0.24, y + 2.28, 0.74, 0.05, PPT_THEME["cyan"])
        if i < 2:
            arrow = add_shape_block(slide, x + 2.28, y + 1.08, 0.5, 0.34, PPT_THEME["cyan"], PPT_THEME["cyan"], MSO_SHAPE.RIGHT_ARROW)
            arrow.rotation = 6 if i == 0 else -6

    for i, item in enumerate((bullets[3:5] if len(bullets) > 3 else bullets[:2])):
        add_panel_chip(slide, 2.25 + (i * 2.8), 5.62, 2.3, 0.48, item, PPT_THEME["ice"])


def add_product_dashboard_layout(slide, title, bullets, industry):
    row_items = bullets[:3] if bullets else ["No content available."]
    while len(row_items) < 3:
        row_items.append("")

    add_textbox(slide, 0.74, 1.32, 2.0, 0.26, "Product View", 12, PPT_THEME["muted"], True)

    for i, item in enumerate(row_items[:3]):
        add_pitch_text_card(slide, 0.76, 1.82 + (i * 1.38), 2.7, 1.08, item, PPT_THEME["blue"] if i != 1 else PPT_THEME["cyan"], PPT_THEME["white"], 9.2, True)

    if industry == "SaaS":
        add_saas_image_accent(slide, 3.96, 1.44, 5.44, 5.0)
    add_shape_block(slide, 4.08, 1.58, 5.16, 4.7, PPT_THEME["white"], PPT_THEME["border"], MSO_SHAPE.ROUNDED_RECTANGLE)
    add_block(slide, 4.08, 1.58, 5.16, 0.36, PPT_THEME["navy"])
    add_shape_block(slide, 4.32, 1.72, 0.12, 0.12, PPT_THEME["cyan"], PPT_THEME["cyan"], MSO_SHAPE.OVAL)
    add_shape_block(slide, 4.54, 1.72, 0.12, 0.12, PPT_THEME["pale_blue"], PPT_THEME["pale_blue"], MSO_SHAPE.OVAL)
    header_label = "Revenue Dashboard" if industry == "SaaS" else "PRODUCT VIEW"
    add_textbox(slide, 7.0, 1.68, 1.72, 0.16, header_label, 6.2, PPT_THEME["cyan"], True, PP_ALIGN.RIGHT)

    if industry == "Marketplace":
        add_shape_block(slide, 6.22, 2.46, 0.78, 0.78, PPT_THEME["cyan"], PPT_THEME["cyan"], MSO_SHAPE.OVAL)
        add_textbox(slide, 6.31, 2.76, 0.6, 0.12, "CORE", 5.6, PPT_THEME["navy"], True, PP_ALIGN.CENTER)
        for x, y, label in [(4.82, 2.28, "Demand"), (7.58, 2.34, "Supply"), (5.08, 4.68, "Ops"), (7.36, 4.76, "Payments")]:
            add_shape_block(slide, x, y, 0.88, 0.54, PPT_THEME["ice"], PPT_THEME["border"], MSO_SHAPE.ROUNDED_RECTANGLE)
            add_textbox(slide, x + 0.08, y + 0.2, 0.72, 0.12, label, 5.8, PPT_THEME["navy"], True, PP_ALIGN.CENTER)
        for x, y, w, rot in [(5.62, 2.62, 0.78, -8), (6.94, 2.66, 0.82, 8), (5.76, 4.82, 0.98, -28), (6.74, 4.88, 0.86, 24)]:
            line = add_block(slide, x, y, w, 0.035, PPT_THEME["cyan"])
            line.rotation = rot
    elif industry == "Food / Delivery":
        for i, label in enumerate(["Order", "Kitchen", "Route", "Door"]):
            x = 4.72 + (i * 1.0)
            y = 2.48 + (0.48 if i % 2 else 0)
            add_shape_block(slide, x, y, 0.78, 0.56, PPT_THEME["ice"], PPT_THEME["border"], MSO_SHAPE.ROUNDED_RECTANGLE)
            add_textbox(slide, x + 0.08, y + 0.2, 0.62, 0.12, label, 5.7, PPT_THEME["navy"], True, PP_ALIGN.CENTER)
            if i < 3:
                connector = add_block(slide, x + 0.82, y + 0.28, 0.42, 0.045, PPT_THEME["cyan"])
                connector.rotation = 22 if i % 2 == 0 else -22
        add_shape_block(slide, 5.08, 4.62, 3.3, 0.56, PPT_THEME["white"], PPT_THEME["border"], MSO_SHAPE.ROUNDED_RECTANGLE)
        add_block(slide, 5.28, 4.88, 2.02, 0.06, PPT_THEME["blue"])
    elif industry == "SaaS":
        add_shape_block(slide, 4.08, 1.94, 0.58, 4.34, PPT_THEME["ice"], PPT_THEME["border"], MSO_SHAPE.RECTANGLE)
        for i in range(4):
            y = 2.34 + (i * 0.58)
            add_shape_block(slide, 4.28, y, 0.16, 0.16, PPT_THEME["cyan"] if i == 0 else PPT_THEME["blue"], PPT_THEME["cyan"] if i == 0 else PPT_THEME["blue"], MSO_SHAPE.OVAL)
            add_block(slide, 4.22, y + 0.32, 0.28, 0.03, PPT_THEME["border"])

        kpis = [("ARR", "$2.4M"), ("NRR", "118%"), ("Pipeline", "$640K")]
        for i, (label, value) in enumerate(kpis):
            x = 4.96 + (i * 1.22)
            add_shape_block(slide, x, 2.18, 1.08, 0.86, PPT_THEME["ice"], PPT_THEME["border"], MSO_SHAPE.ROUNDED_RECTANGLE)
            add_textbox(slide, x + 0.14, 2.38, 0.78, 0.12, label, 6.4, PPT_THEME["muted"], True)
            add_textbox(slide, x + 0.14, 2.66, 0.78, 0.18, value, 11.0, PPT_THEME["navy"], True)

        add_shape_block(slide, 4.96, 3.28, 2.52, 2.14, PPT_THEME["white"], PPT_THEME["border"], MSO_SHAPE.ROUNDED_RECTANGLE)
        add_textbox(slide, 5.16, 3.5, 1.42, 0.15, "Revenue Trend", 7.2, PPT_THEME["navy"], True)
        add_block(slide, 5.18, 5.16, 1.92, 0.025, PPT_THEME["border"])
        for i, h in enumerate([0.42, 0.66, 0.54, 0.84, 1.02, 1.16]):
            add_block(slide, 5.28 + (i * 0.28), 5.14 - h, 0.13, h, PPT_THEME["blue"] if i % 2 else PPT_THEME["cyan"])

        add_shape_block(slide, 7.64, 3.28, 1.28, 2.14, PPT_THEME["ice"], PPT_THEME["border"], MSO_SHAPE.ROUNDED_RECTANGLE)
        add_textbox(slide, 7.84, 3.5, 0.74, 0.15, "Activity", 7.2, PPT_THEME["navy"], True)
        for i, label in enumerate(["Renewal Risk", "Pipeline Alert", "Pricing Approval"]):
            y = 3.88 + (i * 0.46)
            add_shape_block(slide, 7.78, y, 0.13, 0.13, PPT_THEME["cyan"] if i != 1 else PPT_THEME["blue"], PPT_THEME["cyan"] if i != 1 else PPT_THEME["blue"], MSO_SHAPE.OVAL)
            add_textbox(slide, 7.96, y - 0.02, 0.86, 0.18, label, 5.8, PPT_THEME["slate"], True)
            add_block(slide, 7.96, y + 0.23, 0.5 + (i * 0.08), 0.025, PPT_THEME["border"])
    elif industry == "AI Startup":
        for i, (x, y) in enumerate([(4.52, 2.34), (5.72, 2.34), (6.92, 2.34)]):
            add_shape_block(slide, x, y, 0.9, 0.58, PPT_THEME["ice"], PPT_THEME["border"], MSO_SHAPE.ROUNDED_RECTANGLE)
            add_block(slide, x + 0.16, y + 0.22, 0.48 + (i * 0.06), 0.05, PPT_THEME["blue"])
            add_block(slide, x + 0.16, y + 0.38, 0.3, 0.04, PPT_THEME["cyan"])
        add_shape_block(slide, 4.52, 3.32, 2.1, 1.64, PPT_THEME["white"], PPT_THEME["border"], MSO_SHAPE.ROUNDED_RECTANGLE)
        for i, h in enumerate([0.42, 0.8, 0.58, 1.06, 0.9]):
            add_block(slide, 4.82 + (i * 0.28), 4.72 - h, 0.14, h, PPT_THEME["cyan"] if i % 2 else PPT_THEME["blue"])
        add_shape_block(slide, 6.94, 3.32, 1.7, 1.64, PPT_THEME["ice"], PPT_THEME["border"], MSO_SHAPE.ROUNDED_RECTANGLE)
        for i in range(4):
            add_block(slide, 7.18, 3.7 + (i * 0.25), 0.96 - (i * 0.08), 0.05, PPT_THEME["blue"] if i % 2 else PPT_THEME["cyan"])
    else:
        add_shape_block(slide, 5.02, 2.22, 1.22, 2.6, PPT_THEME["ice"], PPT_THEME["border"], MSO_SHAPE.ROUNDED_RECTANGLE)
        add_shape_block(slide, 5.28, 2.66, 0.7, 1.42, PPT_THEME["white"], PPT_THEME["cyan"], MSO_SHAPE.ROUNDED_RECTANGLE)
        for x, y in [(6.74, 2.38), (7.66, 2.78), (6.9, 4.0), (7.84, 4.36)]:
            add_shape_block(slide, x, y, 0.72, 0.54, PPT_THEME["ice"], PPT_THEME["border"], MSO_SHAPE.ROUNDED_RECTANGLE)
            add_block(slide, x + 0.12, y + 0.22, 0.42, 0.05, PPT_THEME["blue"])


def add_business_engine_layout(slide, title, bullets, industry):
    kind = get_slide_kind(title)
    row_items = bullets[:4] if bullets else ["No content available."]
    section_label = "Differentiation Stack" if kind == "competitive" else "Revenue Engine"
    add_textbox(slide, 0.74, 1.32, 2.6, 0.26, section_label, 12, PPT_THEME["muted"], True)

    if kind == "competitive":
        stack_labels = ["SMB Focus", "All-in-One Platform", "Fast Iteration"]
        stack_items = row_items[:3]
        while len(stack_items) < 3:
            stack_items.append("")
        for i, label in enumerate(stack_labels):
            y = 1.86 + (i * 1.24)
            add_shape_block(slide, 1.08, y, 7.74, 1.02, PPT_THEME["white"], PPT_THEME["border"], MSO_SHAPE.ROUNDED_RECTANGLE)
            add_block(slide, 1.08, y, 0.12, 1.02, PPT_THEME["blue"] if i != 1 else PPT_THEME["cyan"])
            add_shape_block(slide, 1.48, y + 0.32, 0.36, 0.36, PPT_THEME["ice"], PPT_THEME["cyan"], MSO_SHAPE.OVAL)
            add_textbox(slide, 1.6, y + 0.43, 0.12, 0.12, str(i + 1), 6.4, PPT_THEME["navy"], True, PP_ALIGN.CENTER)
            add_textbox(slide, 2.08, y + 0.22, 1.86, 0.22, label, 11.4, PPT_THEME["navy"], True)
            add_textbox(slide, 4.06, y + 0.18, 4.18, 0.52, str(stack_items[i] or ""), 8.4, PPT_THEME["slate"])
    else:
        engine_items = row_items[:3]
        while len(engine_items) < 3:
            engine_items.append("")
        engine_labels = ["Subscription Revenue", "Retention-Led Growth", "Premium Upsell"]
        add_shape_block(slide, 3.56, 2.22, 2.36, 0.9, PPT_THEME["navy"], PPT_THEME["blue"], MSO_SHAPE.ROUNDED_RECTANGLE)
        add_textbox(slide, 3.88, 2.52, 1.72, 0.2, "Core Model", 11.0, PPT_THEME["white"], True, PP_ALIGN.CENTER)
        card_positions = [(0.94, 2.02), (6.48, 2.02), (3.42, 4.24)]
        for i, label in enumerate(engine_labels):
            x, y = card_positions[i]
            add_shape_block(slide, x, y, 2.58, 1.3, PPT_THEME["white"], PPT_THEME["border"], MSO_SHAPE.ROUNDED_RECTANGLE)
            add_block(slide, x, y, 2.58, 0.1, PPT_THEME["blue"] if i != 1 else PPT_THEME["cyan"])
            add_textbox(slide, x + 0.18, y + 0.24, 2.08, 0.22, label, 10.2, PPT_THEME["navy"], True)
            add_textbox(slide, x + 0.18, y + 0.62, 2.16, 0.46, str(engine_items[i] or ""), 7.8, PPT_THEME["slate"])


def add_financial_funding_layout(slide, title, bullets, industry):
    kind = get_slide_kind(title)
    row_items = bullets[:4] if bullets else ["No content available."]
    label = "Financial Signal" if kind == "financial" else "Investor Lens" if kind == "investor_lens" else "Funding Plan"
    add_textbox(slide, 0.74, 1.32, 2.4, 0.26, label, 12, PPT_THEME["muted"], True)

    if kind == "financial":
        final_projection = projection_df.iloc[-1]
        metrics = [
            ("Revenue", final_projection["Revenue"]),
            ("Net Income", final_projection["Net Income"]),
            ("Ending Cash", final_projection["Ending Cash"]),
        ]
        for i, (label, value) in enumerate(metrics):
            x = 0.86 + (i * 2.08)
            add_shape_block(slide, x, 2.02, 1.74, 1.32, PPT_THEME["white"], PPT_THEME["border"], MSO_SHAPE.ROUNDED_RECTANGLE)
            add_block(slide, x, 2.02, 1.74, 0.1, PPT_THEME["blue"] if i != 1 else PPT_THEME["cyan"])
            add_textbox(slide, x + 0.16, 2.34, 1.18, 0.18, label.upper(), 6.5, PPT_THEME["muted"], True)
            add_textbox(slide, x + 0.16, 2.68, 1.34, 0.28, f"${value / 1000:,.0f}K", 16, PPT_THEME["navy"], True)
        add_shape_block(slide, 0.98, 3.92, 3.62, 1.56, PPT_THEME["white"], PPT_THEME["border"], MSO_SHAPE.ROUNDED_RECTANGLE)
        add_textbox(slide, 1.22, 4.14, 1.4, 0.16, "3-Year Revenue", 7.4, PPT_THEME["navy"], True)
        year_values = [projection_df.loc[i, "Revenue"] for i in range(min(3, len(projection_df)))]
        max_value = max(year_values) if year_values else 1
        for i, value in enumerate(year_values):
            height = 0.28 + (0.72 * (value / max_value if max_value else 0))
            x = 1.46 + (i * 0.84)
            add_block(slide, x, 5.12 - height, 0.38, height, PPT_THEME["blue"] if i % 2 == 0 else PPT_THEME["cyan"])
            add_textbox(slide, x - 0.08, 5.2, 0.54, 0.14, f"Y{i + 1}", 6, PPT_THEME["muted"], True, PP_ALIGN.CENTER)
        add_block(slide, 1.28, 5.14, 2.54, 0.025, PPT_THEME["border"])

        insight_box = add_shape_block(slide, 4.92, 3.92, 3.96, 1.56, PPT_THEME["ice"], PPT_THEME["border"], MSO_SHAPE.ROUNDED_RECTANGLE)
        tf = insight_box.text_frame
        tf.word_wrap = True
        tf.margin_left = PPTInches(0.2)
        tf.margin_right = PPTInches(0.2)
        tf.margin_top = PPTInches(0.15)
        tf.margin_bottom = PPTInches(0.1)
        p = tf.paragraphs[0]
        p.text = "Investor Insight"
        p.font.size = Pt(9)
        p.font.bold = True
        p.font.color.rgb = PPT_THEME["navy"]
        for item in row_items[:3]:
            bullet_p = tf.add_paragraph()
            bullet_p.text = str(item or "")
            bullet_p.font.size = Pt(7.8)
            bullet_p.font.color.rgb = PPT_THEME["slate"]
            bullet_p.space_after = Pt(2)
    elif kind == "investor_lens":
        labels = ["Credibility", "Risk", "Evidence needed"]
        unique_items = []
        seen_items = set()
        for item in row_items:
            normalized = re.sub(r"\s+", " ", str(item or "").strip()).lower()
            if normalized and normalized not in seen_items:
                unique_items.append(item)
                seen_items.add(normalized)
        for i, label in enumerate(labels):
            x = 0.86 + (i * 2.86)
            add_shape_block(slide, x, 2.04, 2.36, 2.54, PPT_THEME["white"], PPT_THEME["border"], MSO_SHAPE.ROUNDED_RECTANGLE)
            add_block(slide, x, 2.04, 2.36, 0.1, PPT_THEME["blue"] if i != 1 else PPT_THEME["cyan"])
            add_textbox(slide, x + 0.2, 2.38, 1.68, 0.24, label, 12, PPT_THEME["navy"], True)
            add_textbox(slide, x + 0.2, 3.0, 1.84, 1.12, str(unique_items[i] if i < len(unique_items) else ""), 8.8, PPT_THEME["slate"])
        if len(unique_items) > 3:
            add_pitch_text_card(slide, 1.34, 5.12, 7.2, 0.78, unique_items[3], PPT_THEME["cyan"], PPT_THEME["ice"], 8.8, True)
    else:
        for i, item in enumerate(row_items[:4]):
            x = 0.92 + (i * 2.12)
            add_shape_block(slide, x, 2.02, 1.72, 2.9, PPT_THEME["white"], PPT_THEME["border"], MSO_SHAPE.ROUNDED_RECTANGLE)
            add_block(slide, x, 2.02, 1.72, 0.1, PPT_THEME["blue"] if i % 2 == 0 else PPT_THEME["cyan"])
            add_textbox(slide, x + 0.18, 2.4, 1.18, 0.16, f"Milestone {i + 1}", 6.4, PPT_THEME["muted"], True)
            add_textbox(slide, x + 0.18, 2.9, 1.28, 1.28, str(item or ""), 8.8, PPT_THEME["slate"])
            if i < 3:
                add_shape_block(slide, x + 1.8, 3.26, 0.34, 0.24, PPT_THEME["cyan"], PPT_THEME["cyan"], MSO_SHAPE.RIGHT_ARROW)


def render_content_slide(slide, title, bullets, industry):
    kind = get_slide_kind(title)
    if kind in ("problem", "market"):
        add_hero_insight_layout(slide, title, bullets, industry)
    elif kind in ("solution", "gtm"):
        add_process_flow_layout(slide, title, bullets, industry)
    elif kind == "product":
        add_product_dashboard_layout(slide, title, bullets, industry)
    elif kind in ("business_model", "competitive"):
        add_business_engine_layout(slide, title, bullets, industry)
    elif kind in ("financial", "investor_lens", "funding"):
        add_financial_funding_layout(slide, title, bullets, industry)
    else:
        add_hero_insight_layout(slide, title, bullets, industry)


def add_content_rows(slide, title, bullets):
    add_bullet_card(slide, bullets, 0.62, 1.42, 4.15, 4.95)
    add_slide_insight_panel(slide, title, bullets, selected_ppt_industry)

# Title Slide
slide = prs.slides.add_slide(prs.slide_layouts[6])
cover_dark = selected_ppt_industry not in ("Consumer Product", "Food / Delivery")
cover_bg = PPT_THEME["navy"] if cover_dark else PPT_THEME["soft"]
cover_text = PPT_THEME["white"] if cover_dark else PPT_THEME["navy"]
cover_muted = PPT_THEME["pale_blue"] if cover_dark else PPT_THEME["slate"]
cover_card = PPT_THEME["navy_2"] if cover_dark else PPT_THEME["white"]
add_block(slide, 0, 0, 10, 7.5, cover_bg)
add_block(slide, 6.08, 0, 3.92, 7.5, PPT_THEME["navy_2"] if cover_dark else PPT_THEME["ice"])
add_theme_cover_visual(slide, selected_ppt_industry)
add_block(slide, 0.65, 0.62, 1.18, 0.08, PPT_THEME["cyan"])
add_block(slide, 0.65, 0.78, 0.52, 0.08, PPT_THEME["blue"])
add_textbox(slide, 0.65, 1.02, 2.7, 0.22, "TURBOPITCH INVESTOR DECK", 7.5, PPT_THEME["cyan"], True)
add_industry_badge(slide, 0.65, 1.34, get_industry_badge_label(selected_ppt_industry), dark=cover_dark)
hero_name = get_ppt_hero_name()
hero_size = 34 if len(hero_name) <= 20 else 30 if len(hero_name) <= 30 else 26
add_textbox(slide, 0.65, 1.82, 5.15, 0.96, hero_name, hero_size, cover_text, True)
add_textbox(
    slide,
    0.68,
    3.0,
    4.95,
    0.7,
    get_ppt_subtitle(),
    16,
    cover_muted,
)
summary_box = add_shape_block(slide, 0.68, 4.28, 4.95, 0.86, cover_card, PPT_THEME["blue"], MSO_SHAPE.ROUNDED_RECTANGLE)
tf = summary_box.text_frame
tf.word_wrap = True
tf.margin_left = PPTInches(0.2)
tf.margin_right = PPTInches(0.2)
tf.margin_top = PPTInches(0.14)
p = tf.paragraphs[0]
p.text = "Financial snapshot, investor readiness, and generated pitch narrative in one branded deck."
p.font.size = Pt(13)
p.font.color.rgb = cover_text

title_metrics = [
    ("Revenue", f"${projection_df.loc[2, 'Revenue'] / 1000:,.0f}K"),
    ("Net Income", f"${projection_df.loc[2, 'Net Income'] / 1000:,.0f}K"),
    ("Ending Cash", f"${projection_df.loc[2, 'Ending Cash'] / 1000:,.0f}K"),
]
for i, (label, value) in enumerate(title_metrics):
    chip_x = 0.68 + (i * 1.64)
    chip = add_shape_block(slide, chip_x, 5.48, 1.42, 0.52, cover_card, PPT_THEME["border"] if not cover_dark else PPT_THEME["blue"], MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = chip.text_frame
    tf.margin_left = PPTInches(0.12)
    tf.margin_right = PPTInches(0.08)
    tf.margin_top = PPTInches(0.08)
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = label.upper()
    p.font.size = Pt(5.8)
    p.font.bold = True
    p.font.color.rgb = PPT_THEME["cyan"] if cover_dark else PPT_THEME["muted"]
    value_p = tf.add_paragraph()
    value_p.text = value
    value_p.font.size = Pt(11.5)
    value_p.font.bold = True
    value_p.font.color.rgb = cover_text
add_textbox(slide, 7.72, 6.18, 1.7, 0.28, "INVESTOR READY", 8, PPT_THEME["cyan"], True, PP_ALIGN.RIGHT)
add_logo(slide, 8.6, 0.52, 0.68)
add_footer(slide, dark=cover_dark, branded=True)

# Financial Overview Slide
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_block(slide, 0, 0, 10, 7.5, PPT_THEME["soft"])
add_theme_content_visual(slide, selected_ppt_industry)
add_dark_header(slide, "Financial Overview", "Projection Dashboard")

year_card_data = [
    ("Year 1", projection_df.loc[0]),
    ("Year 2", projection_df.loc[1]),
    ("Year 3", projection_df.loc[2]),
]

for i, (year_label, year_row) in enumerate(year_card_data):
    card_x = 0.58 + (i * 1.5)
    add_metric_card(
        slide,
        card_x,
        1.34,
        1.36,
        3.65,
        year_label,
        [
            ("Revenue", f"${year_row['Revenue'] / 1000:,.0f}K"),
            ("Net Income", f"${year_row['Net Income'] / 1000:,.0f}K"),
            ("Ending Cash", f"${year_row['Ending Cash'] / 1000:,.0f}K"),
        ],
        accent_color=PPT_THEME["blue"],
    )

ppt_chart = create_ppt_financial_chart_image(projection_df, PPT_THEME)
chart_panel = add_shape_block(slide, 5.1, 1.28, 4.34, 5.0, PPT_THEME["white"], PPT_THEME["border"], MSO_SHAPE.ROUNDED_RECTANGLE)
add_block(slide, 5.1, 1.28, 4.34, 0.1, PPT_THEME["blue"])
tf = chart_panel.text_frame
tf.margin_left = PPTInches(0.18)
tf.margin_top = PPTInches(0.18)
p = tf.paragraphs[0]
p.text = "Revenue / Net Income / Ending Cash"
p.font.size = Pt(10)
p.font.bold = True
p.font.color.rgb = PPT_THEME["navy"]
slide.shapes.add_picture(ppt_chart, PPTInches(5.22), PPTInches(1.82), width=PPTInches(4.04))
add_footer(slide)

# Reality Check Slide
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_block(slide, 0, 0, 10, 7.5, PPT_THEME["soft"])
add_theme_content_visual(slide, selected_ppt_industry)
add_dark_header(slide, "Investor Readiness Check", "Investor Memo")

status_colors = {
    "Green": PPT_THEME["green"],
    "Yellow": PPT_THEME["yellow"],
    "Red": PPT_THEME["red"],
}
status_text_colors = {
    "Green": PPT_THEME["green_text"],
    "Yellow": PPT_THEME["yellow_text"],
    "Red": PPT_THEME["red_text"],
}

overall_status = reality_engine_output["overall"]
overall_box = add_shape_block(slide, 0.58, 1.28, 8.78, 1.02, status_colors.get(overall_status, PPT_THEME["ice"]), PPT_THEME["border"], MSO_SHAPE.ROUNDED_RECTANGLE)
add_block(slide, 0.55, 1.32, 0.12, 1.0, status_text_colors.get(overall_status, PPT_THEME["blue"]))
tf = overall_box.text_frame
tf.word_wrap = True
tf.margin_left = PPTInches(0.26)
tf.margin_right = PPTInches(0.18)
tf.margin_top = PPTInches(0.12)
overall_p = tf.paragraphs[0]
overall_p.text = f"Overall: {overall_status}"
overall_p.font.size = Pt(15)
overall_p.font.bold = True
overall_p.font.color.rgb = PPT_THEME["navy"]
summary_p = tf.add_paragraph()
summary_p.text = reality_engine_output["summary"]
summary_p.font.size = Pt(10)
summary_p.font.color.rgb = PPT_THEME["slate"]

credible_checks = [
    f"{label}: {item['message']}"
    for label, item in reality_engine_output["checks"].items()
    if item["status"] == "Green"
]
challenged_checks = [
    f"{label}: {item['status']} - {item['message']}"
    for label, item in reality_engine_output["checks"].items()
    if item["status"] != "Green"
]
recommended_fixes = [
    f"Prepare evidence for {label.lower()}: {item['message']}"
    for label, item in reality_engine_output["checks"].items()
    if item["status"] != "Green"
]

sections = [
    ("What looks credible", credible_checks or ["No Green checks yet."], PPT_THEME["green"]),
    ("What investors may challenge", challenged_checks or ["No major challenges flagged."], PPT_THEME["yellow"]),
    ("Recommended fix", recommended_fixes or ["Keep the assumptions supported with clear evidence."], PPT_THEME["ice"]),
]

for i, (section_title, section_items, fill_color) in enumerate(sections):
    add_bullet_panel(
        slide,
        0.58 + (i * 3.0),
        2.72,
        2.78,
        3.44,
        section_title,
        section_items[:4],
        fill_color,
        PPT_THEME["navy"],
    )
add_footer(slide)

# Pitch Deck Content Slides
for title, bullets in deck_slides:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_content_slide_frame(slide, title, "Pitch Deck")
    render_content_slide(slide, title, bullets, selected_ppt_industry)
    add_footer(slide)

ppt_buffer = io.BytesIO()
prs.save(ppt_buffer)
ppt_buffer.seek(0)

# ---------------- Downloads UI ----------------
if st.session_state.tp_step == 4:
    with tab8:
        st.download_button(
            label="Download Business Plan (Word)",
            data=doc_buffer,
            file_name="turbopitch_business_plan.docx",
            mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        )
        st.download_button(
            label="Download Financial Model (Excel)",
            data=excel_buffer,
            file_name="turbopitch_financial_model.xlsx",
            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        )
        st.download_button(
            label="Download Pitch Deck (PowerPoint)",
            data=ppt_buffer,
            file_name="turbopitch_pitch_deck.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
elif st.session_state.tp_step == 5:
    st.download_button(
        label="Download Business Plan (Word)",
        data=doc_buffer,
        file_name="turbopitch_business_plan.docx",
        mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    )
    st.download_button(
        label="Download Financial Model (Excel)",
        data=excel_buffer,
        file_name="turbopitch_financial_model.xlsx",
        mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    )
    st.download_button(
        label="Download Pitch Deck (PowerPoint)",
        data=ppt_buffer,
        file_name="turbopitch_pitch_deck.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )

# ==================================================
# NAVIGATION BUTTONS
# ==================================================

st.divider()
render_nav("bottom")

